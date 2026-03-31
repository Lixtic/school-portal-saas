from django import forms
from django.shortcuts import render, redirect, get_object_or_404
from django.db import models
from django.db.models import Q, Count, Sum, Avg
from django.db.utils import OperationalError, ProgrammingError
import django
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from tenants.decorators import require_addon, require_plan
from teachers.addon_utils import requires_addon, requires_addon_freemium, check_freemium_limit, has_addon, get_free_generation_count, FREE_GENERATION_LIMIT
from django.http import JsonResponse, HttpResponse
from decimal import Decimal, InvalidOperation
from django.utils import timezone
import os
import csv
import json
import logging
import random
import string
from datetime import date
from teachers.models import Teacher, DutyWeek, LessonPlan
from academics.models import ClassSubject, AcademicYear, Timetable, SchoolInfo, Resource, Class, Subject, SchemeOfWork
from students.models import Student, Grade, ClassExercise, StudentExerciseScore, ExamType
from students.utils import normalize_term
from .forms import ResourceForm, LessonPlanForm, TeacherCreateForm, TeacherCSVImportForm #, HomeworkForm
from .models import LessonGenerationSession

logger = logging.getLogger(__name__)
from accounts.models import User
from academics.tutor_models import generate_teacher_id_card, export_id_card_to_pdf, export_multiple_id_cards_to_pdf, TutorSession, TutorMessage
from parents.models import Parent
from communication.models import Conversation, Message
from django.views.decorators.http import require_POST
# from parents.models import Homework


def _classify_ai_error(exception):
    raw_message = str(exception or '').strip()
    message_lower = raw_message.lower()

    if (
        'insufficient_quota' in message_lower
        or 'openai http 429' in message_lower
        or ('quota' in message_lower and 'openai' in message_lower)
    ):
        return {
            'error_code': 'ai_quota_exceeded',
            'message': 'Aura is temporarily unavailable. Please try again later or contact your administrator.',
            'http_status': 503,
            'retryable': False,
        }

    if (
        'api key' in message_lower
        or 'missing openai_api_key' in message_lower
        or 'not configured' in message_lower
    ):
        return {
            'error_code': 'ai_not_configured',
            'message': 'Aura AI is currently unavailable due to configuration. Please contact your administrator.',
            'http_status': 503,
            'retryable': False,
        }

    if 'network error' in message_lower or 'timed out' in message_lower:
        return {
            'error_code': 'ai_network_error',
            'message': 'Aura AI is temporarily unreachable. Please check your connection and try again shortly.',
            'http_status': 503,
            'retryable': True,
        }

    if 'openai http 5' in message_lower:
        return {
            'error_code': 'ai_upstream_error',
            'message': 'Aura is temporarily unavailable. Please try again shortly.',
            'http_status': 503,
            'retryable': True,
        }

    return {
        'error_code': 'ai_generation_failed',
        'message': raw_message or 'Aura could not complete this request right now.',
        'http_status': 500,
        'retryable': True,
    }


def _ai_json_error_response(exception):
    error_info = _classify_ai_error(exception)
    from tenants.ai_quota import QuotaExceeded
    if isinstance(exception, QuotaExceeded):
        return JsonResponse(
            {
                'status': 'error',
                'error_code': 'quota_exceeded',
                'message': exception.user_message,
                'used': exception.used,
                'limit': exception.limit,
                'addon_boost': exception.addon_boost,
                'retryable': False,
            },
            status=429,
        )
    return JsonResponse(
        {
            'status': 'error',
            'message': error_info['message'],
            'error_code': error_info['error_code'],
            'retryable': error_info['retryable'],
        },
        status=error_info['http_status'],
    )


# Homework views moved to 'homework' app

@login_required
def teacher_ai_insights(request):
    """Dashboard for teachers to view AI Tutor usage and common student questions"""
    if request.user.user_type not in ['admin', 'teacher']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    # Filter sessions - show all for admin, show relevant classes for teacher
    current_year = AcademicYear.objects.filter(is_current=True).first()
    
    sessions_qs = TutorSession.objects.select_related('student', 'student__user', 'student__current_class', 'subject')
    
    if request.user.user_type == 'teacher':
        teacher = get_object_or_404(Teacher, user=request.user)
        # Get classes taught by teacher
        teacher_classes = ClassSubject.objects.filter(teacher=teacher, class_name__academic_year=current_year).values_list('class_name', flat=True)
        # Also managed classes
        managed_classes = Class.objects.filter(class_teacher=teacher, academic_year=current_year).values_list('id', flat=True)
        
        relevant_class_ids = set(list(teacher_classes) + list(managed_classes))
        sessions_qs = sessions_qs.filter(student__current_class__id__in=relevant_class_ids)

    # 1. Total Stats
    total_sessions = sessions_qs.count()
    if total_sessions == 0:
         return render(request, 'teachers/ai_insights.html', {'no_data': True})

    active_students_count = sessions_qs.values('student').distinct().count()
    avg_msgs_per_session = sessions_qs.aggregate(avg=Avg('message_count'))['avg'] or 0

    # 2. Most Active Students
    top_students = sessions_qs.values(
        'student__user__first_name', 'student__user__last_name', 'student__current_class__name'
    ).annotate(
        session_count=Count('id'),
        total_msgs=Sum('message_count')
    ).order_by('-session_count')[:5]

    # 3. Common Topics (Flatten topics_discussed)
    # This is slightly complex with JSONField in older Django, but simple in Python
    # Since JSONField query support varies, Python processing for top topics is acceptable for small scale
    recent_sessions = sessions_qs.order_by('-started_at')[:50]
    all_topics = []
    for s in recent_sessions:
        if isinstance(s.topics_discussed, list):
            all_topics.extend([str(t).lower().strip() for t in s.topics_discussed if t])
    
    from collections import Counter
    topic_counts = Counter(all_topics).most_common(8)

    # 4. Recent Activity Feed
    recent_activity = sessions_qs.order_by('-started_at')[:10]

    # 5. Scheme Progress: cross-reference session topics against SchemeOfWork topics
    scheme_progress = []
    try:
        if request.user.user_type == 'teacher':
            teacher_schemes = SchemeOfWork.objects.filter(
                class_subject__teacher=teacher,
                class_subject__class_name__academic_year=current_year,
            ).select_related('class_subject__subject', 'class_subject__class_name')
        else:
            teacher_schemes = SchemeOfWork.objects.filter(
                academic_year=current_year,
            ).select_related('class_subject__subject', 'class_subject__class_name')[:20]

        all_session_content = set()
        for s in recent_sessions:
            if s.title:
                all_session_content.add(s.title.lower())
            if isinstance(s.topics_discussed, list):
                for t in s.topics_discussed:
                    all_session_content.add(str(t).lower())

        for scheme in teacher_schemes:
            topics = scheme.get_topics()
            if not topics:
                continue
            covered = sum(
                1 for t in topics
                if any(t.lower() in content or content in t.lower() for content in all_session_content)
            )
            scheme_progress.append({
                'scheme': scheme,
                'topics': topics,
                'total': len(topics),
                'covered': covered,
                'pct': round(covered / len(topics) * 100) if topics else 0,
            })
    except Exception:
        pass

    context = {
        'total_sessions': total_sessions,
        'active_students_count': active_students_count,
        'avg_msgs': round(avg_msgs_per_session, 1),
        'top_students': top_students,
        'common_topics': topic_counts,
        'recent_activity': recent_activity,
        'scheme_progress': scheme_progress,
    }
    return render(request, 'teachers/ai_insights.html', context)


@login_required
def teacher_list(request):
    if request.user.user_type != 'admin':
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    if request.method == 'POST':
        action = request.POST.get('action')
        teacher_id = request.POST.get('teacher_id')
        teacher = get_object_or_404(Teacher, id=teacher_id)

        if action == 'assign_class':
            class_id = request.POST.get('class_id')
            if class_id:
                class_obj = get_object_or_404(Class, id=class_id)
                class_obj.class_teacher = teacher
                class_obj.save()
                messages.success(request, f'{teacher.user.get_full_name()} assigned to {class_obj.name}.')
            else:
                messages.error(request, 'Please select a class to assign.')
        elif action == 'clear_assignments':
            removed_count = teacher.managed_classes.update(class_teacher=None)
            messages.success(request, f'Cleared {removed_count} class assignment(s) for {teacher.user.get_full_name()}.')

        return redirect('teachers:teacher_list')
    
    query = request.GET.get('q', '')
    teachers = Teacher.objects.select_related('user').annotate(
        classes_count=Count('managed_classes', distinct=True),
        subjects_count=Count('classsubject', distinct=True)
    )
    
    if query:
        teachers = teachers.filter(
            Q(user__first_name__icontains=query) |
            Q(user__last_name__icontains=query) |
            Q(employee_id__icontains=query) |
            Q(user__email__icontains=query)
        )
    
    teachers = teachers.order_by('user__first_name', 'user__last_name')
    classes = Class.objects.select_related('academic_year').order_by('name')
    
    context = {
        'teachers': teachers,
        'query': query,
        'classes': classes,
    }
    return render(request, 'teachers/teacher_list.html', context)


@login_required
def import_teachers_csv(request):
    """Import teachers from CSV file"""
    if request.user.user_type != 'admin':
        messages.error(request, 'Access denied. Only admins can import teachers.')
        return redirect('teachers:teacher_list')
    
    if request.method == 'POST':
        form = TeacherCSVImportForm(request.POST, request.FILES)
        if form.is_valid():
            csv_file = request.FILES['csv_file']
            
            # Validate file extension
            if not csv_file.name.endswith('.csv'):
                messages.error(request, 'Please upload a valid CSV file.')
                return redirect('teachers:import_csv')
            
            # Read and process CSV
            try:
                decoded_file = csv_file.read().decode('utf-8').splitlines()
                reader = csv.DictReader(decoded_file)
                
                imported_count = 0
                skipped_count = 0
                errors = []
                
                for row_num, row in enumerate(reader, start=2):  # Start at 2 (header is row 1)
                    try:
                        # Handle both naming conventions: "First Name" or "first_name"
                        first_name = (row.get('First Name') or row.get('first_name') or '').strip()
                        last_name = (row.get('Last Name') or row.get('last_name') or '').strip()
                        employee_id = (row.get('Employee ID') or row.get('employee_id') or '').strip()
                        email = (row.get('Email') or row.get('email') or '').strip()
                        phone = (row.get('Phone') or row.get('phone') or '').strip()
                        qualification = (row.get('Qualification') or row.get('qualification') or '').strip()
                        date_of_joining = (row.get('Date of Joining') or row.get('date_of_joining') or '').strip()
                        age = (row.get('Age') or row.get('age') or '').strip()
                        
                        if not first_name:
                            errors.append(f"Row {row_num}: Missing first name")
                            skipped_count += 1
                            continue
                        
                        # Last name is optional
                        if not last_name:
                            last_name = ''
                        
                        # Generate username
                        if last_name:
                            base_username = f"{first_name.lower()}.{last_name.lower()}"
                        else:
                            base_username = first_name.lower()
                        username = base_username.replace(' ', '').replace('-', '')
                        counter = 1
                        while User.objects.filter(username=username).exists():
                            username = f"{base_username.replace(' ', '').replace('-', '')}{counter}"
                            counter += 1
                        
                        # Use provided employee ID or generate one
                        if employee_id and employee_id.upper() not in ['N/A', 'NA', '']:
                            # Clean up employee ID
                            employee_id = employee_id.strip().replace(' ', '')
                            # Check if it already exists
                            if Teacher.objects.filter(employee_id=employee_id).exists():
                                errors.append(f"Row {row_num}: Employee ID '{employee_id}' already exists")
                                skipped_count += 1
                                continue
                        else:
                            # Generate employee ID
                            prefix = 'TCH'
                            for _ in range(100):  # Try up to 100 times
                                suffix = ''.join(random.choices(string.digits, k=4))
                                candidate = f"{prefix}{suffix}"
                                if not Teacher.objects.filter(employee_id=candidate).exists():
                                    employee_id = candidate
                                    break
                        
                        if not employee_id:
                            errors.append(f"Row {row_num}: Could not generate unique employee ID")
                            skipped_count += 1
                            continue
                        
                        # Prepare email
                        if not email or email.upper() in ['N/A', 'NA', '']:
                            email = f"{username}@school.local"
                        
                        # Check if user with this email already exists
                        if User.objects.filter(email=email).exists():
                            email = f"{username}{random.randint(1, 999)}@school.local"
                        
                        # Parse age for date of birth
                        try:
                            teacher_age = int(age) if age else 30
                        except ValueError:
                            teacher_age = 30
                        
                        # Calculate date of birth
                        dob_year = max(1900, date.today().year - teacher_age)
                        date_of_birth = date(dob_year, 1, 1)
                        
                        # Parse date of joining
                        if date_of_joining and date_of_joining.upper() not in ['N/A', 'NA', '']:
                            try:
                                # Try parsing common date formats
                                from datetime import datetime
                                for fmt in ['%Y-%m-%d', '%d/%m/%Y', '%m/%d/%Y', '%d-%m-%Y']:
                                    try:
                                        joining_date = datetime.strptime(date_of_joining, fmt).date()
                                        break
                                    except ValueError:
                                        continue
                                else:
                                    joining_date = date.today()
                            except Exception:
                                joining_date = date.today()
                        else:
                            joining_date = date.today()
                        
                        # Prepare qualification
                        if not qualification or qualification.upper() in ['N/A', 'NA', '']:
                            qualification = 'Bachelor of Education'
                        
                        # Create user
                        user = User.objects.create(
                            username=username,
                            first_name=first_name,
                            last_name=last_name,
                            email=email,
                            user_type='teacher'
                        )
                        user.set_unusable_password()
                        user.save()
                        
                        # Create teacher
                        teacher = Teacher.objects.create(
                            user=user,
                            employee_id=employee_id,
                            date_of_birth=date_of_birth,
                            date_of_joining=joining_date,
                            qualification=qualification
                        )
                        
                        imported_count += 1
                        
                    except Exception as e:
                        errors.append(f"Row {row_num}: {str(e)}")
                        skipped_count += 1
                        continue
                
                # Show results
                if imported_count > 0:
                    messages.success(request, f'Successfully imported {imported_count} teacher(s).')
                
                if skipped_count > 0:
                    messages.warning(request, f'Skipped {skipped_count} row(s) due to errors.')
                
                if errors and len(errors) <= 10:
                    for error in errors:
                        messages.error(request, error)
                elif errors:
                    messages.error(request, f'{len(errors)} errors occurred. First 10: {", ".join(errors[:10])}')
                
                return redirect('teachers:teacher_list')
                
            except Exception as e:
                messages.error(request, f'Error processing CSV file: {str(e)}')
                return redirect('teachers:import_csv')
    else:
        form = TeacherCSVImportForm()
    
    return render(request, 'teachers/import_csv.html', {'form': form})


@login_required
def teacher_detail(request, teacher_id):
    if request.user.user_type != 'admin':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = get_object_or_404(Teacher.objects.select_related('user'), id=teacher_id)
    
    # Get managed classes
    managed_classes = teacher.managed_classes.select_related('academic_year')
    
    # Get subjects taught
    class_subjects = ClassSubject.objects.filter(teacher=teacher).select_related(
        'class_name', 'subject', 'class_name__academic_year'
    )
    
    # Get lesson plans
    lesson_plans = LessonPlan.objects.filter(teacher=teacher).select_related(
        'subject', 'school_class'
    ).order_by('-date_added')[:5]
    

    context = {
        'teacher': teacher,
        'managed_classes': managed_classes,
        'class_subjects': class_subjects,
        'lesson_plans': lesson_plans,
    }
    return render(request, 'teachers/teacher_detail.html', context)


@login_required
@requires_addon('grade-insight-dashboard')
def analytics_dashboard(request):
    if request.user.user_type not in ['teacher', 'admin']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    # Try to get the teacher object
    teacher = None
    if request.user.user_type == 'teacher':
        teacher = get_object_or_404(Teacher, user=request.user)
    
    # Get classes taught by the teacher
    classes = Class.objects.none()
    selected_class = None
    students_data = []
    
    if teacher:
        # Teachers see their assigned classes
        classes = Class.objects.filter(classsubject__teacher=teacher).distinct()
    elif request.user.user_type == 'admin':
        # Admins see all classes
        classes = Class.objects.all()

    class_id = request.GET.get('class_id')
    if class_id:
        selected_class = get_object_or_404(Class, id=class_id)
    elif classes.exists():
        selected_class = classes.first()
    
    # Aggregated class-level insights
    top_strength = None
    common_misconceptions_list = []
    ai_engagement_count = 0
    
    if selected_class:
        from academics.tutor_models import LearnerMemory
        from collections import Counter
        
        students = Student.objects.filter(current_class=selected_class).select_related('user')
        
        # Prefetch LearnerMemory and last session in bulk
        memory_map = {}
        try:
            memories = LearnerMemory.objects.filter(student__in=students)
            memory_map = {m.student_id: m for m in memories}
        except (OperationalError, ProgrammingError):
            pass  # Table may not exist yet
        
        # Prefetch last TutorSession per student — single aggregated query
        last_session_map = {}
        session_count_map = {}
        try:
            from django.db.models import Max
            session_stats = (
                TutorSession.objects.filter(student__in=students)
                .values('student_id')
                .annotate(
                    last_started=Max('started_at'),
                    total_sessions=Count('id'),
                )
            )
            for row in session_stats:
                last_session_map[row['student_id']] = row['last_started']
                session_count_map[row['student_id']] = row['total_sessions']
        except (OperationalError, ProgrammingError):
            pass
        
        # Prefetch average grade per student (current academic year)
        current_year = AcademicYear.objects.filter(is_current=True).first()
        grade_avg_map = {}
        if current_year:
            grade_avgs = (
                Grade.objects.filter(
                    student__in=students,
                    academic_year=current_year,
                    total_score__gt=0,
                )
                .values('student_id')
                .annotate(avg_score=Avg('total_score'))
            )
            grade_avg_map = {g['student_id']: float(g['avg_score']) for g in grade_avgs}
        
        # Collect class-wide topics for aggregation
        all_mastered = []
        all_misconceptions = []
        
        for student in students:
            memory = memory_map.get(student.id)
            grade_avg = grade_avg_map.get(student.id)
            last_session_time = last_session_map.get(student.id)
            sessions = session_count_map.get(student.id, 0)
            
            has_ai_data = memory is not None and (memory.total_sessions_analysed or 0) > 0
            has_grade_data = grade_avg is not None
            
            # ── Calculate mastery score ──
            mastery = None
            misconception = None
            suggested_intervention = None
            
            if has_ai_data:
                mastered = list(memory.mastered_topics or [])
                misconceptions = list(memory.active_misconceptions or [])
                weaknesses = list(memory.weaknesses or [])
                strengths = list(memory.strengths or [])
                
                all_mastered.extend(mastered)
                all_misconceptions.extend(misconceptions)
                
                # AI mastery heuristic: ratio of mastered to total known topics
                total_signals = len(mastered) + len(misconceptions) + len(weaknesses)
                if total_signals > 0:
                    ai_mastery = round(len(mastered) / total_signals * 100)
                else:
                    ai_mastery = 70  # Sessions exist but no clear signals yet
                
                # Blend with grade average if available
                if has_grade_data:
                    mastery = round(ai_mastery * 0.4 + grade_avg * 0.6)
                else:
                    mastery = ai_mastery
                
                # Pick top misconception/weakness for intervention
                if misconceptions:
                    misconception = misconceptions[0]
                    suggested_intervention = f"Focus on: {misconceptions[0]}"
                elif weaknesses:
                    misconception = weaknesses[0]
                    suggested_intervention = f"Review: {weaknesses[0]}"
            elif has_grade_data:
                mastery = round(grade_avg)
            
            # Determine status
            if mastery is not None:
                if mastery >= 80:
                    status = 'green'
                elif mastery >= 60:
                    status = 'yellow'
                    if not misconception and has_grade_data and grade_avg < 70:
                        suggested_intervention = "Grades trending down — consider review session"
                else:
                    status = 'red'
                    if not suggested_intervention:
                        suggested_intervention = "Low performance — schedule 1-on-1 review"
            else:
                status = 'secondary'  # No data at all
                mastery = 0
            
            if has_ai_data or last_session_time:
                ai_engagement_count += 1
            
            # Format "last active"
            last_active = None
            if last_session_time:
                delta = timezone.now() - last_session_time
                if delta.days > 0:
                    last_active = f"{delta.days}d ago"
                elif delta.seconds >= 3600:
                    last_active = f"{delta.seconds // 3600}h ago"
                elif delta.seconds >= 60:
                    last_active = f"{delta.seconds // 60}m ago"
                else:
                    last_active = "just now"
            
            students_data.append({
                'name': student.user.get_full_name(),
                'id': student.id,
                'mastery': mastery,
                'status': status,
                'misconception': misconception,
                'intervention': suggested_intervention,
                'last_active': last_active,
                'has_ai_data': has_ai_data,
                'session_count': sessions,
            })
        
        # ── Class-wide aggregated insights ──
        if all_mastered:
            strength_counts = Counter(all_mastered).most_common(1)
            top_strength = strength_counts[0][0] if strength_counts else None
        
        if all_misconceptions:
            misc_counts = Counter(all_misconceptions).most_common(5)
            common_misconceptions_list = [{'topic': t, 'count': c} for t, c in misc_counts]
            
    # Calculate class aggregate stats
    avg_mastery = 0
    intervention_count = 0
    scored_students = [s for s in students_data if s['status'] != 'secondary']
    if scored_students:
        avg_mastery = sum(s['mastery'] for s in scored_students) / len(scored_students)
        intervention_count = sum(1 for s in students_data if s['status'] == 'red')

    # ── Grade distribution for the selected class ──
    grade_distribution = {'A': 0, 'B': 0, 'C': 0, 'D': 0, 'F': 0}
    if selected_class and current_year:
        for s in students_data:
            m = s.get('mastery', 0)
            if s['status'] != 'secondary':
                if m >= 80:   grade_distribution['A'] += 1
                elif m >= 70: grade_distribution['B'] += 1
                elif m >= 60: grade_distribution['C'] += 1
                elif m >= 50: grade_distribution['D'] += 1
                else:         grade_distribution['F'] += 1

    # ── Class-level attendance rate (last 30 days) ──
    class_attendance_rate = None
    if selected_class:
        from students.models import Attendance as _Att
        import datetime as _dt
        thirty_ago = timezone.now().date() - _dt.timedelta(days=30)
        att_total = _Att.objects.filter(
            student__current_class=selected_class, date__gte=thirty_ago
        ).count()
        att_present = _Att.objects.filter(
            student__current_class=selected_class, date__gte=thirty_ago, status='present'
        ).count()
        class_attendance_rate = round(att_present / att_total * 100, 1) if att_total else None

    # ── Lesson plan completion (current week) ──
    lesson_plan_completion = None
    if teacher and selected_class:
        try:
            iso = timezone.now().isocalendar()
            current_week_num = iso[1]
            submitted = LessonPlan.objects.filter(
                teacher=teacher, school_class=selected_class, week_number=current_week_num
            ).count()
            expected = ClassSubject.objects.filter(
                teacher=teacher, class_name=selected_class
            ).count()
            lesson_plan_completion = {
                'submitted': submitted,
                'expected': expected,
                'rate': round(submitted / expected * 100) if expected > 0 else 0,
            }
        except Exception:
            pass

    context = {
        'classes': classes,
        'selected_class': selected_class,
        'students_data': students_data,
        'avg_mastery': round(avg_mastery, 1),
        'intervention_count': intervention_count,
        'top_strength': top_strength,
        'common_misconceptions_list': common_misconceptions_list,
        'ai_engagement_count': ai_engagement_count,
        'total_students': len(students_data),
        'grade_distribution': grade_distribution,
        'class_attendance_rate': class_attendance_rate,
        'lesson_plan_completion': lesson_plan_completion,
    }
    return render(request, 'teachers/analytics_dashboard.html', context)




@login_required
@require_plan('pro', 'enterprise')
@requires_addon('grade-insight-dashboard')
def generate_remedial_lesson(request):
    """
    Called from analytics dashboard to 'Generate Lesson' for a misconception.
    Redirects to lesson plan creator with a pre-filled topic.
    """
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
        
    topic = request.GET.get('topic', 'General Review')
    
    # We simply redirect to the creation form, passing the topic.
    # The lesson_plan_create view handles the 'Artificial Intelligence' part by pre-filling data.
    from django.urls import reverse
    import urllib.parse
    
    encoded_topic = urllib.parse.quote(topic)
    url = reverse('teachers:lesson_plan_create')
    return redirect(f"{url}?topic={encoded_topic}")






@login_required
def assign_class_teacher(request, class_id):
    if request.user.user_type != 'admin':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    class_obj = get_object_or_404(Class, id=class_id)
    
    if request.method == 'POST':
        # Check if this is an AJAX request to add a new teacher
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest' and request.POST.get('action') == 'add_teacher':
            from django.http import JsonResponse
            import random
            import string
            
            try:
                # Extract form data
                first_name = request.POST.get('first_name', '').strip()
                last_name = request.POST.get('last_name', '').strip()
                email = request.POST.get('email', '').strip()
                phone = request.POST.get('phone', '').strip()
                employee_id = request.POST.get('employee_id', '').strip()
                qualification = request.POST.get('qualification', '').strip()
                date_of_birth = request.POST.get('date_of_birth', '').strip()
                date_of_joining = request.POST.get('date_of_joining', '').strip()
                
                # Validate required fields
                if not first_name or not last_name or not email:
                    return JsonResponse({'success': False, 'error': 'First name, last name, and email are required'})
                
                # Auto-generate Employee ID if not provided
                if not employee_id:
                    while True:
                        suffix = ''.join(random.choices(string.digits, k=4))
                        employee_id = f"TCH{suffix}"
                        if not Teacher.objects.filter(employee_id=employee_id).exists() and not User.objects.filter(username=employee_id).exists():
                            break
                
                # Check if username already exists
                username = employee_id
                if User.objects.filter(username=username).exists():
                    return JsonResponse({'success': False, 'error': f'Employee ID {employee_id} already exists'})
                
                # Check if email already exists
                if User.objects.filter(email=email).exists():
                    return JsonResponse({'success': False, 'error': f'Email {email} already exists'})
                
                # Create user
                user = User.objects.create_user(
                    username=username,
                    email=email,
                    password=username,  # Default password is employee ID
                    first_name=first_name,
                    last_name=last_name,
                    user_type='teacher',
                    phone=phone
                )
                
                # Create teacher
                from datetime import datetime, date
                teacher = Teacher(
                    user=user,
                    employee_id=employee_id,
                    qualification=qualification or 'Not provided'
                )
                
                # Set dates if provided
                if date_of_birth:
                    try:
                        teacher.date_of_birth = datetime.strptime(date_of_birth, '%Y-%m-%d').date()
                    except ValueError:
                        pass
                
                if date_of_joining:
                    try:
                        teacher.date_of_joining = datetime.strptime(date_of_joining, '%Y-%m-%d').date()
                    except ValueError:
                        teacher.date_of_joining = date.today()
                else:
                    teacher.date_of_joining = date.today()
                
                teacher.save()
                
                return JsonResponse({
                    'success': True,
                    'message': f'Teacher {user.get_full_name()} added successfully',
                    'teacher_id': teacher.id,
                    'teacher_name': user.get_full_name(),
                    'employee_id': employee_id
                })
                
            except Exception as e:
                return JsonResponse({'success': False, 'error': str(e)})
        
        # Regular form submission for assigning teacher
        teacher_id = request.POST.get('teacher_id')
        if teacher_id:
            try:
                teacher = Teacher.objects.get(id=teacher_id)
                class_obj.class_teacher = teacher
                class_obj.save()
                messages.success(request, f'{teacher.user.get_full_name()} assigned as class teacher for {class_obj.name}')
            except Teacher.DoesNotExist:
                messages.error(request, 'Invalid teacher selected')
        else:
            class_obj.class_teacher = None
            class_obj.save()
            messages.success(request, f'Class teacher removed from {class_obj.name}')
        
        return redirect('academics:manage_classes')
    
    teachers = Teacher.objects.select_related('user').order_by('user__first_name', 'user__last_name')
    
    from datetime import date
    context = {
        'class_obj': class_obj,
        'teachers': teachers,
        'today': date.today().strftime('%Y-%m-%d')
    }
    return render(request, 'teachers/assign_class_teacher.html', context)


@login_required
def edit_teacher(request, teacher_id):
    from .forms import TeacherEditForm
    
    if request.user.user_type != 'admin':
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if request.method == 'POST':
        form = TeacherEditForm(request.POST, request.FILES, instance=teacher)
        if form.is_valid():
            form.save()
            messages.success(request, f'Teacher profile for {teacher.user.get_full_name()} updated successfully.')
            return redirect('teachers:teacher_detail', teacher_id=teacher.id)
    else:
        form = TeacherEditForm(instance=teacher)
    
    context = {
        'form': form,
        'teacher': teacher
    }
    return render(request, 'teachers/edit_teacher.html', context)



@login_required
def teacher_classes(request):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = Teacher.objects.get(user=request.user)
    class_subjects = ClassSubject.objects.filter(teacher=teacher).select_related(
        'class_name', 'subject'
    )
    
    return render(request, 'teachers/my_classes.html', {'class_subjects': class_subjects})


def _is_basic9_class_name(class_name):
    normalized = (class_name or '').strip().upper()
    tokens = ['BASIC 9', 'BASIC9', 'JHS 3', 'JHS3', 'FORM 3', 'GRADE 9']
    return any(token in normalized for token in tokens)


@login_required
def enter_grades(request):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = Teacher.objects.get(user=request.user)
    class_subjects = ClassSubject.objects.filter(teacher=teacher).select_related('class_name', 'subject')
    selected_cs_id = request.GET.get('class_subject_id')
    exam_types = ExamType.objects.filter(is_active=True).order_by('name')
    selected_exam_type_id = request.GET.get('exam_type_id')
    
    if request.method == 'POST':
        term = normalize_term(request.POST.get('term', 'first'))
        class_subject_id = request.POST.get('class_subject_id') or request.POST.get('class_subject')
        exam_type_id = (request.POST.get('exam_type_id') or '').strip()
        student_ids = request.POST.getlist('student_id[]') or request.POST.getlist('student_id')

        academic_year = AcademicYear.objects.filter(is_current=True).first()
        if not academic_year:
            messages.error(request, 'No academic year is marked as current.')
            return redirect('teachers:enter_grades')

        try:
            cs = class_subjects.get(id=class_subject_id)
        except ClassSubject.DoesNotExist:
            messages.error(request, 'Invalid class/subject selection.')
            return redirect('teachers:enter_grades')

        is_basic9_class = _is_basic9_class_name(cs.class_name.name)
        exam_type = None
        if is_basic9_class:
            if not exam_type_id:
                messages.error(request, 'Please select an exam type for Basic 9 grade entry.')
                return redirect('teachers:enter_grades')
            exam_type = ExamType.objects.filter(id=exam_type_id, is_active=True).first()
            if not exam_type:
                messages.error(request, 'Invalid or inactive exam type selected.')
                return redirect('teachers:enter_grades')

        created_count = 0
        for student_id in student_ids:
            overall_raw = request.POST.get(f'overall_score_{student_id}')
            class_score_raw = request.POST.get(f'class_score_{student_id}', '')
            exams_score_raw = request.POST.get(f'exams_score_{student_id}', '')

            class_score = exams_score = None

            # Prefer the explicit class/exam fields; fall back to overall if needed
            try:
                if class_score_raw != '' and exams_score_raw != '':
                    class_score = Decimal(class_score_raw)
                    exams_score = Decimal(exams_score_raw)
                elif overall_raw not in (None, ''):
                    overall = Decimal(overall_raw)
                    overall = max(Decimal('0'), min(overall, Decimal('100')))
                    class_score = overall * Decimal('0.3')
                    exams_score = overall * Decimal('0.7')
                else:
                    continue
            except (InvalidOperation, TypeError):
                continue

            # Server-side validation of maxima
            class_score = max(Decimal('0'), min(class_score, Decimal('30')))
            exams_score = max(Decimal('0'), min(exams_score, Decimal('70')))

            Grade.objects.update_or_create(
                student_id=student_id,
                subject=cs.subject,
                academic_year=academic_year,
                term=term,
                exam_type=exam_type,
                defaults={
                    'class_score': class_score,
                    'exams_score': exams_score,
                    'created_by': request.user
                }
            )
            created_count += 1

        if created_count == 0:
            messages.warning(request, 'No grades were saved. Please ensure you loaded students and entered scores.')
        else:
            if exam_type:
                messages.success(request, f'Grades entered/updated for {created_count} students in {cs.class_name} - {cs.subject} ({exam_type.name}).')
            else:
                messages.success(request, f'Grades entered/updated for {created_count} students in {cs.class_name} - {cs.subject}.')
        return redirect('teachers:enter_grades')
    
    # Auto-select if teacher has exactly one assignment
    default_cs_id = selected_cs_id
    if not default_cs_id and class_subjects.count() == 1:
        default_cs_id = str(class_subjects.first().id)

    return render(request, 'teachers/enter_grades.html', {
        'class_subjects': class_subjects,
        'selected_cs_id': selected_cs_id,
        'default_cs_id': default_cs_id,
        'exam_types': exam_types,
        'selected_exam_type_id': selected_exam_type_id,
    })


@login_required
@require_POST
def scan_grades_sheet(request):
    if request.user.user_type != 'teacher':
        return JsonResponse({'status': 'error', 'message': 'Access denied'}, status=403)
        
    if 'image' not in request.FILES:
        return JsonResponse({'status': 'error', 'message': 'No image provided'}, status=400)
        
    try:
        import base64
        import json
        from academics.ai_tutor import _post_chat_completion, _get_openai_api_key
        
        image_file = request.FILES['image']
        image_b64 = base64.b64encode(image_file.read()).decode('utf-8')
        mime_type = image_file.content_type
        
        prompt = """You are an AI assistant helping a teacher digitize a handwritten or printed student grade sheet.
Please extract the list of student names and their corresponding grades. 
Look for either 'overall_score', or specific 'class_score' and 'exams_score' if broken down.
Return ONLY valid JSON in the following format:
[
  {
     "name": "Student Full Name",
     "overall_score": 85,
     "class_score": 25,
     "exams_score": 60
  }
]
If a score is missing or unable to be read, set it to null. Ensure the response contains no markdown formatting around the output, just raw JSON text. No '''json.
"""

        payload = {
            "model": "gpt-4o-mini",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{image_b64}"
                            }
                        }
                    ]
                }
            ],
            "temperature": 0.1,
            "max_tokens": 1000
        }

        response_data = _post_chat_completion(payload, _get_openai_api_key())
        
        if 'error' in response_data:
             return JsonResponse({'status': 'error', 'message': f"AI Error: {response_data['error']}"}, status=500)

        choices = response_data.get('choices', [])
        if not choices:
            return JsonResponse({'status': 'error', 'message': 'Failed to process image. No response from AI.'}, status=500)
            
        content = choices[0]['message']['content'].strip()
        
        # Strip potential markdown formatting
        if content.startswith('```json'):
            content = content[7:]
        if content.startswith('```'):
            content = content[3:]
        if content.endswith('```'):
            content = content[:-3]
            
        grades_data = json.loads(content.strip())
        return JsonResponse({'status': 'success', 'data': grades_data})
        
    except json.JSONDecodeError:
        return JsonResponse({'status': 'error', 'message': 'AI returned malformed data format. Please try again.'}, status=500)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': f'Server error: {str(e)}'}, status=500)


@login_required
def get_students(request, class_id):
    subject_id = request.GET.get('subject_id')
    term = normalize_term(request.GET.get('term', 'first'))
    exam_type_id = request.GET.get('exam_type_id')
    class_obj = Class.objects.filter(id=class_id).first()
    is_basic9_class = _is_basic9_class_name(class_obj.name) if class_obj else False

    students = Student.objects.filter(current_class_id=class_id).select_related('user')

    academic_year = AcademicYear.objects.filter(is_current=True).first()
    grades_by_student = {}
    if academic_year and subject_id:
        grades = Grade.objects.filter(
            student__current_class_id=class_id,
            subject_id=subject_id,
            academic_year=academic_year,
            term=term,
        )
        if is_basic9_class and exam_type_id:
            grades = grades.filter(exam_type_id=exam_type_id)
        else:
            grades = grades.filter(exam_type__isnull=True)
        grades_by_student = {g.student_id: g for g in grades}

    data = []
    for s in students:
        g = grades_by_student.get(s.id)
        data.append({
            'id': s.id,
            'name': s.user.get_full_name(),
            'admission_number': s.admission_number,
            'class_score': str(g.class_score) if g else '',
            'exams_score': str(g.exams_score) if g else '',
            'total_score': str(g.total_score) if g else '',
            'grade': g.grade if g else '',
            'remarks': g.remarks if g else '',
        })

    return JsonResponse(data, safe=False)


@login_required
def teacher_schedule(request):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    try:
        teacher = Teacher.objects.get(user=request.user)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher profile not found.')
        return redirect('dashboard')

    # Get all schedule entries for this teacher/class-subject
    timetable_qs = Timetable.objects.filter(class_subject__teacher=teacher).select_related(
        'class_subject', 'class_subject__class_name', 'class_subject__subject'
    ).order_by('day', 'start_time')
    
    # Organize by day for simpler template loop
    days_data = []
    # Timetable.DAY_CHOICES is ((0,'Monday'), ...)
    for code, name in Timetable.DAY_CHOICES:
        # Filter in python to avoid N queries, 
        # or just filter in loop if list is short. 
        # Since timetable_qs is likely small ( < 50 items), list filter is fine.
        entries = [t for t in timetable_qs if t.day == code]
        days_data.append({
            'name': name,
            'entries': entries
        })
            
    return render(request, 'teachers/schedule.html', {'days': days_data})

@login_required
def print_duty_roster(request):
    current_year = AcademicYear.objects.filter(is_current=True).first()
    if not current_year:
        current_year = AcademicYear.objects.first()

    today = timezone.now().date()

    req_term = request.GET.get('term')
    if req_term:
        term = req_term
    else:
        current_duty = DutyWeek.objects.filter(
            academic_year=current_year,
            start_date__lte=today,
            end_date__gte=today
        ).first()
        term = current_duty.term if current_duty else (
            'second' if 1 <= today.month <= 4 else
            ('third' if 5 <= today.month <= 8 else 'first')
        )

    year_id = request.GET.get('year', current_year.id if current_year else None)
    year = get_object_or_404(AcademicYear, id=year_id) if year_id else None

    weeks = DutyWeek.objects.filter(
        academic_year=year, term=term
    ).prefetch_related(
        'assignments', 'assignments__teacher', 'assignments__teacher__user'
    ).order_by('week_number')

    # Personalised "My Duty" banner
    my_next_duty = None
    teacher = getattr(request.user, 'teacher', None)
    if teacher:
        # Currently on duty this week?
        my_next_duty = DutyWeek.objects.filter(
            assignments__teacher=teacher,
            start_date__lte=today,
            end_date__gte=today,
        ).prefetch_related('assignments__teacher__user').first()
        # Otherwise, next upcoming duty
        if not my_next_duty:
            my_next_duty = DutyWeek.objects.filter(
                assignments__teacher=teacher,
                start_date__gt=today,
            ).prefetch_related('assignments__teacher__user').order_by('start_date').first()

    context = {
        'weeks': weeks,
        'year': year,
        'term': term,
        'today': today,
        'teacher': teacher,
        'my_next_duty': my_next_duty,
        'school_info': SchoolInfo.objects.first(),
        'available_terms': ['first', 'second', 'third'],
        'academic_years': AcademicYear.objects.all(),
    }
    return render(request, 'teachers/duty_roster_pdf.html', context)


@login_required
def generate_duty_weeks(request):
    """Admin-only: auto-create DutyWeek rows for a full term."""
    if request.user.user_type != 'admin':
        messages.error(request, "Only admins can generate duty weeks.")
        return redirect('teachers:duty_roster')

    if request.method != 'POST':
        return redirect('teachers:duty_roster')

    from datetime import timedelta
    from django.urls import reverse as url_reverse

    year_id   = request.POST.get('year')
    term      = request.POST.get('term', 'first')
    start_str = request.POST.get('start_date', '')
    try:
        num_weeks = max(1, min(int(request.POST.get('num_weeks') or 13), 20))
    except (ValueError, TypeError):
        num_weeks = 13

    year = get_object_or_404(AcademicYear, id=year_id)

    try:
        start = date.fromisoformat(start_str)
    except (ValueError, TypeError):
        messages.error(request, "Invalid start date.")
        return redirect(f"{url_reverse('teachers:duty_roster')}?term={term}&year={year_id}")

    created = skipped = 0
    for i in range(num_weeks):
        week_start = start + timedelta(weeks=i)
        week_end   = week_start + timedelta(days=4)   # Mon → Fri
        _, was_created = DutyWeek.objects.get_or_create(
            academic_year=year,
            term=term,
            week_number=i + 1,
            defaults={'start_date': week_start, 'end_date': week_end},
        )
        if was_created:
            created += 1
        else:
            skipped += 1

    if created:
        messages.success(request, f"{created} duty week(s) generated for {term.title()} Term {year.name}.")
    if skipped:
        messages.info(request, f"{skipped} week(s) already existed and were skipped.")

    return redirect(f"{url_reverse('teachers:duty_roster')}?term={term}&year={year_id}")


@login_required
@requires_addon('exercise-maker')
def manage_exercises(request, class_subject_id):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = Teacher.objects.get(user=request.user)
    class_subject = get_object_or_404(ClassSubject, id=class_subject_id, teacher=teacher)
    
    term = request.GET.get('term', 'first')
    
    exercises = ClassExercise.objects.filter(
        class_subject=class_subject,
        term=term
    ).order_by('-date_assigned')
    
    if request.method == 'POST':
        title = request.POST.get('title')
        max_marks = request.POST.get('max_marks')
        term_input = request.POST.get('term')
        
        if title and max_marks:
            ClassExercise.objects.create(
                class_subject=class_subject,
                term=term_input,
                title=title,
                max_marks=max_marks
            )
            messages.success(request, 'Exercise created successfully.')
            return redirect(f"{request.path}?term={term_input}")

    # Helper for term names
    term_choices = Grade.TERM_CHOICES

    return render(request, 'teachers/manage_exercises.html', {
        'class_subject': class_subject,
        'exercises': exercises,
        'selected_term': term,
        'term_choices': term_choices,
    })


@login_required
def enter_exercise_scores(request, exercise_id):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = Teacher.objects.get(user=request.user)
    exercise = get_object_or_404(ClassExercise, id=exercise_id, class_subject__teacher=teacher)
    
    students = Student.objects.filter(current_class=exercise.class_subject.class_name).order_by('user__last_name')
    
    if request.method == 'POST':
        try:
            with django.db.transaction.atomic():
                updated_count = 0
                for student in students:
                    score_val = request.POST.get(f'score_{student.id}')
                    remarks_val = request.POST.get(f'remarks_{student.id}', '')
                    
                    if score_val:
                        try:
                            score = Decimal(score_val)
                            score = max(Decimal('0'), min(score, exercise.max_marks))
                            
                            StudentExerciseScore.objects.update_or_create(
                                student=student,
                                exercise=exercise,
                                defaults={'score': score, 'remarks': remarks_val}
                            )
                            updated_count += 1
                        except (InvalidOperation, TypeError):
                            continue
                
                # RECALCULATION LOGIC
                all_exercises = ClassExercise.objects.filter(
                    class_subject=exercise.class_subject,
                    term=exercise.term
                )
                total_possible = sum(ex.max_marks for ex in all_exercises)
                
                if total_possible > 0:
                    current_year = AcademicYear.objects.filter(is_current=True).first()
                    if not current_year:
                        current_year = AcademicYear.objects.order_by('-start_date').first()

                    # Pre-fetch all scores for this subject/term to avoid N+1 inside loop if possible 
                    # but simple iteration is fine for class size < 50
                    
                    for student in students:
                        student_scores = StudentExerciseScore.objects.filter(
                            student=student,
                            exercise__in=all_exercises
                        ).aggregate(total=models.Sum('score'))['total'] or Decimal('0')
                        
                        # (Obtained / Possible) * 30
                        new_class_score = (student_scores / total_possible) * Decimal('30')
                        
                        grade_obj, _ = Grade.objects.get_or_create(
                            student=student,
                            subject=exercise.class_subject.subject,
                            academic_year=current_year,
                            term=exercise.term,
                            exam_type=None,
                            defaults={'created_by': request.user}
                        )
                        grade_obj.class_score = new_class_score
                        grade_obj.save()

            messages.success(request, f'Scores saved for {updated_count} students. Class assessment totals updated.')
        except Exception as e:
            messages.error(request, f"Error saving scores: {e}")
            
        return redirect('teachers:enter_exercise_scores', exercise_id=exercise.id)

    existing_scores = StudentExerciseScore.objects.filter(exercise=exercise)
    score_map = {s.student_id: s for s in existing_scores}
    
    student_list = []
    for s in students:
        student_list.append({
            'student': s,
            'score_obj': score_map.get(s.id)
        })

    return render(request, 'teachers/enter_exercise_scores.html', {
        'exercise': exercise,
        'student_list': student_list,
    })


@login_required
def search_students(request):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    query = request.GET.get('q', '').strip()
    class_filter = request.GET.get('class_id', '').strip()
    gender_filter = request.GET.get('gender', '').strip()

    teacher = Teacher.objects.get(user=request.user)

    # Determine which fee structures this teacher is assigned to collect
    from finance.models import FeeStructure as _FeeStructure
    collector_structure_ids = set(
        _FeeStructure.objects.filter(assigned_collector=teacher).values_list('id', flat=True)
    )

    students = Student.objects.select_related('user', 'current_class').distinct()

    if query:
        students = students.filter(
            Q(user__first_name__icontains=query) |
            Q(user__last_name__icontains=query) |
            Q(admission_number__icontains=query)
        )
    if class_filter:
        students = students.filter(current_class__id=class_filter)
    if gender_filter:
        students = students.filter(gender=gender_filter)

    students = students.order_by('user__last_name', 'user__first_name')[:50]

    all_classes = Class.objects.filter(
        academic_year__is_current=True
    ).order_by('name')

    searched = bool(query or class_filter or gender_filter)

    return render(request, 'teachers/search_results.html', {
        'query': query,
        'students': students if searched else [],
        'school_name': SchoolInfo.objects.first().name if SchoolInfo.objects.exists() else 'School',
        'all_classes': all_classes,
        'class_filter': class_filter,
        'gender_filter': gender_filter,
        'collector_structure_ids': collector_structure_ids,
        'searched': searched,
    })


@login_required
def curriculum_library(request):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    teacher = Teacher.objects.get(user=request.user)

    resource_fields_available = False
    try:
        from django.db import connection
        with connection.cursor() as cursor:
            cols = [col.name for col in connection.introspection.get_table_description(cursor, Resource._meta.db_table)]
        resource_fields_available = 'resource_type' in cols and 'curriculum' in cols
    except (OperationalError, ProgrammingError) as exc:
        logger.debug('Resource field introspection failed: %s', exc)
        resource_fields_available = False

    try:
        qs = Resource.objects.filter(
            Q(target_audience__in=['teachers', 'all']) | Q(class_subject__teacher=teacher)
        ).order_by('-uploaded_at')
        if resource_fields_available:
            qs = qs.filter(resource_type='curriculum')
        else:
            qs = qs.filter(class_subject__isnull=True)
            qs = qs.only('id', 'title', 'description', 'file', 'link', 'uploaded_at')
        resources = list(qs)
    except (OperationalError, ProgrammingError):
        resources = []
        resource_fields_available = False

    curriculum_options = [('all', 'All Curricula')]
    if resource_fields_available:
        curriculum_options = [('all', 'All Curricula')] + list(Resource.CURRICULUM_CHOICES)

    selected_curriculum = request.GET.get('curriculum', 'all')
    if resource_fields_available and selected_curriculum != 'all':
        resources = [r for r in resources if getattr(r, 'curriculum', None) == selected_curriculum]

    return render(request, 'teachers/curriculum_library.html', {
        'resources': resources,
        'resource_fields_available': resource_fields_available,
        'curriculum_options': curriculum_options,
        'selected_curriculum': selected_curriculum,
    })


@login_required
def class_resources(request, class_subject_id):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = Teacher.objects.get(user=request.user)
    class_subject = get_object_or_404(ClassSubject, id=class_subject_id, teacher=teacher)

    # Detect availability of new fields on the DB
    resource_fields_available = False
    try:
        from django.db import connection
        with connection.cursor() as cursor:
            cols = [col.name for col in connection.introspection.get_table_description(cursor, Resource._meta.db_table)]
        resource_fields_available = 'resource_type' in cols and 'curriculum' in cols
    except (OperationalError, ProgrammingError) as exc:
        logger.debug('Resource field introspection failed: %s', exc)
        resource_fields_available = False

    resource_type_filter = request.GET.get('type', 'all')
    curriculum_filter = request.GET.get('curriculum', 'all')
    
    if request.method == 'POST':
        form = ResourceForm(request.POST, request.FILES)
        if form.is_valid():
            resource = form.save(commit=False)
            resource.class_subject = class_subject
            resource.uploaded_by = request.user
            resource.save()
            messages.success(request, 'Resource uploaded successfully.')
            return redirect('teachers:class_resources', class_subject_id=class_subject.id)
    else:
        form = ResourceForm(initial={
            'class_subject': class_subject.id,
            'resource_type': 'teaching',
            'curriculum': 'ges_jhs_new',
        })

    if 'class_subject' in form.fields:
        form.fields['class_subject'].widget = forms.HiddenInput()
        form.fields['class_subject'].initial = class_subject.id

    base_qs = Resource.objects.filter(class_subject=class_subject).order_by('-uploaded_at')
    if not resource_fields_available:
        resources = base_qs.only('id', 'title', 'description', 'file', 'link', 'uploaded_at')
        curriculum_options = [('all', 'All Curricula')]
        resource_type_filter = 'all'
        curriculum_filter = 'all'
    else:
        resources = base_qs
        if resource_type_filter in ['curriculum', 'teaching']:
            resources = resources.filter(resource_type=resource_type_filter)
        if curriculum_filter != 'all':
            resources = resources.filter(curriculum=curriculum_filter)
        curriculum_options = [('all', 'All Curricula')] + list(Resource.CURRICULUM_CHOICES)
    
    return render(request, 'teachers/class_resources.html', {
        'class_subject': class_subject,
        'resources': resources,
        'form': form,
        'resource_type_filter': resource_type_filter,
        'curriculum_filter': curriculum_filter,
        'curriculum_options': curriculum_options,
        'resource_fields_available': resource_fields_available,
    })

@login_required
def delete_resource(request, resource_id):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
        
    resource = get_object_or_404(Resource, id=resource_id, class_subject__teacher__user=request.user)
    class_subject_id = resource.class_subject.id
    resource.delete()
    messages.success(request, 'Resource deleted successfully.')
    return redirect('teachers:class_resources', class_subject_id=class_subject_id)



@login_required
def lesson_plan_list(request):
    is_teacher = request.user.user_type == 'teacher'
    is_admin = request.user.user_type == 'admin'
    if not is_teacher and not is_admin:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = None
    if is_teacher:
        teacher = get_object_or_404(Teacher, user=request.user)
    
    week = request.GET.get('week', '').strip()
    subject_id = request.GET.get('subject', '').strip()
    class_id = request.GET.get('school_class', '').strip()
    query = request.GET.get('q', '').strip()
    aura_only = request.GET.get('aura_only', '').strip()

    teacher_subjects = []
    teacher_classes = []
    teacher_list = []
    selected_teacher = request.GET.get('teacher', '').strip()

    if is_teacher:
        assigned_pairs = ClassSubject.objects.filter(teacher=teacher).select_related('class_name', 'subject')
        teacher_subjects = [pair.subject for pair in assigned_pairs]
        teacher_classes = [pair.class_name for pair in assigned_pairs]
    else:
        current_year = AcademicYear.objects.filter(is_current=True).first()
        class_qs = Class.objects.all()
        if current_year:
            class_qs = class_qs.filter(academic_year=current_year)
        teacher_subjects = Subject.objects.all().order_by('name')
        teacher_classes = class_qs.order_by('name')
        teacher_list = Teacher.objects.select_related('user').order_by('user__first_name', 'user__last_name')

    # de-duplicate while preserving order
    seen_subject_ids = set()
    teacher_subjects = [s for s in teacher_subjects if not (s.id in seen_subject_ids or seen_subject_ids.add(s.id))]
    seen_class_ids = set()
    teacher_classes = [c for c in teacher_classes if not (c.id in seen_class_ids or seen_class_ids.add(c.id))]

    try:
        # Force evaluation to catch DB errors if table doesn't exist yet
        lesson_plans_qs = LessonPlan.objects.select_related('subject', 'school_class', 'teacher', 'teacher__user')
        if is_teacher:
            lesson_plans_qs = lesson_plans_qs.filter(teacher=teacher)
        elif selected_teacher:
            lesson_plans_qs = lesson_plans_qs.filter(teacher_id=selected_teacher)
        if week:
            lesson_plans_qs = lesson_plans_qs.filter(week_number=week)
        if subject_id:
            lesson_plans_qs = lesson_plans_qs.filter(subject_id=subject_id)
        if class_id:
            lesson_plans_qs = lesson_plans_qs.filter(school_class_id=class_id)
        if query:
            lesson_plans_qs = lesson_plans_qs.filter(
                Q(topic__icontains=query) |
                Q(objectives__icontains=query) |
                Q(presentation__icontains=query)
            )
        if aura_only:
            lesson_plans_qs = lesson_plans_qs.filter(introduction__icontains='PULSE CHECK')
        lesson_plans_qs = lesson_plans_qs.order_by('school_class__name', '-date_added')
        lesson_plans = list(lesson_plans_qs)
        aura_drafted_count = sum(1 for p in lesson_plans if 'PULSE CHECK' in (p.introduction or ''))
    except (OperationalError, ProgrammingError):
        lesson_plans = []
        teacher_subjects = []
        teacher_classes = []
        aura_drafted_count = 0
        messages.warning(request, "Lesson Plan system is initializing. Please try again later.")
        
    return render(request, 'teachers/lesson_plan_list.html', {
        'lesson_plans': lesson_plans,
        'aura_drafted_count': aura_drafted_count,
        'selected_week': week,
        'selected_subject': subject_id,
        'selected_class': class_id,
        'selected_teacher': selected_teacher,
        'selected_query': query,
        'selected_aura_only': aura_only,
        'teacher_subjects': teacher_subjects,
        'teacher_classes': teacher_classes,
        'teacher_list': teacher_list,
        'is_teacher': is_teacher,
        'is_admin': is_admin,
    })
        
@login_required

@login_required
def lesson_plan_create(request):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = get_object_or_404(Teacher, user=request.user)
    
    # Check for AI generation request
    initial_data = {}
    
    # CASE 1: Full AI Content passed via POST (from Chat) or huge GET
    # Ideally we'd use POST for large content, but if coming from a link, it's GET.
    # We might need to parse the unstructured text from Aura-T.
    
    ai_content = request.POST.get('ai_content') or request.GET.get('ai_content')
    topic = request.GET.get('topic') # Legacy/Simple Fallback
    indicator = request.GET.get('indicator', '').strip()  # Full indicator text from SOW chip
    # Pre-fill subject/class when coming from Scheme of Work topic chip
    prefill_subject_id = request.GET.get('subject_id')
    prefill_class_id   = request.GET.get('class_id')

    if ai_content:
        # PARSING LOGIC: Extract fields from the Aura-T standardized format
        import re
        
        def extract_section(header, text):
            # Allow for optional suffix inside the bold tag (like duration "[10 mins]")
            # Pattern matches: **Header[optional text]**:? content...
            # We use re.escape(header) but need to allow trailing chars before the closing **
            # Matches: **PHASE 1: STARTER [10mins]** or **PHASE 1: STARTER**
            pattern = re.compile(rf"\*\*{re.escape(header)}.*?\*\*\s*(.*?)(?=\n\*\*|\Z)", re.DOTALL | re.IGNORECASE)
            match = pattern.search(text)
            return match.group(1).strip() if match else ""

        def extract_phase(phase_num, text):
            """Extract full phase content up to the next phase-level boundary."""
            boundaries = {1: r'\*\*PHASE\s*2', 2: r'\*\*PHASE\s*3', 3: r'\*\*(?:Resources|Homework)'}
            end = boundaries.get(phase_num, r'\Z')
            pat = re.compile(rf"\*\*PHASE\s*{phase_num}[^*]*\*\*\s*(.*?)(?={end}|\Z)", re.DOTALL | re.IGNORECASE)
            m = pat.search(text)
            return m.group(1).strip() if m else ""

        # Attempt to map Aura-T output format to LessonPlan model fields
        parsed_topic = extract_section("Strand", ai_content)
        parsed_sub_strand = extract_section("Sub Strand", ai_content)
        if not parsed_topic:
            parsed_topic = topic or "New Lesson Plan"
            
        initial_data = {
            'topic': parsed_topic,
            'sub_strand': parsed_sub_strand,
            'week_number': extract_section("WEEK", ai_content) or 1,
            'objectives': extract_section("Content Standard", ai_content) + "\n\n" + extract_section("Indicator", ai_content),
            'teaching_materials': extract_section("Resources", ai_content),
            'introduction': extract_phase(1, ai_content),
            'presentation': extract_phase(2, ai_content),
            'evaluation': extract_phase(3, ai_content) or extract_section("Assessment", ai_content),
            'homework': extract_section("Homework", ai_content) or "Refer to textbook exercises.",
            'core_competencies': extract_section("Core Competencies", ai_content), # If you add this field to model later
        }
        messages.success(request, "Lesson plan draft populated from Aura-T session.")

    elif topic:
        # CASE 2: Pre-fill from topic chip (with indicator when available)
        if indicator:
            objectives = (
                f"Learning Indicator: {indicator}\n\n"
                f"By the end of the lesson, students will be able to:\n"
                f"1. Demonstrate understanding of: {indicator}\n"
                f"2. Apply the concept to solve related problems.\n"
                f"3. Connect {topic} to real-world examples."
            )
        else:
            objectives = (
                f"By the end of the lesson, students will be able to:\n"
                f"1. Understand the core concepts of {topic}.\n"
                f"2. Apply {topic} to solve simple problems.\n"
                f"3. Analyze real-world examples of {topic}."
            )
        initial_data = {
            'topic': topic,
            'objectives': objectives,
            'introduction': f"Begin with a 5-minute warm-up activity related to {topic}. Ask students what they already know about it to gauge prior knowledge.",
            'presentation': f"1. Define {topic} and key terminology.\n2. Demonstrate the main concept using visual aids.\n3. Walk through 2-3 step-by-step examples on the board.\n4. Facilitate a class discussion to check for understanding.",
            'evaluation': f"Distribute a short worksheet with 5 problems on {topic}. Circulate to provide individual assistance.",
            'homework': f"Read the chapter on {topic} and complete exercises 1-10 on page 42.",
            'teaching_materials': f"Textbook, Whiteboard, Marker, Projector (optional), Handouts on {topic}"
        }
        if indicator:
            messages.info(request, f"Lesson plan pre-filled with indicator: {indicator[:80]}{'…' if len(indicator) > 80 else ''}")
        else:
            messages.info(request, f"AI has drafted a lesson plan block for '{topic}'. Please review and edit.")

    # Pre-select subject / class when redirected from Scheme of Work
    if prefill_subject_id and not initial_data.get('subject'):
        initial_data['subject'] = prefill_subject_id
    if prefill_class_id and not initial_data.get('school_class'):
        initial_data['school_class'] = prefill_class_id

    if request.method == 'POST' and 'topic' in request.POST: # Standard form submit
        form = LessonPlanForm(request.POST, teacher=teacher)
        if form.is_valid():
            lesson_plan = form.save(commit=False)
            lesson_plan.teacher = teacher
            lesson_plan.save()
            b7_meta_raw = request.POST.get('b7_meta', '')
            indicator_prompt = (request.POST.get('indicator_prompt') or '').strip()
            if b7_meta_raw:
                try:
                    import json as _json
                    lesson_plan.b7_meta = _json.loads(b7_meta_raw)
                    if indicator_prompt:
                        lesson_plan.b7_meta['indicator'] = indicator_prompt
                    lesson_plan.save(update_fields=['b7_meta'])
                except (ValueError, TypeError):
                    pass
            elif indicator_prompt:
                lesson_plan.b7_meta = {'indicator': indicator_prompt}
                lesson_plan.save(update_fields=['b7_meta'])
            messages.success(request, 'Lesson plan created successfully.')
            return redirect('teachers:lesson_plan_list')
    else:
        form = LessonPlanForm(teacher=teacher, initial=initial_data)
        
    # Freemium usage info for lesson generation
    from teachers.addon_utils import get_free_generation_count, FREE_GENERATION_LIMIT, has_addon
    lesson_gen_used = get_free_generation_count(request.user, 'lesson_gen') if request.user.user_type == 'teacher' else 0
    has_planner_pro = has_addon(request.user, 'smart-planner-pro') if request.user.user_type == 'teacher' else True
    lesson_gen_remaining = max(0, FREE_GENERATION_LIMIT - lesson_gen_used)

    return render(request, 'teachers/lesson_plan_form.html', {
        'form': form,
        'title': 'Create Lesson Plan',
        'prefill_indicator': indicator,
        'auto_ges': request.GET.get('auto_ges') == '1',
        'lesson_gen_used': lesson_gen_used,
        'lesson_gen_limit': FREE_GENERATION_LIMIT,
        'lesson_gen_remaining': lesson_gen_remaining,
        'has_planner_pro': has_planner_pro,
    })


@login_required
def lesson_plan_edit(request, pk):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = get_object_or_404(Teacher, user=request.user)
    lesson_plan = get_object_or_404(LessonPlan, pk=pk, teacher=teacher)
    
    if request.method == 'POST':
        form = LessonPlanForm(request.POST, instance=lesson_plan, teacher=teacher)
        if form.is_valid():
            plan = form.save()
            # Save B7 extra meta fields if provided (period, strand, hidden_rows, etc.)
            b7_meta_raw = request.POST.get('b7_meta', '')
            if b7_meta_raw:
                try:
                    import json as _json
                    plan.b7_meta = _json.loads(b7_meta_raw)
                    plan.save(update_fields=['b7_meta'])
                except (ValueError, TypeError):
                    pass
            messages.success(request, 'Lesson plan updated successfully.')
            
            # Check for next URL to return to same view (e.g. from inline edit)
            from django.utils.http import url_has_allowed_host_and_scheme
            next_url = request.POST.get('next') or request.GET.get('next')
            if next_url and url_has_allowed_host_and_scheme(next_url, allowed_hosts={request.get_host()}):
                return redirect(next_url)
                
            return redirect('teachers:lesson_plan_list')
    else:
        form = LessonPlanForm(instance=lesson_plan, teacher=teacher)

    from teachers.addon_utils import get_free_generation_count, FREE_GENERATION_LIMIT, has_addon
    lesson_gen_used = get_free_generation_count(request.user, 'lesson_gen')
    has_planner_pro = has_addon(request.user, 'smart-planner-pro')
    lesson_gen_remaining = max(0, FREE_GENERATION_LIMIT - lesson_gen_used)

    return render(request, 'teachers/lesson_plan_form.html', {
        'form': form,
        'title': 'Edit Lesson Plan',
        'lesson_gen_used': lesson_gen_used,
        'lesson_gen_limit': FREE_GENERATION_LIMIT,
        'lesson_gen_remaining': lesson_gen_remaining,
        'has_planner_pro': has_planner_pro,
    })

@login_required
def lesson_plan_detail(request, pk):
    is_teacher = request.user.user_type == 'teacher'
    is_admin = request.user.user_type == 'admin'
    if not is_teacher and not is_admin:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    if is_teacher:
        teacher = get_object_or_404(Teacher, user=request.user)
        lesson_plan = get_object_or_404(LessonPlan, pk=pk, teacher=teacher)
    else:
        lesson_plan = get_object_or_404(LessonPlan, pk=pk)
    
    intro = lesson_plan.introduction or ''
    pres  = lesson_plan.presentation or ''
    evl   = lesson_plan.evaluation or ''
    has_pulse       = 'PULSE CHECK' in intro
    has_checkpoints = 'CHECKPOINT QUESTIONS' in pres or 'CQ1:' in pres
    has_sprint      = 'MASTERY SPRINT' in evl or 'MS1:' in evl
    has_insight     = 'TEACHER INSIGHT' in evl
    aura_score      = sum([has_pulse, has_checkpoints, has_sprint, has_insight])

    return render(request, 'teachers/lesson_plan_detail.html', {
        'lesson_plan': lesson_plan,
        'is_admin': is_admin,
        'has_pulse': has_pulse,
        'has_checkpoints': has_checkpoints,
        'has_sprint': has_sprint,
        'has_insight': has_insight,
        'aura_score': aura_score,
    })


@login_required
def lesson_plan_duplicate(request, pk):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    teacher = get_object_or_404(Teacher, user=request.user)
    lesson_plan = get_object_or_404(LessonPlan, pk=pk, teacher=teacher)

    duplicate_plan = LessonPlan.objects.create(
        teacher=teacher,
        subject=lesson_plan.subject,
        school_class=lesson_plan.school_class,
        week_number=lesson_plan.week_number,
        topic=f"{lesson_plan.topic} (Copy)",
        objectives=lesson_plan.objectives,
        teaching_materials=lesson_plan.teaching_materials,
        introduction=lesson_plan.introduction,
        presentation=lesson_plan.presentation,
        evaluation=lesson_plan.evaluation,
        homework=lesson_plan.homework,
    )

    messages.success(request, 'Lesson plan duplicated. You can now edit the copy.')
    return redirect('teachers:lesson_plan_edit', pk=duplicate_plan.pk)


@login_required
def lesson_plan_print(request, pk):
    is_teacher = request.user.user_type == 'teacher'
    is_admin = request.user.user_type == 'admin'
    if not is_teacher and not is_admin:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    if is_teacher:
        teacher = get_object_or_404(Teacher, user=request.user)
        lesson_plan = get_object_or_404(LessonPlan, pk=pk, teacher=teacher)
    else:
        lesson_plan = get_object_or_404(LessonPlan, pk=pk)
    
    # Default to B7 weekly template for a consistent print experience.
    template_format = (request.GET.get('template') or 'b7').strip().lower()
    b7_meta = getattr(lesson_plan, 'b7_meta', None) or {}
    b7_context = {
        'term': b7_meta.get('term', '2'),
        'week_ending': b7_meta.get('week_ending', lesson_plan.date_added.strftime('%Y-%m-%d') if lesson_plan.date_added else ''),
        'period': b7_meta.get('period', '1'),
        'duration': b7_meta.get('duration', '60 Mins'),
        'strand': b7_meta.get('strand', ''),
        'class_size': b7_meta.get('class_size', ''),
        'indicator': b7_meta.get('indicator', 'See Content Standard'),
        'lesson_of': b7_meta.get('lesson_of', '1 of 3'),
        'perf_indicator': b7_meta.get('perf_indicator', 'Learners can relate the lesson to real life situations.'),
        'core_competencies': b7_meta.get('core_competencies', 'CP 5.1, CC 8.1'),
        'references': b7_meta.get('references', f"National {lesson_plan.subject.name} Curriculum"),
        'keywords': b7_meta.get('keywords', lesson_plan.topic),
    }
    
    if template_format in {'b7', 'weekly'}:
        template_name = 'teachers/lesson_plan_print_b7.html'
    else:
        # Default to standard GES detail view with print mode
        template_name = 'teachers/lesson_plan_detail.html'
    
    import json as _json
    return render(request, template_name, {
        'lesson_plan': lesson_plan,
        'print_mode': True,
        'current_template': template_format,
        'b7_meta_json': _json.dumps(getattr(lesson_plan, 'b7_meta', None) or {}),
        'b7': b7_context,
        'can_edit': is_teacher,
        'is_admin': is_admin,
    })


@login_required
@requires_addon('quick-report-writer')
def lesson_plan_pdf(request, pk):
    """Read-only B7 weekly layout optimized for browser Save as PDF."""
    is_teacher = request.user.user_type == 'teacher'
    is_admin = request.user.user_type == 'admin'
    if not is_teacher and not is_admin:
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    if is_teacher:
        teacher = get_object_or_404(Teacher, user=request.user)
        lesson_plan = get_object_or_404(LessonPlan, pk=pk, teacher=teacher)
    else:
        lesson_plan = get_object_or_404(LessonPlan, pk=pk)

    b7_meta = getattr(lesson_plan, 'b7_meta', None) or {}
    b7_context = {
        'term': b7_meta.get('term', '2'),
        'week_ending': b7_meta.get('week_ending', lesson_plan.date_added.strftime('%Y-%m-%d') if lesson_plan.date_added else ''),
        'period': b7_meta.get('period', '1'),
        'duration': b7_meta.get('duration', '60 Mins'),
        'strand': b7_meta.get('strand', ''),
        'class_size': b7_meta.get('class_size', ''),
        'indicator': b7_meta.get('indicator', 'See Content Standard'),
        'lesson_of': b7_meta.get('lesson_of', '1 of 3'),
        'perf_indicator': b7_meta.get('perf_indicator', 'Learners can relate the lesson to real life situations.'),
        'core_competencies': b7_meta.get('core_competencies', 'CP 5.1, CC 8.1'),
        'references': b7_meta.get('references', f"National {lesson_plan.subject.name} Curriculum"),
        'keywords': b7_meta.get('keywords', lesson_plan.topic),
    }

    return render(request, 'teachers/lesson_plan_pdf_b7.html', {
        'lesson_plan': lesson_plan,
        'b7': b7_context,
    })

@login_required
@require_plan('pro', 'enterprise')
@requires_addon_freemium('smart-planner-pro')
def aura_t_api(request):
    """
    Handles AJAX requests for lesson plan generation, differentiation, and assignment creation.
    """
    if request.method == "POST":
        try:
            import json
            from teachers.services.aura_gen_engine import AuraGenEngine
            
            data = json.loads(request.body)
            action = data.get('action')
            
            teacher = get_object_or_404(Teacher, user=request.user)
            
            from tenants.ai_quota import check_and_consume

            if action == 'generate_lesson':
                topic = data.get('topic')
                sub_strand = (data.get('sub_strand') or '').strip()
                indicator = (data.get('indicator') or '').strip()
                class_id = data.get('class_id')
                subject_id = data.get('subject_id')

                if not topic:
                    return JsonResponse({"status": "error", "message": "Strand is required"}, status=400)

                subject_name = "General Studies"
                class_name = "General"
                if subject_id:
                    from academics.models import Subject
                    subject = get_object_or_404(Subject, pk=subject_id)
                    subject_name = subject.name
                if class_id:
                    from academics.models import Class
                    school_class = get_object_or_404(Class, pk=class_id)
                    class_name = school_class.name

                check_and_consume(request.tenant, request.user.id, 'lesson_gen')
                result = AuraGenEngine.generate_lesson_plan(
                    topic=topic,
                    subject=subject_name,
                    grade_level=class_name,
                    sub_strand=sub_strand,
                    indicator=indicator,
                )
                return JsonResponse({"status": "success", "data": result})
            elif action == 'differentiate':
                # Placeholder for differentiation logic
                return JsonResponse({
                    "status": "success", 
                    "content": "Differentiation suggestions ready"
                })

            elif action == 'generate_slides':
                topic = data.get('topic')
                class_id = data.get('class_id')
                subject_id = data.get('subject_id')

                if not topic:
                    return JsonResponse({"status": "error", "message": "Topic is required"}, status=400)

                subject_name = "General Studies"
                class_name = "General"
                if subject_id:
                    from academics.models import Subject
                    subject = get_object_or_404(Subject, pk=subject_id)
                    subject_name = subject.name
                if class_id:
                    from academics.models import Class
                    school_class = get_object_or_404(Class, pk=class_id)
                    class_name = school_class.name

                check_and_consume(request.tenant, request.user.id, 'slide_gen')
                result = AuraGenEngine.generate_slides_outline(topic, subject_name, class_name)
                return JsonResponse({"status": "success", "data": result})

            elif action == 'generate_exercises':
                topic = data.get('topic')
                class_id = data.get('class_id')
                subject_id = data.get('subject_id')

                if not topic:
                    return JsonResponse({"status": "error", "message": "Topic is required"}, status=400)

                subject_name = "General Studies"
                class_name = "General"
                if subject_id:
                    from academics.models import Subject
                    subject = get_object_or_404(Subject, pk=subject_id)
                    subject_name = subject.name
                if class_id:
                    from academics.models import Class
                    school_class = get_object_or_404(Class, pk=class_id)
                    class_name = school_class.name

                check_and_consume(request.tenant, request.user.id, 'exercise_gen')
                result = AuraGenEngine.generate_interactive_exercises(topic, subject_name, class_name)
                return JsonResponse({"status": "success", "data": result})

            elif action == 'generate_exercises_and_assign':
                topic = data.get('topic')
                class_id = data.get('class_id')
                subject_id = data.get('subject_id')
                due_date = data.get('due_date')

                if not topic or not class_id or not subject_id:
                    return JsonResponse({"status": "error", "message": "Topic, class, and subject are required"}, status=400)

                from academics.models import Class, Subject
                from homework.models import Homework, Question, Choice
                from django.utils import timezone
                from datetime import timedelta, datetime

                school_class = get_object_or_404(Class, pk=class_id)
                subject = get_object_or_404(Subject, pk=subject_id)

                check_and_consume(request.tenant, request.user.id, 'exercise_gen')
                result = AuraGenEngine.generate_interactive_exercises(topic, subject.name, school_class.name)

                try:
                    if due_date:
                        due_date_value = datetime.strptime(due_date, '%Y-%m-%d').date()
                    else:
                        due_date_value = timezone.now().date() + timedelta(days=7)
                except ValueError:
                    due_date_value = timezone.now().date() + timedelta(days=7)

                homework = Homework.objects.create(
                    title=f"{topic} - Interactive Exercise",
                    description=f"AI-generated interactive exercise on {topic}.",
                    teacher=teacher,
                    subject=subject,
                    target_class=school_class,
                    due_date=due_date_value
                )

                exercises = result.get('exercises', [])
                if not exercises:
                    exercises = AuraGenEngine.generate_interactive_exercises(topic, subject.name, school_class.name).get('exercises', [])

                if isinstance(exercises, dict):
                    exercises = exercises.get('items', [])
                normalized_exercises = []
                for item in exercises:
                    if isinstance(item, str):
                        normalized_exercises.append({'type': 'short', 'prompt': item, 'answer': '', 'dok_level': 1})
                    elif isinstance(item, dict):
                        if not item.get('prompt') and item.get('question'):
                            item['prompt'] = item.get('question')
                        normalized_exercises.append(item)
                exercises = normalized_exercises

                if len(exercises) < 5:
                    fallback = AuraGenEngine._mock_exercises(topic)
                    exercises = (exercises + fallback)[:5]
                elif len(exercises) > 10:
                    exercises = exercises[:10]

                type_map = {
                    'multiple_choice': 'mcq',
                    'multiple choice': 'mcq',
                    'mcq': 'mcq',
                    'quiz': 'mcq',
                    'short_answer': 'short',
                    'short answer': 'short',
                    'short': 'short',
                }
                questions_created = 0
                for item in exercises:
                    raw_type = (item.get('type') or '').lower()
                    ex_type = type_map.get(raw_type, raw_type or 'short')
                    prompt = item.get('prompt') or 'Untitled question'
                    dok_level = item.get('dok_level')
                    try:
                        dok_level = int(dok_level)
                    except (TypeError, ValueError):
                        dok_level = 1
                    if dok_level not in [1, 2, 3, 4]:
                        dok_level = 1
                    if ex_type not in ['mcq', 'short']:
                        ex_type = 'short'

                    question = Question.objects.create(
                        homework=homework,
                        text=prompt,
                        points=1,
                        question_type=ex_type,
                        dok_level=dok_level,
                        correct_answer=item.get('answer', '')
                    )
                    questions_created += 1

                    if ex_type == 'mcq':
                        options = item.get('options') or []
                        correct = item.get('answer')
                        if not options and correct:
                            options = [
                                str(correct),
                                'Option B',
                                'Option C',
                                'Option D'
                            ]

                        if not options:
                            options = ['Option A', 'Option B', 'Option C', 'Option D']

                        def normalize_text(value):
                            return str(value or '').strip().lower()

                        correct_norm = normalize_text(correct)
                        correct_index = None
                        label_map = {'a': 0, 'b': 1, 'c': 2, 'd': 3}
                        if correct_norm in label_map:
                            correct_index = label_map[correct_norm]
                        elif correct_norm.isdigit():
                            idx = int(correct_norm) - 1
                            if 0 <= idx < len(options):
                                correct_index = idx

                        has_correct = False
                        for idx, opt in enumerate(options):
                            opt_norm = normalize_text(opt)
                            is_correct = False

                            if correct_index is not None and idx == correct_index:
                                is_correct = True
                            elif correct_norm and opt_norm == correct_norm:
                                is_correct = True

                            if is_correct:
                                has_correct = True

                            Choice.objects.create(
                                question=question,
                                text=opt,
                                is_correct=is_correct
                            )

                        if not has_correct:
                            first_choice = question.choices.first()
                            if first_choice:
                                first_choice.is_correct = True
                                first_choice.save(update_fields=['is_correct'])

                return JsonResponse({
                    "status": "success",
                    "data": {
                        "homework_id": homework.id,
                        "homework_title": homework.title,
                        "questions_created": questions_created
                    }
                })
            
            elif action == 'generate_assignment':
                lesson_id = data.get('lesson_id')
                topic = data.get('topic')

                check_and_consume(request.tenant, request.user.id, 'assignment_gen')
                if lesson_id:
                    lesson_plan = get_object_or_404(LessonPlan, pk=lesson_id, teacher=teacher)
                    result = AuraGenEngine.generate_assignment_package(lesson_plan)
                elif topic and data.get('class_id'):
                    # Create a temporary lesson object for the engine if topic is provided manually
                    # This is a simplified path
                    from academics.models import Class, Subject
                    school_class = get_object_or_404(Class, pk=data.get('class_id'))
                    subject = get_object_or_404(Subject, pk=data.get('subject_id')) # Assuming subject passed
                    
                    # Mock lesson plan for the engine
                    mock_lesson = LessonPlan(
                        teacher=teacher,
                        subject=subject,
                        school_class=school_class,
                        topic=topic,
                        objectives=f"Understand {topic}",
                        week_number=1 # Dummy
                    )
                    result = AuraGenEngine.generate_assignment_package(mock_lesson, topic_prompt=topic)
                else:
                    return JsonResponse({"status": "error", "message": "Missing lesson_id or topic/class_id"}, status=400)
                
                return JsonResponse({
                    "status": "success", 
                    "data": result
                })

            elif action == 'save_assignment_package':
                from homework.models import Homework as HW
                from academics.models import Class, Subject
                from datetime import datetime, timedelta
                from django.utils import timezone as tz

                title        = data.get('title', '').strip() or 'AI Generated Assignment'
                description  = data.get('description', '').strip() or ''
                class_id     = data.get('class_id')
                subject_id   = data.get('subject_id')
                due_date_str = data.get('due_date', '')

                if not class_id or not subject_id:
                    return JsonResponse({"status": "error", "message": "Class and subject are required."}, status=400)

                school_class = get_object_or_404(Class, pk=class_id)
                subject      = get_object_or_404(Subject, pk=subject_id)

                try:
                    due_date = datetime.strptime(due_date_str, '%Y-%m-%d').date()
                except (ValueError, TypeError):
                    due_date = tz.now().date() + timedelta(days=7)

                hw = HW.objects.create(
                    title=title,
                    description=description,
                    teacher=teacher,
                    subject=subject,
                    target_class=school_class,
                    due_date=due_date,
                )

                # Persist AI-generated questions as Question objects
                from homework.models import Question as HWQuestion
                questions_payload = data.get('questions', [])
                if isinstance(questions_payload, list):
                    for i, q_text in enumerate(questions_payload):
                        if isinstance(q_text, str) and q_text.strip():
                            HWQuestion.objects.create(
                                homework=hw,
                                text=q_text.strip(),
                                question_type='short',
                                dok_level=1,
                                points=1,
                            )

                schema = request.tenant.schema_name
                return JsonResponse({
                    "status": "success",
                    "homework_id": hw.pk,
                    "redirect_url": f"/{schema}/homework/{hw.pk}/",
                })

            elif action == 'regenerate':
                import re as _re
                plan_id = data.get('plan_id')
                if not plan_id:
                    return JsonResponse({"status": "error", "message": "plan_id required"}, status=400)

                regen_plan = get_object_or_404(LessonPlan, pk=plan_id, teacher=teacher)
                topic_r = regen_plan.topic
                subject_r = regen_plan.subject.name if regen_plan.subject else "General Studies"
                class_r = regen_plan.school_class.name if regen_plan.school_class else "General"

                check_and_consume(request.tenant, request.user.id, 'lesson_gen')
                regen_result = AuraGenEngine.generate_lesson_plan(topic_r, subject_r, class_r)
                lesson_body = (regen_result.get('lesson_plan') or '').strip()
                if not lesson_body:
                    return JsonResponse({"status": "error",
                                         "message": regen_result.get('error', 'Generation failed')}, status=500)

                def _extract_r(header, text):
                    pat = _re.compile(
                        rf"\*\*{_re.escape(header)}.*?\*\*\s*(.*?)(?=\n\*\*|\Z)",
                        _re.DOTALL | _re.IGNORECASE
                    )
                    m = pat.search(text)
                    return m.group(1).strip() if m else ''

                def _extract_phase(phase_num, text):
                    """Extract full phase content up to the next phase-level boundary."""
                    boundaries = {
                        1: r'\*\*PHASE\s*2',
                        2: r'\*\*PHASE\s*3',
                        3: r'\*\*(?:Resources|Homework)',
                    }
                    end = boundaries.get(phase_num, r'\Z')
                    pat = _re.compile(
                        rf"\*\*PHASE\s*{phase_num}[^*]*\*\*\s*(.*?)(?={end}|\Z)",
                        _re.DOTALL | _re.IGNORECASE
                    )
                    m = pat.search(text)
                    return m.group(1).strip() if m else ''

                regen_plan.objectives = (
                    _extract_r('Content Standard', lesson_body) + '\n' +
                    _extract_r('Indicator', lesson_body)
                ).strip() or regen_plan.objectives
                regen_plan.introduction = _extract_phase(1, lesson_body) or regen_plan.introduction
                regen_plan.presentation = _extract_phase(2, lesson_body) or regen_plan.presentation
                regen_plan.evaluation = (
                    _extract_phase(3, lesson_body) or
                    _extract_r('Assessment', lesson_body) or
                    regen_plan.evaluation
                )
                hw_new = _extract_r('Homework', lesson_body)
                if hw_new:
                    regen_plan.homework = hw_new
                regen_plan.save()

                return JsonResponse({
                    "status": "success",
                    "message": "Lesson plan upgraded to Aura-T v3.",
                    "plan_id": regen_plan.pk,
                })

            return JsonResponse({"status": "error", "message": "Unknown action"}, status=400)
            
        except Exception as e:
            return _ai_json_error_response(e)
            
    return JsonResponse({"status": "error", "message": "Invalid method"}, status=405)


@login_required
@requires_addon_freemium('smart-planner-pro')
def ges_lesson_api(request):
    """
    Separate GES Weekly Notes generator endpoint.
    This flow does not use Aura-T formatting/parsing.
    """
    if request.method != 'POST':
        return JsonResponse({"status": "error", "message": "Invalid method"}, status=405)

    if request.user.user_type != 'teacher':
        return JsonResponse({"status": "error", "message": "Access denied"}, status=403)

    try:
        import json
        from teachers.services.ges_lesson_engine import GESLessonEngine
        from tenants.ai_quota import check_and_consume

        data = json.loads(request.body)
        topic = (data.get('topic') or '').strip()
        indicator = (data.get('indicator') or '').strip()
        sub_strand = (data.get('sub_strand') or '').strip()
        section = (data.get('section') or '').strip()
        current_text = (data.get('current_text') or '').strip()
        class_id = data.get('class_id')
        subject_id = data.get('subject_id')
        week_number = data.get('week_number')

        if not topic:
            return JsonResponse({"status": "error", "message": "Topic is required"}, status=400)
        if not indicator:
            return JsonResponse({"status": "error", "message": "Target indicator is required"}, status=400)

        try:
            week_number = int(week_number or 1)
        except (TypeError, ValueError):
            week_number = 1
        if week_number < 1:
            week_number = 1

        subject_name = 'General Studies'
        class_name = 'General'
        if subject_id:
            from academics.models import Subject
            subject = get_object_or_404(Subject, pk=subject_id)
            subject_name = subject.name
        if class_id:
            from academics.models import Class
            school_class = get_object_or_404(Class, pk=class_id)
            class_name = school_class.name

        check_and_consume(request.tenant, request.user.id, 'lesson_gen')
        if section:
            result = GESLessonEngine.generate_section(
                topic=topic,
                indicator=indicator,
                subject=subject_name,
                grade_level=class_name,
                week_number=week_number,
                section=section,
                current_text=current_text,
                sub_strand=sub_strand,
            )
            return JsonResponse({"status": "success", "data": result})

        result = GESLessonEngine.generate_weekly_notes(
            topic=topic,
            indicator=indicator,
            subject=subject_name,
            grade_level=class_name,
            week_number=week_number,
            sub_strand=sub_strand,
        )
        return JsonResponse({"status": "success", "data": result})
    except Exception as e:
        return _ai_json_error_response(e)

@login_required
@require_plan('basic', 'pro', 'enterprise')
def aura_command_center(request):
    """
    Aura-T Command Center - teachers go to the unified AI sessions page,
    admins get the full multi-teacher plan overview.
    """
    if request.user.user_type == 'teacher':
        return redirect('teachers:ai_sessions_list')

    is_admin = request.user.user_type == 'admin'
    if not is_admin:
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    teacher = None
    # Admin-only path: show all teachers' plans

    plans_qs = LessonPlan.objects.select_related('subject', 'school_class', 'teacher', 'teacher__user')
    plans = list(plans_qs)

    for p in plans:
        intro = p.introduction or ''
        pres  = p.presentation or ''
        evl   = p.evaluation or ''
        p.has_pulse       = 'PULSE CHECK' in intro
        p.has_checkpoints = 'CHECKPOINT QUESTIONS' in pres or 'CQ1:' in pres
        p.has_sprint      = 'MASTERY SPRINT' in evl or 'MS1:' in evl
        p.has_insight     = 'TEACHER INSIGHT' in evl
        p.aura_score      = sum([p.has_pulse, p.has_checkpoints, p.has_sprint, p.has_insight])
        p.teacher_name    = p.teacher.user.get_full_name() if p.teacher and p.teacher.user else 'Unknown'
        p.pulse_count     = 0
        p.has_active_pulse = False

    # ── Per-plan pulse session counts & active detection ──────────────────
    try:
        from academics.pulse_models import PulseSession
        plan_pks = [p.pk for p in plans]
        pulse_counts  = {}
        active_pk_set = set()
        for row in PulseSession.objects.filter(lesson_plan_id__in=plan_pks).values('lesson_plan_id', 'status'):
            lpk = row['lesson_plan_id']
            pulse_counts[lpk] = pulse_counts.get(lpk, 0) + 1
            if row['status'] == 'active':
                active_pk_set.add(lpk)
        for p in plans:
            p.pulse_count      = pulse_counts.get(p.pk, 0)
            p.has_active_pulse = p.pk in active_pk_set
    except Exception:
        pass

    total         = len(plans)
    v3_ready      = sum(1 for p in plans if p.aura_score == 4)
    partial_count = sum(1 for p in plans if 0 < p.aura_score < 4)
    no_aura_count = sum(1 for p in plans if p.aura_score == 0)
    v3_pct        = round(v3_ready / total * 100) if total else 0

    subjects = sorted(set(p.subject.name for p in plans if p.subject))
    classes  = sorted(set(p.school_class.name for p in plans if p.school_class))

    # Admin: sort by teacher name then score
    plans.sort(key=lambda p: (p.teacher_name, -p.aura_score, -p.pk))
    prev_tname = None
    for p in plans:
        p.teacher_changed = p.teacher_name != prev_tname
        prev_tname = p.teacher_name

    return render(request, 'teachers/aura_command_center.html', {
        'plans':         plans,
        'total':         total,
        'v3_ready':      v3_ready,
        'v3_pct':        v3_pct,
        'partial_count': partial_count,
        'no_aura_count': no_aura_count,
        'is_teacher':    False,
        'is_admin':      True,
        'teacher':       None,
        'subjects':      subjects,
        'classes':       classes,
    })


@login_required
@require_plan('basic', 'pro', 'enterprise')
def aura_flight_manual(request):
    """
    Aura-T Teacher Flight Manual - quick-start guide for the AI-integrated lesson.
    """
    is_teacher = request.user.user_type == 'teacher'
    is_admin   = request.user.user_type == 'admin'
    if not is_teacher and not is_admin:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    return render(request, 'teachers/aura_flight_manual.html', {
        'is_teacher': is_teacher,
        'is_admin': is_admin,
    })


@login_required
@requires_addon('quick-report-writer')
def lesson_plan_cards_print(request, pk):
    """
    Printable Student Nuggets + Mastery Sprint cards for a lesson plan.
    Cut-and-distribute format: support card | extension card, then exit ticket strip.
    """
    is_teacher = request.user.user_type == 'teacher'
    is_admin   = request.user.user_type == 'admin'
    if not is_teacher and not is_admin:
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    if is_teacher:
        teacher = get_object_or_404(Teacher, user=request.user)
        lesson_plan = get_object_or_404(LessonPlan, pk=pk, teacher=teacher)
    else:
        lesson_plan = get_object_or_404(LessonPlan, pk=pk)

    return render(request, 'teachers/lesson_plan_cards_print.html', {
        'lesson_plan': lesson_plan,
    })


@login_required
def lesson_plan_delete(request, pk):
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = get_object_or_404(Teacher, user=request.user)
    lesson_plan = get_object_or_404(LessonPlan, pk=pk, teacher=teacher)
    
    if request.method == 'POST':
        lesson_plan.delete()
        messages.success(request, 'Lesson plan deleted successfully.')
        return redirect('teachers:lesson_plan_list')
        
    return render(request, 'teachers/lesson_plan_confirm_delete.html', {'lesson_plan': lesson_plan})

@login_required
def add_teacher(request):
    if request.user.user_type != 'admin':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
        
    if request.method == 'POST':
        form = TeacherCreateForm(request.POST)
        if form.is_valid():
            employee_id = form.cleaned_data['employee_id']
            if not employee_id:
                # Auto-generate Employee ID (TCHR + 4 Digits)
                import random
                import string
                while True:
                    suffix = ''.join(random.choices(string.digits, k=4))
                    employee_id = f"TCHR{suffix}"
                    if not Teacher.objects.filter(employee_id=employee_id).exists() and not User.objects.filter(username=employee_id).exists():
                        break

            username = employee_id
            if User.objects.filter(username=username).exists():
                messages.error(request, f"User {username} already exists")
                return render(request, 'teachers/add_teacher.html', {'form': form})
                
            user = User.objects.create_user(
                username=username,
                email=form.cleaned_data['email'],
                password=username,  # Default password is employee ID
                first_name=form.cleaned_data['first_name'],
                last_name=form.cleaned_data['last_name'],
                user_type='teacher',
                phone=form.cleaned_data.get('phone', '')
            )
            
            teacher = form.save(commit=False)
            teacher.user = user
            teacher.employee_id = employee_id # Ensure it's set
            teacher.save()
            form.save_m2m() # Save subjects
            messages.success(request, f"Teacher {user.get_full_name()} added successfully with Employee ID: {employee_id}")
            return redirect('teachers:teacher_list')
    else:
        form = TeacherCreateForm()
        
    return render(request, 'teachers/add_teacher.html', {'form': form})


# =====================
# ID CARD GENERATION
# =====================

@login_required
def teacher_id_card(request, teacher_id):
    """Generate and download teacher ID card (PNG)"""
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    # Permission check
    if request.user.user_type == 'teacher':
        teacher_profile = get_object_or_404(Teacher, user=request.user)
        if teacher_profile.id != teacher_id:
            messages.error(request, 'You can only download your own ID card')
            return redirect('dashboard')
    elif request.user.user_type not in ['admin']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    # Generate card
    try:
        card = generate_teacher_id_card(teacher)
        
        # Return as PNG
        response = HttpResponse(content_type='image/png')
        response['Content-Disposition'] = f'attachment; filename="teacher_id_{teacher.id}.png"'
        card.save(response, 'PNG')
        return response
    except Exception as e:
        messages.error(request, f'Error generating ID card: {str(e)}')
        return redirect('dashboard')


@login_required
def teacher_id_card_pdf(request, teacher_id):
    """Generate and download teacher ID card (PDF)"""
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    # Permission check
    if request.user.user_type == 'teacher':
        teacher_profile = get_object_or_404(Teacher, user=request.user)
        if teacher_profile.id != teacher_id:
            messages.error(request, 'You can only download your own ID card')
            return redirect('dashboard')
    elif request.user.user_type not in ['admin']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    # Generate card and convert to PDF
    try:
        card = generate_teacher_id_card(teacher)
        pdf_buffer = export_id_card_to_pdf(card)
        
        response = HttpResponse(pdf_buffer, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="teacher_id_{teacher.id}.pdf"'
        return response
    except Exception as e:
        messages.error(request, f'Error generating ID card: {str(e)}')
        return redirect('dashboard')


@login_required
def bulk_teacher_id_cards_pdf(request):
    """Generate bulk ID cards for all teachers (PDF)"""
    if request.user.user_type not in ['admin']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    # Get parameters
    teacher_ids = request.GET.get('ids', '').split(',')
    teacher_ids = [tid.strip() for tid in teacher_ids if tid.strip()]
    
    try:
        if teacher_ids:
            teachers = Teacher.objects.filter(id__in=teacher_ids).select_related('user')
        else:
            teachers = Teacher.objects.all().select_related('user')
        
        if not teachers.exists():
            messages.error(request, 'No teachers found')
            return redirect('teachers:teacher_list')
        
        # Generate cards
        card_list = []
        for teacher in teachers:
            try:
                card = generate_teacher_id_card(teacher)
                card_list.append((teacher.user.get_full_name(), card))
            except Exception:
                pass  # Skip teachers with errors
        
        if not card_list:
            messages.error(request, 'Could not generate any ID cards')
            return redirect('teachers:teacher_list')
        
        # Export to PDF
        pdf_buffer = export_multiple_id_cards_to_pdf(card_list)
        
        response = HttpResponse(pdf_buffer, content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="teacher_ids_all.pdf"'
        return response
        
    except Exception as e:
        messages.error(request, f'Error generating ID cards: {str(e)}')
        return redirect('teachers:teacher_list')


import json
from django.views.decorators.http import require_POST
from communication.models import Conversation, Message
from parents.models import Parent


# --- Power Words - Teacher Command Center ---

@login_required
def power_words_dashboard(request):
    """
    Teacher Command Center: per-student Power Word clouds, weekly growth,
    Aura-T academic verb trend insights, and one-click action buttons.

    Access: teacher (own classes) or admin (all classes).
    """
    if request.user.user_type not in ['teacher', 'admin']:
        messages.error(request, 'Access denied.')
        return redirect('dashboard')

    from collections import defaultdict
    from academics.tutor_models import PowerWord
    from django.utils import timezone as tz

    teacher = None
    if request.user.user_type == 'teacher':
        teacher = Teacher.objects.filter(user=request.user).first()

    # ── Select class ────────────────────────────────────────────────
    if teacher:
        classes = Class.objects.filter(classsubject__teacher=teacher).distinct().order_by('name')
    else:
        classes = Class.objects.all().order_by('name')

    class_id = request.GET.get('class_id')
    selected_class = None
    if class_id:
        selected_class = get_object_or_404(Class, id=class_id)
    elif classes.exists():
        selected_class = classes.first()

    # ── Derive week numbers for "this week" and "last week" ──────────
    now = tz.now()
    iso = now.isocalendar()
    current_year, current_week = iso[0], iso[1]
    last_week = current_week - 1 if current_week > 1 else 52
    last_week_year = current_year if current_week > 1 else current_year - 1

    # ── Academic verbs to detect in Power Words ──────────────────────
    ACADEMIC_VERBS = {
        'analyze', 'analyse', 'compare', 'evaluate', 'synthesize', 'synthesise',
        'describe', 'explain', 'justify', 'identify', 'classify', 'predict',
        'infer', 'conclude', 'hypothesize', 'hypothesise', 'illustrate',
        'summarize', 'summarise', 'argue', 'debate', 'interpret', 'calculate',
        'demonstrate', 'apply', 'construct', 'define', 'distinguish',
    }

    student_cards = []

    if selected_class:
        students = Student.objects.filter(
            current_class=selected_class
        ).select_related('user').order_by('user__last_name', 'user__first_name')

        # Bulk-fetch Power Words for all students in the class
        all_words = PowerWord.objects.filter(
            student__in=students
        ).select_related('student').order_by('-last_heard')

        # Group by student
        words_by_student = defaultdict(list)
        for pw in all_words:
            words_by_student[pw.student_id].append(pw)

        for student in students:
            student_words = words_by_student.get(student.id, [])

            # This-week and last-week splits
            this_week_words = [
                w for w in student_words
                if w.year == current_year and w.week == current_week
            ]
            last_week_words = [
                w for w in student_words
                if w.year == last_week_year and w.week == last_week
            ]

            # Academic verb percentage (this week vs last week)
            def verb_pct(word_list):
                if not word_list:
                    return 0
                verbs = sum(1 for w in word_list if w.word.lower() in ACADEMIC_VERBS)
                return round(verbs / len(word_list) * 100)

            verb_pct_this = verb_pct(this_week_words)
            verb_pct_last = verb_pct(last_week_words)
            verb_delta = verb_pct_this - verb_pct_last

            # Aura-T insight string
            aura_insight = None
            all_historical_verbs = sum(1 for w in student_words if w.word.lower() in ACADEMIC_VERBS)
            total = len(student_words)
            if total >= 3:
                overall_verb_pct = round(all_historical_verbs / total * 100)
                if verb_delta > 10:
                    aura_insight = (
                        f"{student.user.first_name} is using academic verbs (Analyze, Compare) "
                        f"{verb_delta}% more often than last week."
                    )
                elif verb_delta < -10:
                    aura_insight = (
                        f"{student.user.first_name}'s academic verb usage dropped {abs(verb_delta)}% "
                        "this week — may need a vocabulary review."
                    )
                elif overall_verb_pct >= 30:
                    aura_insight = (
                        f"{student.user.first_name} consistently uses academic verbs in {overall_verb_pct}% "
                        "of Power Words — strong academic language command!"
                    )

            # Word cloud: (word, font_size) pairs sized by used_count
            max_count = max((w.used_count for w in student_words), default=1)
            word_cloud = [
                {
                    'word': w.word.title(),
                    'used_count': w.used_count,
                    'subject': w.subject,
                    'font_size': max(12, min(32, int(12 + (w.used_count / max_count) * 20))),
                    'is_verb': w.word.lower() in ACADEMIC_VERBS,
                    'session_type': w.session_type,
                }
                for w in student_words[:30]  # cap at 30 for display
            ]

            student_cards.append({
                'student': student,
                'total_words': total,
                'new_this_week': len(this_week_words),
                'new_this_week_list': [w.word.title() for w in this_week_words[:8]],
                'verb_pct_this': verb_pct_this,
                'verb_delta': verb_delta,
                'aura_insight': aura_insight,
                'word_cloud': word_cloud,
            })

    context = {
        'classes': classes,
        'selected_class': selected_class,
        'student_cards': student_cards,
        'current_week': current_week,
        'current_year': current_year,
    }
    return render(request, 'teachers/power_words_dashboard.html', context)


@require_POST
@login_required
def power_words_action(request):
    """
    Handle teacher actions on the Power Words dashboard:
      - send_congratulations: Creates an in-app notification for the student.
      - add_to_report:         Flags the student's top Power Words for term report.
    """
    if request.user.user_type not in ['teacher', 'admin']:
        return JsonResponse({'error': 'Access denied'}, status=403)

    try:
        body = json.loads(request.body.decode('utf-8'))
    except Exception:
        return JsonResponse({'error': 'Invalid JSON'}, status=400)

    action = body.get('action')
    student_id = body.get('student_id')

    if not action or not student_id:
        return JsonResponse({'error': 'action and student_id required'}, status=400)

    try:
        student = Student.objects.select_related('user').get(id=student_id)
    except Student.DoesNotExist:
        return JsonResponse({'error': 'Student not found'}, status=404)

    from academics.tutor_models import PowerWord
    from django.utils import timezone as tz
    from announcements.models import Notification

    if action == 'send_congratulations':
        # Build uplifting note based on top 3 recent words
        recent_words = PowerWord.objects.filter(student=student).order_by('-last_heard')[:3]
        word_list = ', '.join(w.word.title() for w in recent_words) if recent_words else 'your Power Words'
        teacher_name = request.user.get_full_name() or request.user.username
        msg = (
            f"🌟 Great work, {student.user.first_name}! "
            f"{teacher_name} noticed you've been mastering academic vocabulary: "
            f"{word_list}. Keep it up!"
        )
        try:
            Notification.objects.create(
                recipient=student.user,
                message=msg,
                alert_type='general',
            )
            return JsonResponse({'status': 'ok', 'message': 'Congratulations note sent!'})
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    elif action == 'add_to_report':
        # Mark recent Power Words as confirmed by teacher so they appear in reports
        try:
            from django.utils import timezone as tz
            iso = tz.now().isocalendar()
            updated = PowerWord.objects.filter(
                student=student,
                year=iso[0],
                week=iso[1],
            ).update(confirmed_by_teacher=True)
            return JsonResponse({
                'status': 'ok',
                'message': f"{updated} Power Word(s) flagged for term report.",
                'count': updated,
            })
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)

    return JsonResponse({'error': f"Unknown action: {action}"}, status=400)


@require_POST
@login_required
@requires_addon('grade-insight-dashboard')
def boost_intervention(request):
    """
    AJAX endpoint to trigger an AI intervention 'Boost'.
    1. Creates a remedial assignment (ClassExercise)
    2. Notifies the parent via internal message system
    """
    if request.user.user_type != 'teacher':
        return JsonResponse({'status': 'error', 'message': 'Access denied'}, status=403)
        
    try:
        data = json.loads(request.body)
        student_id = data.get('student_id')
        intervention_text = data.get('interventionText')

        student = get_object_or_404(Student, id=student_id)
        # Note: request.user might not be Teacher instance directly, but user linked to Teacher
        try:
            teacher = Teacher.objects.get(user=request.user)
        except Teacher.DoesNotExist:
             return JsonResponse({'status': 'error', 'message': 'Teacher profile not found'}, status=404)
        
        # 1. Assign Remedial Task
        # Find a subject this teacher teaches to this student
        # We look for a ClassSubject where teacher=teacher and class_name=student.current_class
        class_subject = ClassSubject.objects.filter(
            teacher=teacher, 
            class_name=student.current_class
        ).first()

        log_message = []
        
        if class_subject:
            # Create a specialized ClassExercise
            exercise = ClassExercise.objects.create(
                class_subject=class_subject,
                title=f"Boost: {intervention_text[:20]}...",
                description=f"AI-suggested intervention: {intervention_text}",
                max_marks=10,
                term='first' # Defaulting to first term
            )
            # Assign score entry (initializes it for student)
            StudentExerciseScore.objects.create(
                student=student,
                exercise=exercise,
                score=0,
                remarks="Assigned via AI Boost"
            )
            log_message.append(f"Assigned '{exercise.title}' in {class_subject.subject.name}")
        else:
            log_message.append("No class subject found to assign work")

        # 2. Notify Parent
        # Student -> Parents (ManyToMany)
        # We pick the first one
        parent = student.parents.first() # Using related_name='parents' from Parent model? 
        # Wait, Parent model has: children = models.ManyToManyField('students.Student', related_name='parents')
        # So yes, student.parents.all() returns QuerySet of Parent objects.
        
        if parent:
            parent_user = parent.user
            # Create or get conversation
            # Using specific method from Conversation model or get_or_create
            conversation, created = Conversation.get_or_create_between(request.user, parent_user)
            
            msg_content = (
                f"Dear {parent_user.last_name}, we noticed {student.user.first_name} could use a little extra help with "
                f"'{intervention_text}'. I've assigned a short practice module (Boost) to support them. "
                f"Please encourage them to complete it effectively. Thanks!"
            )
            
            Message.objects.create(
                conversation=conversation,
                sender=request.user,
                content=msg_content
            )
            log_message.append(f"Notified {parent_user.get_full_name()}")
        else:
            log_message.append("No parent linked for notification")

        return JsonResponse({
            'status': 'success', 
            'message': f"Intervention deployed! {'. '.join(log_message)}."
        })
        
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=500)


# ==========================================
# AI Chat Sessions
# ==========================================

from django.http import JsonResponse
from django.conf import settings
import json

@login_required
@login_required
@require_addon('ai-lesson-planner')
def ai_sessions_list(request):
    """
    Aura-T Command Centre (combined):
    AI Copilot sessions, student matrix, lesson plan upgrade tracker,
    and live Pulse monitoring — all in one bento-grid interface.
    """
    if not hasattr(request.user, 'teacher'):
         messages.error(request, "Access restricted to teachers.")
         return redirect('dashboard')
         
    teacher = request.user.teacher
    sessions = LessonGenerationSession.objects.filter(teacher=teacher).order_by('-updated_at')
    
    # --- REDESIGN: REAL DATA INJECTION ---
    now = timezone.localtime()
    weekday = now.weekday()  # 0=Monday, 6=Sunday
    current_time = now.time()
    
    # 1. Fetch Next Active Class
    # Query: Find classes for this teacher today that haven't ended yet
    next_class_session = Timetable.objects.filter(
        class_subject__teacher=teacher,
        day=weekday,
        end_time__gte=current_time
    ).select_related('class_subject__class_name', 'class_subject__subject').order_by('start_time').first()
    
    # 2. Determine Scope Context (Which class to show in matrix?)
    active_class = None
    if next_class_session:
        active_class = next_class_session.class_subject.class_name
    elif teacher.managed_classes.exists():
        # Fallback to the first class they manage (e.g. homeroom)
        active_class = teacher.managed_classes.first()
    else:
        # Fallback to any class they teach
        first_subject = ClassSubject.objects.filter(teacher=teacher).select_related('class_name').first()
        if first_subject:
            active_class = first_subject.class_name
            
    # 3. Build Student Matrix Data
    student_data = []
    struggling_count = 0
    total_students = 0

    if active_class:
        # Get up to 40 students for the visualization grid
        students = list(Student.objects.filter(current_class=active_class).select_related('user')[:40])
        student_ids = [s.id for s in students]
        total_students = len(students)

        # Real engagement: query today's TutorSession activity
        from datetime import date as _date
        today = _date.today()
        active_today_ids = set(
            TutorSession.objects.filter(
                student_id__in=student_ids,
                started_at__date=today
            ).values_list('student_id', flat=True)
        )
        # Recent sessions (last 3 days) for "idle" — logged in but not today
        recent_ids = set(
            TutorSession.objects.filter(
                student_id__in=student_ids,
                started_at__date__gte=today - timezone.timedelta(days=3)
            ).values_list('student_id', flat=True)
        ) - active_today_ids

        # Struggling: students with last session showing low engagement (msg_count < 4)
        struggling_set = set(
            TutorSession.objects.filter(
                student_id__in=student_ids,
                message_count__lt=4,
                started_at__date=today
            ).values_list('student_id', flat=True)
        )

        # Latest session ID per student (for "peek" link)
        latest_sessions = {}
        for ts in TutorSession.objects.filter(
            student_id__in=student_ids
        ).order_by('student_id', '-started_at').distinct('student_id').values('student_id', 'id'):
            latest_sessions[ts['student_id']] = ts['id']

        for s in students:
            if s.id in struggling_set:
                status = 'struggling'
                struggling_count += 1
            elif s.id in active_today_ids:
                status = 'active'
            elif s.id in recent_ids:
                status = 'idle'
            else:
                status = 'offline'

            student_data.append({
                'id': s.id,
                'name': s.user.get_full_name() or s.user.username,
                'status': status,
                'tooltip': f"{s.admission_number} | {status.title()}",
                'latest_session_id': latest_sessions.get(s.id),
            })

    # ── Lesson Plan Upgrade Data (Command Center) ─────────────────────────
    plans_qs = LessonPlan.objects.filter(teacher=teacher).select_related(
        'subject', 'school_class'
    )
    plans = list(plans_qs)

    for p in plans:
        intro = p.introduction or ''
        pres  = p.presentation or ''
        evl   = p.evaluation or ''
        p.has_pulse       = 'PULSE CHECK' in intro
        p.has_checkpoints = 'CHECKPOINT QUESTIONS' in pres or 'CQ1:' in pres
        p.has_sprint      = 'MASTERY SPRINT' in evl or 'MS1:' in evl
        p.has_insight     = 'TEACHER INSIGHT' in evl
        p.aura_score      = sum([p.has_pulse, p.has_checkpoints, p.has_sprint, p.has_insight])
        p.pulse_count     = 0
        p.has_active_pulse = False

    # Bulk-query pulse session counts + active state
    try:
        from academics.pulse_models import PulseSession as PS
        plan_pks = [p.pk for p in plans]
        pulse_counts  = {}
        active_pk_set = set()
        for row in PS.objects.filter(lesson_plan_id__in=plan_pks).values('lesson_plan_id', 'status'):
            lpk = row['lesson_plan_id']
            pulse_counts[lpk] = pulse_counts.get(lpk, 0) + 1
            if row['status'] == 'active':
                active_pk_set.add(lpk)
        for p in plans:
            p.pulse_count      = pulse_counts.get(p.pk, 0)
            p.has_active_pulse = p.pk in active_pk_set
    except Exception:
        pass

    total_plans   = len(plans)
    v3_ready      = sum(1 for p in plans if p.aura_score == 4)
    partial_count = sum(1 for p in plans if 0 < p.aura_score < 4)
    no_aura_count = sum(1 for p in plans if p.aura_score == 0)
    v3_pct        = round(v3_ready / total_plans * 100) if total_plans else 0
    plan_subjects = sorted(set(p.subject.name for p in plans if p.subject))
    plan_classes  = sorted(set(p.school_class.name for p in plans if p.school_class))
    plans.sort(key=lambda p: (-p.aura_score, -p.pk))

    from academics.ai_tutor import get_openai_chat_model
    active_ai_model = get_openai_chat_model()

    context = {
        'sessions': sessions,
        'next_class': next_class_session,
        'active_class': active_class,
        'student_data': student_data,
        'struggling_count': struggling_count,
        'student_count': total_students,
        'current_time': now,
        'active_ai_model': active_ai_model,
        # plan upgrade
        'plans': plans,
        'total_plans': total_plans,
        'v3_ready': v3_ready,
        'v3_pct': v3_pct,
        'partial_count': partial_count,
        'no_aura_count': no_aura_count,
        'plan_subjects': plan_subjects,
        'plan_classes': plan_classes,
    }
    return render(request, 'teachers/ai_sessions_list.html', context)

@login_required
@require_addon('ai-lesson-planner')
def ai_session_new(request):
    if not hasattr(request.user, 'teacher'):
         return redirect('dashboard')
    
    teacher = request.user.teacher
    session = LessonGenerationSession.objects.create(teacher=teacher)
    return redirect('teachers:ai_session_detail', session_id=session.id)

@login_required
@require_addon('ai-lesson-planner')
def ai_session_detail(request, session_id):
    if not hasattr(request.user, 'teacher'):
         return redirect('dashboard')
         
    teacher = request.user.teacher
    session = get_object_or_404(LessonGenerationSession, id=session_id, teacher=teacher)
    
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            user_message = data.get('message')
            
            if user_message:
                messages_list = list(session.messages)
                messages_list.append({"role": "user", "content": user_message})
                session.messages = messages_list
                session.save()
                
                from academics.ai_tutor import _post_chat_completion, get_openai_chat_model
                api_key = settings.OPENAI_API_KEY
                
                system_prompt = """You are an expert lesson planner for teachers. Help them create detailed, engaging lesson plans.
When generating a lesson plan, ALWAYS structure your output to exactly match this template (do not include markdown tables, strictly use this vertical text format with bold headings):

[School Name] [Phone Numbers]
**TERM:** [Term]
**WEEKLY LESSON PLAN – [Class Name]**
**WEEK:** [Week Number]

**Week Ending:** [Date]
**DAY:** [Day of week]
**Subject:** [Subject]
**Duration:** [Duration in minutes]
**Strand:** [Strand description]
**Sub Strand:** [Sub strand description]
**Class:** [Class Name]
**Class Size:** [Number]

**Content Standard:** 
[Details here]

**Indicator:** 
[Details here]

**Lesson:** 
[Lesson Number/Total]

**Performance Indicator:** 
[Details here]

**Core Competencies:** 
[Details here]

**Key words:** 
[Keywords here]

**Reference:** 
[Reference book or curriculum link]

**Phase/Duration** | **Learners Activities** | **Resources**

**PHASE 1: STARTER [Time]**
[Provide starter activity details here...]

**PHASE 2: NEW LEARNING [Time]**
[Provide main teaching activity details here...]

**PHASE 3: REFLECTION [Time]**
[Provide reflection and summary details here...]"""
                api_msgs = [{"role": "system", "content": system_prompt}] + messages_list
                
                payload = {
                    "model": get_openai_chat_model(),
                    "messages": api_msgs,
                    "temperature": 0.7
                }
                
                response_data = _post_chat_completion(payload, api_key)
                
                if "error" in response_data:
                    raise Exception(response_data["error"].get("message", "API Error"))
                    
                ai_content = response_data.get("choices", [{}])[0].get("message", {}).get("content", "")
                
                messages_list.append({"role": "assistant", "content": ai_content})
                session.messages = messages_list
                
                if session.title == "New Session" and len(messages_list) >= 2:
                    try:
                        title_payload = {
                            "model": get_openai_chat_model(),
                            "messages": api_msgs + [{"role": "user", "content": "Suggest a short 3-5 word title for this chat."}],
                            "max_tokens": 15
                        }
                        title_resp = _post_chat_completion(title_payload, api_key)
                        new_title = title_resp.get("choices", [{}])[0].get("message", {}).get("content", "").strip().strip('"')
                        if new_title:
                            session.title = new_title
                    except Exception:
                        pass
                
                session.save()
                return JsonResponse({'status': 'success', 'reply': ai_content, 'title': session.title})
        except Exception as e:
            return _ai_json_error_response(e)
    
    class_subjects = ClassSubject.objects.filter(teacher=teacher).select_related('class_name', 'subject')
    classes = sorted(list(set(cs.class_name for cs in class_subjects)), key=lambda c: c.name)
    subjects = sorted(list(set(cs.subject for cs in class_subjects)), key=lambda s: s.name)

    return render(request, 'teachers/ai_session_detail.html', {
        'session': session,
        'classes': classes,
        'subjects': subjects,
    })

@login_required
@require_addon('ai-lesson-planner')
def ai_session_rename(request, session_id):
    if request.method == 'POST':
        teacher = request.user.teacher
        session = get_object_or_404(LessonGenerationSession, id=session_id, teacher=teacher)
        new_title = request.POST.get('title')
        if new_title:
            session.title = new_title
            session.save()
            messages.success(request, "Session renamed.")
    return redirect('teachers:ai_sessions_list')

@login_required
@require_addon('ai-lesson-planner')
def ai_session_delete(request, session_id):
    teacher = request.user.teacher
    session = get_object_or_404(LessonGenerationSession, id=session_id, teacher=teacher)
    session.delete()
    messages.success(request, "Session deleted.")
    return redirect('teachers:ai_sessions_list')

@login_required
@require_plan('basic', 'pro', 'enterprise')
def assignment_creator(request):
    """
    View for Aura-T Assignment Creator Interface.
    Allows teachers to generate assignments, differentiation, and rubrics.
    """
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    
    teacher = get_object_or_404(Teacher, user=request.user)
    recent_lessons = LessonPlan.objects.filter(teacher=teacher).order_by('-date_added')[:5]
    
    # Get classes and subjects for manual entry
    class_subjects = ClassSubject.objects.filter(teacher=teacher).select_related('class_name', 'subject')
    classes = sorted(list(set(cs.class_name for cs in class_subjects)), key=lambda c: c.name)
    subjects = sorted(list(set(cs.subject for cs in class_subjects)), key=lambda s: s.name)
    
    return render(request, 'teachers/assignment_creator.html', {
        'recent_lessons': recent_lessons,
        'classes': classes,
        'subjects': subjects
    })


# ─────────────────────────────────────────────────────────────
# SCHEME OF WORK — upload, list, delete
# ─────────────────────────────────────────────────────────────

@login_required
def scheme_of_work_list(request):
    """List all schemes of work for the logged-in teacher."""
    if request.user.user_type not in ['admin', 'teacher']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    teacher = get_object_or_404(Teacher, user=request.user)
    current_year = AcademicYear.objects.filter(is_current=True).first()

    class_subjects = ClassSubject.objects.filter(
        teacher=teacher,
        class_name__academic_year=current_year
    ).select_related('class_name', 'subject') if current_year else ClassSubject.objects.filter(teacher=teacher).select_related('class_name', 'subject')

    schemes = SchemeOfWork.objects.filter(
        class_subject__in=class_subjects
    ).select_related('class_subject__class_name', 'class_subject__subject', 'academic_year').order_by('-uploaded_at')

    return render(request, 'teachers/scheme_of_work.html', {
        'schemes': schemes,
        'class_subjects': class_subjects,
        'current_year': current_year,
    })


@login_required
def scheme_of_work_upload(request):
    """Upload or replace a scheme of work for a class/subject/term."""
    if request.user.user_type not in ['admin', 'teacher']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    teacher = get_object_or_404(Teacher, user=request.user)
    current_year = AcademicYear.objects.filter(is_current=True).first()

    if not current_year:
        messages.error(request, 'No active academic year found.')
        return redirect('teachers:scheme_of_work_list')

    class_subjects = ClassSubject.objects.filter(
        teacher=teacher, class_name__academic_year=current_year
    ).select_related('class_name', 'subject')

    if request.method == 'POST':
        cs_id = request.POST.get('class_subject')
        term = request.POST.get('term')
        image = request.FILES.get('image')

        if not cs_id or not term or not image:
            messages.error(request, 'Please fill all fields and upload an image.')
        else:
            cs = get_object_or_404(ClassSubject, id=cs_id, teacher=teacher)

            # Upsert — one scheme per class/subject/term/year
            scheme, created = SchemeOfWork.objects.get_or_create(
                class_subject=cs,
                term=term,
                academic_year=current_year,
                defaults={'uploaded_by': teacher, 'image': image, 'extracted_topics': '[]'},
            )
            if not created:
                # Replace image
                scheme.image = image
                scheme.extracted_topics = '[]'
                scheme.uploaded_by = teacher
                scheme.save()

            # Extract topics and indicators using GPT-4 Vision
            try:
                from academics.ai_tutor import extract_scheme_of_work_data
                # Use absolute path for local storage; fall back to URL for remote
                # backends (Cloudinary etc.) that don't support .path
                try:
                    image_ref = scheme.image.path
                except (NotImplementedError, AttributeError, ValueError):
                    image_ref = scheme.image.url
                data = extract_scheme_of_work_data(image_ref)
                import json
                topics = data.get('topics', [])
                indicators = data.get('indicators', {})
                scheme.extracted_topics = json.dumps(topics)
                scheme.extracted_indicators = json.dumps(indicators)
                scheme.save(update_fields=['extracted_topics', 'extracted_indicators'])
                topic_count = len(topics)
                if topic_count > 0:
                    messages.success(request, f'Scheme of work uploaded! {topic_count} topics extracted for Aura.')
                else:
                    messages.warning(request, 'Scheme uploaded but Aura could not extract topics from the image. You can re-upload a clearer image.')
            except Exception as exc:
                messages.warning(request, f'Scheme saved but topic extraction failed: {exc}')

            return redirect('teachers:scheme_of_work_list')

    return render(request, 'teachers/scheme_of_work_upload.html', {
        'class_subjects': class_subjects,
        'term_choices': SchemeOfWork.TERM_CHOICES,
    })


@login_required
def scheme_of_work_delete(request, pk):
    """Delete a scheme of work."""
    if request.user.user_type not in ['admin', 'teacher']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    teacher = get_object_or_404(Teacher, user=request.user)
    scheme = get_object_or_404(
        SchemeOfWork,
        pk=pk,
        class_subject__teacher=teacher
    )
    scheme.delete()
    messages.success(request, 'Scheme of work deleted.')
    return redirect('teachers:scheme_of_work_list')


@login_required
@require_POST
def scheme_of_work_update_topics(request, pk):
    """AJAX: replace extracted_topics list for a scheme."""
    if request.user.user_type not in ['admin', 'teacher']:
        return JsonResponse({'error': 'Access denied'}, status=403)
    teacher = get_object_or_404(Teacher, user=request.user)
    scheme = get_object_or_404(SchemeOfWork, pk=pk, class_subject__teacher=teacher)
    try:
        payload = json.loads(request.body)
        topics = payload.get('topics', [])
        if not isinstance(topics, list):
            raise ValueError('topics must be a list')
        topics = [str(t).strip() for t in topics if str(t).strip()]
        scheme.extracted_topics = json.dumps(topics)
        indicators = payload.get('indicators', {})
        if isinstance(indicators, dict):
            scheme.extracted_indicators = json.dumps(
                {str(k).strip(): str(v).strip() for k, v in indicators.items()}
            )
        scheme.save(update_fields=['extracted_topics', 'extracted_indicators'])
        return JsonResponse({'ok': True, 'count': len(topics)})
    except Exception as exc:
        return JsonResponse({'error': str(exc)}, status=400)


@login_required
def scheme_of_work_indicators_api(request):
    """Return Scheme of Work indicators for a teacher's selected subject/class, with topic match hints."""
    if request.user.user_type not in ['admin', 'teacher']:
        return JsonResponse({'status': 'error', 'message': 'Access denied'}, status=403)

    teacher = get_object_or_404(Teacher, user=request.user)
    subject_id = (request.GET.get('subject_id') or '').strip()
    class_id = (request.GET.get('class_id') or '').strip()
    topic_query = (request.GET.get('topic') or '').strip()

    if not subject_id or not class_id:
        return JsonResponse({'status': 'success', 'items': [], 'matched_indicator': '', 'matched_topic': ''})

    class_subject = ClassSubject.objects.filter(
        teacher=teacher,
        subject_id=subject_id,
        class_name_id=class_id,
    ).select_related('subject', 'class_name').first()

    if not class_subject:
        return JsonResponse({'status': 'success', 'items': [], 'matched_indicator': '', 'matched_topic': ''})

    current_year = AcademicYear.objects.filter(is_current=True).first()
    schemes = SchemeOfWork.objects.filter(class_subject=class_subject)
    if current_year:
        schemes = schemes.filter(academic_year=current_year)
    schemes = schemes.order_by('-updated_at', '-uploaded_at')

    if not schemes.exists() and current_year:
        schemes = SchemeOfWork.objects.filter(class_subject=class_subject).order_by('-updated_at', '-uploaded_at')

    items = []
    seen = set()

    # Import code-only detection from the GES engine
    try:
        from teachers.services.ges_lesson_engine import GESLessonEngine as _GESLE
        _ges_code_only = _GESLE.is_code_only
    except Exception:
        import re as _re
        _GES_CODE_ONLY_RE = _re.compile(r'^[A-Z]{1,3}\d{1,2}(?:\.\d+){2,5}\s*$')
        _ges_code_only = lambda ind: bool(_GES_CODE_ONLY_RE.match((ind or '').strip()))

    for scheme in schemes:
        indicators = scheme.get_indicators()
        topics = scheme.get_topics()

        for topic in topics:
            t = str(topic).strip()
            if not t:
                continue
            indicator = str(indicators.get(t, '')).strip()
            key = (t.lower(), indicator.lower())
            if key in seen:
                continue
            seen.add(key)
            items.append({'topic': t, 'indicator': indicator, 'code_only': _ges_code_only(indicator)})

        # Include indicators whose topics are in the dict but absent in extracted_topics list
        for t_raw, ind_raw in indicators.items():
            t = str(t_raw).strip()
            if not t:
                continue
            indicator = str(ind_raw).strip()
            key = (t.lower(), indicator.lower())
            if key in seen:
                continue
            seen.add(key)
            items.append({'topic': t, 'indicator': indicator, 'code_only': _ges_code_only(indicator)})

    matched_indicator = ''
    matched_topic = ''
    q = topic_query.lower()
    if q:
        # Exact topic first, then contains.
        exact = next((it for it in items if it['topic'].strip().lower() == q), None)
        contains = next((it for it in items if q in it['topic'].strip().lower()), None)
        winner = exact or contains
        if winner:
            matched_topic = winner['topic']
            matched_indicator = winner['indicator']

    return JsonResponse({
        'status': 'success',
        'items': items,
        'matched_indicator': matched_indicator,
        'matched_topic': matched_topic,
    })


@login_required
@require_POST
def scheme_of_work_reextract(request, pk):
    """Re-run GPT-4o Vision extraction on the stored image."""
    if request.user.user_type not in ['admin', 'teacher']:
        return JsonResponse({'error': 'Access denied'}, status=403)
    teacher = get_object_or_404(Teacher, user=request.user)
    scheme = get_object_or_404(SchemeOfWork, pk=pk, class_subject__teacher=teacher)
    try:
        from academics.ai_tutor import extract_scheme_of_work_data
        try:
            image_ref = scheme.image.path
        except (NotImplementedError, AttributeError, ValueError):
            image_ref = scheme.image.url
        data = extract_scheme_of_work_data(image_ref)
        topics = data.get('topics', [])
        indicators = data.get('indicators', {})
        scheme.extracted_topics = json.dumps(topics)
        scheme.extracted_indicators = json.dumps(indicators)
        scheme.save(update_fields=['extracted_topics', 'extracted_indicators'])
        return JsonResponse({'ok': True, 'count': len(topics), 'topics': topics, 'indicators': indicators})
    except Exception as exc:
        return JsonResponse({'error': str(exc)}, status=500)


@login_required
@require_POST
def scheme_of_work_bulk_generate(request, pk):
    """
    Bulk-generate Aura-T lesson plans for all topics in a scheme of work.
    Skips topics that already have a matching lesson plan.
    Returns JSON: {ok, created, skipped, errors: [...], total}
    """
    if request.user.user_type not in ['admin', 'teacher']:
        return JsonResponse({'error': 'Access denied'}, status=403)

    teacher = get_object_or_404(Teacher, user=request.user)
    scheme = get_object_or_404(SchemeOfWork, pk=pk, class_subject__teacher=teacher)

    topics = scheme.get_topics()
    if not topics:
        return JsonResponse({'error': 'No topics found. Extract topics from the scheme first.'}, status=400)

    # Quota check upfront: count only the topics that will actually be generated
    from tenants.ai_quota import check_and_consume, QuotaExceeded
    already_done = sum(
        1 for t in topics
        if LessonPlan.objects.filter(
            teacher=teacher, subject=scheme.class_subject.subject,
            school_class=scheme.class_subject.class_name, topic__iexact=t
        ).exists()
    )
    to_generate = len(topics) - already_done
    if to_generate > 0:
        try:
            check_and_consume(request.tenant, request.user.id, 'bulk_gen', call_count=to_generate)
        except QuotaExceeded as e:
            return JsonResponse({
                'error': e.user_message,
                'error_code': 'quota_exceeded',
                'used': e.used,
                'limit': e.limit,
                'addon_boost': e.addon_boost,
            }, status=429)

    from teachers.services.aura_gen_engine import AuraGenEngine
    import re as _re

    subject     = scheme.class_subject.subject
    school_class = scheme.class_subject.class_name
    subject_name = subject.name
    class_name   = school_class.name

    # Start week numbering after any existing plans for this teacher / subject / class
    existing_max = LessonPlan.objects.filter(
        teacher=teacher, subject=subject, school_class=school_class
    ).aggregate(max_week=models.Max('week_number'))['max_week'] or 0

    def _extract(header, text):
        pat = _re.compile(
            rf"\*\*{_re.escape(header)}.*?\*\*\s*(.*?)(?=\n\*\*|\Z)",
            _re.DOTALL | _re.IGNORECASE
        )
        m = pat.search(text)
        return m.group(1).strip() if m else ''

    created = 0
    skipped = 0
    errors  = []

    for idx, topic in enumerate(topics):
        # Skip if a plan for this topic already exists for this teacher/subject/class
        if LessonPlan.objects.filter(
            teacher=teacher, subject=subject,
            school_class=school_class, topic__iexact=topic
        ).exists():
            skipped += 1
            continue

        try:
            result      = AuraGenEngine.generate_lesson_plan(topic, subject_name, class_name)
            lesson_body = (result.get('lesson_plan') or '').strip()
            if not lesson_body:
                errors.append({'topic': topic, 'error': result.get('error', 'Empty response')})
                continue

            LessonPlan.objects.create(
                teacher=teacher,
                subject=subject,
                school_class=school_class,
                week_number=existing_max + idx + 1,
                topic=topic,
                objectives=(
                    _extract('Content Standard', lesson_body) + '\n' +
                    _extract('Indicator', lesson_body)
                ).strip() or topic,
                teaching_materials=_extract('Teaching/Learning Materials', lesson_body),
                introduction=_extract('PHASE 1: STARTER', lesson_body),
                presentation=(
                    _extract('PHASE 2: NEW LEARNING', lesson_body) or
                    _extract('PHASE 2', lesson_body)
                ),
                evaluation=(
                    _extract('PHASE 3: REFLECTION', lesson_body) or
                    _extract('Assessment', lesson_body)
                ),
                homework=_extract('Homework', lesson_body),
            )
            created += 1
        except Exception as exc:
            errors.append({'topic': topic, 'error': str(exc)})

    return JsonResponse({
        'ok': True,
        'created': created,
        'skipped': skipped,
        'errors': errors,
        'total': len(topics),
    })


@login_required
def scheme_of_work_edit(request, pk):
    """Edit scheme metadata (term / class-subject) and optionally replace the image."""
    if request.user.user_type not in ['admin', 'teacher']:
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    teacher = get_object_or_404(Teacher, user=request.user)
    scheme  = get_object_or_404(SchemeOfWork, pk=pk, class_subject__teacher=teacher)
    current_year = AcademicYear.objects.filter(is_current=True).first()

    class_subjects = ClassSubject.objects.filter(
        teacher=teacher, class_name__academic_year=current_year
    ).select_related('class_name', 'subject') if current_year else \
        ClassSubject.objects.filter(teacher=teacher).select_related('class_name', 'subject')

    if request.method == 'POST':
        cs_id = request.POST.get('class_subject')
        term  = request.POST.get('term')
        new_image = request.FILES.get('image')

        if not cs_id or not term:
            messages.error(request, 'Class/subject and term are required.')
        else:
            cs = get_object_or_404(ClassSubject, id=cs_id, teacher=teacher)
            # Guard uniqueness (only if changing the key fields)
            conflict = SchemeOfWork.objects.filter(
                class_subject=cs, term=term,
                academic_year=scheme.academic_year
            ).exclude(pk=scheme.pk).first()
            if conflict:
                messages.error(request, 'A scheme of work for that class/subject/term already exists.')
            else:
                scheme.class_subject = cs
                scheme.term = term
                update_fields = ['class_subject', 'term']

                if new_image:
                    scheme.image = new_image
                    # Clear old topics so user knows extraction is pending
                    scheme.extracted_topics = '[]'
                    scheme.extracted_indicators = '{}'
                    update_fields += ['image', 'extracted_topics', 'extracted_indicators']
                    scheme.save(update_fields=update_fields)

                    # Re-extract with new image
                    try:
                        from academics.ai_tutor import extract_scheme_of_work_data
                        try:
                            image_ref = scheme.image.path
                        except (NotImplementedError, AttributeError, ValueError):
                            image_ref = scheme.image.url
                        data = extract_scheme_of_work_data(image_ref)
                        topics = data.get('topics', [])
                        indicators = data.get('indicators', {})
                        scheme.extracted_topics = json.dumps(topics)
                        scheme.extracted_indicators = json.dumps(indicators)
                        scheme.save(update_fields=['extracted_topics', 'extracted_indicators'])
                        messages.success(request, f'Scheme updated. {len(topics)} topic(s) re-extracted.')
                    except Exception as exc:
                        messages.warning(request, f'Scheme saved but re-extraction failed: {exc}')
                else:
                    scheme.save(update_fields=update_fields)
                    messages.success(request, 'Scheme updated.')

                return redirect('teachers:scheme_of_work_list')

    return render(request, 'teachers/scheme_of_work_edit.html', {
        'scheme': scheme,
        'class_subjects': class_subjects,
        'term_choices': SchemeOfWork.TERM_CHOICES,
    })


@login_required
@require_POST
def scheme_of_work_dedup_indicators(request, pk):
    """
    AJAX: detect or fix duplicate indicators within a scheme.
    POST body:
      { "action": "scan" }            → returns duplicates list
      { "action": "apply", "indicators": {...} }  → saves new indicators
    """
    if request.user.user_type not in ['admin', 'teacher']:
        return JsonResponse({'error': 'Access denied'}, status=403)

    teacher = get_object_or_404(Teacher, user=request.user)
    scheme  = get_object_or_404(SchemeOfWork, pk=pk, class_subject__teacher=teacher)

    try:
        payload = json.loads(request.body)
    except Exception:
        return JsonResponse({'error': 'Invalid JSON'}, status=400)

    action = payload.get('action', 'scan')

    if action == 'scan':
        import re as _re
        indicators = scheme.get_indicators()  # {topic: indicator_text}

        def _ind_code(s):
            """Extract the leading code segment, e.g. 'B8.2.1.1.1' from a full indicator string."""
            s = s.strip()
            m = _re.match(r'^([A-Za-z][0-9A-Za-z.]*)', s)
            return m.group(1).rstrip('.').upper() if m else s.upper()

        # Build reverse map: normalised_code → [(topic, full_indicator_text), ...]
        rev = {}
        for topic, ind in indicators.items():
            if not ind:
                continue
            code = _ind_code(ind)
            rev.setdefault(code, []).append({'topic': topic, 'indicator': ind})

        duplicates = []
        for code, entries in rev.items():
            if len(entries) > 1:
                duplicates.append({
                    'code': code,
                    # Use the most common full text as the display indicator
                    'indicator': entries[0]['indicator'],
                    'topics': [e['topic'] for e in entries],
                    'entries': entries,   # full detail for the UI
                })

        return JsonResponse({'ok': True, 'duplicates': duplicates, 'count': len(duplicates)})

    elif action == 'apply':
        new_indicators = payload.get('indicators', {})
        if not isinstance(new_indicators, dict):
            return JsonResponse({'error': 'indicators must be a dict'}, status=400)
        scheme.extracted_indicators = json.dumps(
            {str(k).strip(): str(v).strip() for k, v in new_indicators.items()}
        )
        scheme.save(update_fields=['extracted_indicators'])
        return JsonResponse({'ok': True})

    return JsonResponse({'error': 'Unknown action'}, status=400)


@login_required
def student_session_peek(request, session_id):
    """
    Teacher-only read-only view of a student's AI Tutor session transcript.
    Access is permitted only if the teacher teaches or manages the student's class.
    """
    if request.user.user_type not in ['teacher', 'admin']:
        messages.error(request, "Access restricted to teachers.")
        return redirect('dashboard')

    session = get_object_or_404(TutorSession, id=session_id)
    student = session.student

    if request.user.user_type == 'teacher':
        teacher = get_object_or_404(Teacher, user=request.user)
        # Check teacher teaches or manages this student's class
        teaches_class = (
            student.current_class and (
                teacher.managed_classes.filter(id=student.current_class_id).exists() or
                ClassSubject.objects.filter(
                    teacher=teacher,
                    class_name=student.current_class
                ).exists()
            )
        )
        if not teaches_class:
            messages.error(request, "You do not have access to this student's session.")
            return redirect('teachers:teacher_ai_insights')

    messages_qs = TutorMessage.objects.filter(session=session).order_by('created_at')

    return render(request, 'teachers/student_session_peek.html', {
        'session': session,
        'student': student,
        'chat_messages': messages_qs,
    })


@login_required
@require_POST
def submit_to_hod(request):
    """
    Submit a lesson plan / summary to the Head of Department (school admin).
    Creates a Notification for all admin users.
    """
    if request.user.user_type != 'teacher':
        return JsonResponse({'error': 'Access denied'}, status=403)

    teacher = get_object_or_404(Teacher, user=request.user)
    lesson_title = request.POST.get('lesson_title', '').strip() or 'Untitled Lesson Plan'
    lesson_body = request.POST.get('lesson_body', '').strip()

    from announcements.models import Notification
    admin_users = User.objects.filter(user_type='admin')
    created = 0
    for admin_user in admin_users:
        Notification.objects.create(
            user=admin_user,
            title=f"Lesson Plan Submitted: {lesson_title}",
            message=(
                f"Teacher {teacher.user.get_full_name()} submitted a lesson plan for your review.\n\n"
                f"{lesson_body[:500]}" if lesson_body else
                f"Teacher {teacher.user.get_full_name()} submitted a lesson plan for your review."
            ),
            notification_type='announcement',
        )
        created += 1

    return JsonResponse({'ok': True, 'notified': created, 'message': f'Submitted to {created} admin(s).'})


@login_required
def save_aura_t_plan(request):
    """Save the Aura-T command centre lesson plan preview as a LessonPlan record."""
    if request.method != 'POST':
        return JsonResponse({'ok': False, 'error': 'POST required'}, status=405)
    if request.user.user_type != 'teacher':
        return JsonResponse({'ok': False, 'error': 'Teachers only'}, status=403)

    try:
        teacher = Teacher.objects.get(user=request.user)
    except Teacher.DoesNotExist:
        return JsonResponse({'ok': False, 'error': 'Teacher profile not found'}, status=404)

    lesson_title = request.POST.get('lesson_title', '').strip() or 'Aura-T Lesson Plan'
    lesson_body = request.POST.get('lesson_body', '').strip()

    # Parse from body using the same section extractor as lesson_plan_create
    import re as _re

    def _extract(header, text):
        pat = _re.compile(rf"\*\*{_re.escape(header)}.*?\*\*\s*(.*?)(?=\n\*\*|\Z)", _re.DOTALL | _re.IGNORECASE)
        m = pat.search(text)
        return m.group(1).strip() if m else ''

    def _extract_phase(phase_num, text):
        """Extract full phase content up to the next phase-level boundary."""
        boundaries = {1: r'\*\*PHASE\s*2', 2: r'\*\*PHASE\s*3', 3: r'\*\*(?:Resources|Homework)'}
        end = boundaries.get(phase_num, r'\Z')
        pat = _re.compile(rf"\*\*PHASE\s*{phase_num}[^*]*\*\*\s*(.*?)(?={end}|\Z)", _re.DOTALL | _re.IGNORECASE)
        m = pat.search(text)
        return m.group(1).strip() if m else ''

    topic = _extract('Strand', lesson_body) or lesson_title
    sub_strand = _extract('Sub Strand', lesson_body)

    # Get first class the teacher is associated with
    from academics.models import ClassSubject
    class_subject = ClassSubject.objects.filter(teacher=teacher).select_related('school_class', 'subject').first()

    plan = LessonPlan.objects.create(
        teacher=teacher,
        subject=class_subject.subject if class_subject else None,
        school_class=class_subject.school_class if class_subject else None,
        week_number=1,
        topic=topic,
        sub_strand=sub_strand,
        objectives=_extract('Content Standard', lesson_body) + '\n' + _extract('Indicator', lesson_body),
        introduction=_extract_phase(1, lesson_body),
        presentation=_extract_phase(2, lesson_body),
        evaluation=_extract_phase(3, lesson_body) or _extract('Assessment', lesson_body),
        homework=_extract('Homework', lesson_body) or 'Refer to textbook exercises.',
    )

    return JsonResponse({'ok': True, 'plan_id': plan.id, 'detail_url': f'/lesson-plans/{plan.id}/'})


# ─────────────────────────────────────────────────────────────────────────────
# DIGITAL PULSE — real-time class engagement
# ─────────────────────────────────────────────────────────────────────────────

@login_required
@require_POST
def pulse_launch(request, plan_pk):
    """Teacher launches a pulse for a lesson plan.  Returns JSON {session_id}."""
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)

    try:
        from academics.pulse_models import PulseSession, parse_pulse_questions, parse_q3_chips
    except ImportError:
        return JsonResponse({'error': 'Pulse feature not available'}, status=500)

    teacher = None
    if request.user.user_type == 'teacher':
        teacher = get_object_or_404(Teacher, user=request.user)
        plan = get_object_or_404(LessonPlan, pk=plan_pk, teacher=teacher)
    else:
        plan = get_object_or_404(LessonPlan, pk=plan_pk)
        teacher = plan.teacher

    # Close any previous active sessions for this plan
    PulseSession.objects.filter(lesson_plan=plan, status='active').update(status='closed')

    q1, q2, q3 = parse_pulse_questions(plan.introduction or '')
    chips       = parse_q3_chips(plan)

    session = PulseSession.objects.create(
        lesson_plan=plan,
        teacher=teacher,
        q1_text=q1,
        q2_text=q2,
        q3_text=q3,
        q3_chips=chips,
    )
    return JsonResponse({
        'session_id': session.pk,
        'q1': q1, 'q2': q2, 'q3': q3,
        'chips': chips,
        'total_students': session.total_students,
    })


@login_required
def pulse_live(request, session_id):
    """Polling endpoint — teacher gets live response count + who's typing."""
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)

    try:
        from academics.pulse_models import PulseSession
    except ImportError:
        return JsonResponse({'error': 'unavailable'}, status=500)

    try:
        session = get_object_or_404(PulseSession, pk=session_id)
    except (ProgrammingError, OperationalError):
        return JsonResponse({'error': 'Pulse data unavailable. Run migrations.'}, status=503)

    try:
        responded = session.responded_count
        total     = session.total_students
        typing    = [f'{fn} {ln}'.strip() for fn, ln in session.typing_students]
    except (ProgrammingError, OperationalError):
        responded = 0
        total = 0
        typing = []

    # Build aggregate for Q1 / Q2
    q1_true = q1_false = q2_true = q2_false = 0
    q3_counts = {}
    try:
        for r in session.responses.filter(submitted_at__isnull=False):
            if r.q1_answer is True:  q1_true  += 1
            if r.q1_answer is False: q1_false += 1
            if r.q2_answer is True:  q2_true  += 1
            if r.q2_answer is False: q2_false += 1
            if r.q3_answer:
                q3_counts[r.q3_answer] = q3_counts.get(r.q3_answer, 0) + 1
    except (ProgrammingError, OperationalError):
        q1_true = q1_false = q2_true = q2_false = 0
        q3_counts = {}

    return JsonResponse({
        'status': getattr(session, 'status', 'active'),
        'responded': responded,
        'total': total,
        'typing': typing,
        'q1': {'true': q1_true, 'false': q1_false},
        'q2': {'true': q2_true, 'false': q2_false},
        'q3_counts': q3_counts,
    })


@login_required
@require_POST
def pulse_close(request, session_id):
    """Teacher closes an active pulse session."""
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)

    try:
        from academics.pulse_models import PulseSession
    except ImportError:
        return JsonResponse({'error': 'unavailable'}, status=500)

    session = get_object_or_404(PulseSession, pk=session_id)
    if session.status == 'active':
        session.status    = 'closed'
        session.closed_at = timezone.now()
        session.save(update_fields=['status', 'closed_at'])

    from django.urls import reverse
    # Prepend SCRIPT_NAME so path-based tenant prefix (/school1/) is included
    script_name  = request.META.get('SCRIPT_NAME', '').rstrip('/')
    results_path = reverse('teachers:pulse_results', args=[session_id])
    results_url  = script_name + results_path
    return JsonResponse({'ok': True, 'redirect': results_url})


@login_required
def pulse_results(request, session_id):
    """Post-session Pulse results summary for the teacher."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')

    try:
        from academics.pulse_models import PulseSession, PulseResponse
    except ImportError:
        messages.error(request, 'Pulse feature not available.')
        return redirect('teachers:lesson_plan_list')

    session = get_object_or_404(PulseSession, pk=session_id)

    # Ensure the requesting teacher owns the session (admins bypass)
    if request.user.user_type == 'teacher':
        teacher = get_object_or_404(Teacher, user=request.user)
        if session.teacher != teacher:
            messages.error(request, 'Access denied.')
            return redirect('teachers:lesson_plan_list')

    # ── Resolve class + topic from either lesson_plan or presentation ───────
    if session.lesson_plan_id:
        school_class = session.lesson_plan.school_class
        topic = session.lesson_plan.topic
    elif session.target_class_id:
        school_class = session.target_class
        topic = session.presentation.title if session.presentation_id else 'Pulse Check'
    else:
        school_class = None
        topic = session.presentation.title if session.presentation_id else 'Pulse Check'

    # ── All students in the target class ───────────────────────────────────
    all_students = (
        Student.objects.filter(current_class=school_class, user__is_active=True)
        .select_related('user')
        .order_by('user__last_name', 'user__first_name')
    ) if school_class else Student.objects.none()

    # ── Index responses by student_id ──────────────────────────────────────
    response_map = {
        r.student_id: r
        for r in session.responses.select_related('student__user').all()
    }

    # ── Build per-student rows ─────────────────────────────────────────────
    rows = []
    q1_true = q1_false = q1_skip = 0
    q2_true = q2_false = q2_skip = 0
    q3_counts = {}
    no_show_count = 0

    for stu in all_students:
        resp = response_map.get(stu.pk)
        submitted = resp is not None and resp.submitted_at is not None

        if resp and resp.q1_answer is True:   q1_true  += 1
        elif resp and resp.q1_answer is False: q1_false += 1
        else:                                  q1_skip  += 1

        if resp and resp.q2_answer is True:   q2_true  += 1
        elif resp and resp.q2_answer is False: q2_false += 1
        else:                                  q2_skip  += 1

        q3_val = (resp.q3_answer or '').strip() if resp else ''
        if q3_val:
            q3_counts[q3_val] = q3_counts.get(q3_val, 0) + 1

        if not submitted:
            no_show_count += 1

        # ── At-risk logic ──────────────────────────────────────────────
        # Q1 = prerequisite (correct = True), Q2 = misconception (correct = False)
        # Flag if: didn't respond, OR answered Q1 False, OR answered Q2 True
        at_risk = (
            not submitted
            or (resp and resp.q1_answer is False)
            or (resp and resp.q2_answer is True)
        )

        rows.append({
            'student': stu,
            'submitted': submitted,
            'q1': resp.q1_answer if resp else None,
            'q2': resp.q2_answer if resp else None,
            'q3': q3_val,
            'at_risk': at_risk,
        })

    total = len(rows)
    responded = total - no_show_count
    response_pct = round(responded / total * 100) if total else 0

    # Q3 sorted by frequency
    q3_sorted = sorted(q3_counts.items(), key=lambda x: -x[1])

    # At-risk count
    at_risk_count = sum(1 for r in rows if r['at_risk'])

    context = {
        'session': session,
        'plan': session.lesson_plan,
        'topic': topic,
        'rows': rows,
        'total': total,
        'responded': responded,
        'response_pct': response_pct,
        'no_show_count': no_show_count,
        'at_risk_count': at_risk_count,
        'q1_true': q1_true, 'q1_false': q1_false, 'q1_skip': q1_skip,
        'q2_true': q2_true, 'q2_false': q2_false, 'q2_skip': q2_skip,
        'q3_counts': q3_sorted,
        'q3_chips': session.q3_chips,
        'school_class': school_class,
    }
    return render(request, 'teachers/pulse_results.html', context)


@login_required
def pulse_history(request):
    """All past Pulse sessions for the teacher (or all sessions for admins)."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')

    try:
        from academics.pulse_models import PulseSession
    except ImportError:
        messages.error(request, 'Pulse feature not available.')
        return redirect('teachers:aura_command_center')

    is_admin = request.user.user_type == 'admin'

    _related = [
        'lesson_plan__school_class', 'lesson_plan__subject',
        'presentation', 'target_class', 'teacher__user',
    ]
    if is_admin:
        qs = PulseSession.objects.select_related(*_related).order_by('-created_at')
    else:
        teacher = get_object_or_404(Teacher, user=request.user)
        qs = PulseSession.objects.filter(teacher=teacher).select_related(
            *_related
        ).order_by('-created_at')

    sessions = []
    for s in qs:
        total     = s.total_students
        responded = s.responded_count
        pct       = round(responded / total * 100) if total else 0
        # Resolve topic + class regardless of whether it was lesson_plan or presentation based
        if s.lesson_plan_id:
            s_topic = s.lesson_plan.topic
            s_class = s.lesson_plan.school_class
            s_subject = s.lesson_plan.subject
        else:
            s_topic = s.presentation.title if s.presentation_id else 'Pulse Check'
            s_class = s.target_class
            s_subject = s.presentation.subject if s.presentation_id else None
        sessions.append({
            'session':    s,
            'plan':       s.lesson_plan,
            'topic':      s_topic,
            'school_class': s_class,
            'subject':    s_subject,
            'total':      total,
            'responded':  responded,
            'pct':        pct,
        })

    return render(request, 'teachers/pulse_history.html', {
        'sessions':       sessions,
        'is_admin':       is_admin,
        'total_sessions': len(sessions),
        'active_count':   sum(1 for s in sessions if s['session'].status == 'active'),
        'closed_count':   sum(1 for s in sessions if s['session'].status == 'closed'),
    })


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE DECK CREATOR  (Gamma-style)
# ═══════════════════════════════════════════════════════════════════════════════

@login_required
@require_addon('presentations')
def presentation_list(request):
    """Dashboard listing all of the teacher's presentations."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation, Slide
    from django.db.models import OuterRef, Subquery, Count as _Count
    subject_id = (request.GET.get('subject_id') or '').strip()
    class_id = (request.GET.get('class_id') or '').strip()

    deck_qs = Presentation.objects.filter(teacher=teacher)
    if subject_id:
        deck_qs = deck_qs.filter(subject_id=subject_id)
    if class_id:
        deck_qs = deck_qs.filter(school_class_id=class_id)

    first_emoji = Slide.objects.filter(
        presentation=OuterRef('pk')).order_by('order').values('emoji')[:1]
    decks = list(
        deck_qs
        .select_related('subject', 'school_class')
        .annotate(
            cover_emoji=Subquery(first_emoji),
            annotated_slide_count=_Count('slides'),
        )
    )
    total_slides = sum(d.annotated_slide_count for d in decks)
    return render(request, 'teachers/presentations/list.html', {
        'decks': decks,
        'total_decks': len(decks),
        'total_slides': total_slides,
        'selected_subject_id': subject_id,
        'selected_class_id': class_id,
    })


@login_required
@require_addon('presentations')
def presentation_create(request):
    """Create a new blank (or AI-seeded) presentation, then redirect to editor."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation, Slide
    from academics.models import Subject, Class

    prefill_subject_id = (request.GET.get('subject_id') or '').strip()
    prefill_class_id = (request.GET.get('class_id') or '').strip()
    prefill_title = (request.GET.get('title') or '').strip()
    prefill_lesson_topic = (request.GET.get('lesson_topic') or '').strip()
    prefill_smart_seed = (request.GET.get('smart_seed') or '').strip().lower() in ('1', 'true', 'yes', 'on')
    prefill_theme = (request.GET.get('theme') or '').strip().lower()
    valid_themes = {value for value, _label in Presentation.THEME_CHOICES}
    if prefill_theme not in valid_themes:
        prefill_theme = ''

    if not prefill_theme and prefill_subject_id:
        subject_name = (
            Subject.objects.filter(pk=prefill_subject_id)
            .values_list('name', flat=True)
            .first()
            or ''
        ).lower()

        # Choose a visual mood that matches broad subject families.
        theme_rules = [
            (('science', 'biology', 'chemistry', 'physics', 'ict', 'computer', 'technology'), 'ocean'),
            (('math', 'mathematics', 'algebra', 'geometry', 'statistics', 'arithmetic'), 'slate'),
            (('social', 'history', 'geography', 'civic', 'government', 'economics'), 'amber'),
            (('english', 'language', 'literature', 'creative writing', 'reading'), 'aurora'),
            (('french', 'twi', 'akan', 'ewe', 'ga', 'hausa', 'dagbani'), 'rose'),
            (('rme', 'religion', 'ethics', 'guidance'), 'midnight'),
            (('creative', 'art', 'music', 'design', 'performing'), 'coral'),
            (('physical', 'pe', 'sports', 'health', 'career technology'), 'forest'),
        ]

        for keywords, suggested_theme in theme_rules:
            if any(keyword in subject_name for keyword in keywords):
                prefill_theme = suggested_theme
                break

    if not prefill_theme:
        prefill_theme = 'aurora'

    if request.method == 'POST':
        title        = request.POST.get('title', 'Untitled Deck').strip() or 'Untitled Deck'
        theme        = request.POST.get('theme', 'aurora')
        subject_id   = request.POST.get('subject_id')
        class_id     = request.POST.get('class_id')
        lesson_topic = request.POST.get('lesson_topic', '').strip()
        smart_seed   = (request.POST.get('smart_seed') or '').strip().lower() in ('1', 'true', 'yes', 'on')
        subject      = Subject.objects.filter(pk=subject_id).first() if subject_id else None
        school_class = Class.objects.filter(pk=class_id).first()      if class_id   else None

        deck = Presentation.objects.create(
            teacher=teacher, title=title, theme=theme,
            subject=subject, school_class=school_class,
        )
        if smart_seed:
            latest_plan = None
            if subject and school_class:
                latest_plan = (
                    LessonPlan.objects
                    .filter(teacher=teacher, subject=subject, school_class=school_class)
                    .order_by('-updated_at', '-id')
                    .first()
                )

            topic = (latest_plan.topic if latest_plan and latest_plan.topic else lesson_topic) or title
            subtitle_bits = []
            if subject:
                subtitle_bits.append(subject.name)
            if school_class:
                subtitle_bits.append(school_class.name)

            Slide.objects.create(
                presentation=deck,
                order=0,
                layout='title',
                title=topic,
                content=' - '.join(subtitle_bits) if subtitle_bits else 'Starter lesson deck',
                emoji='🚀',
            )

            objectives_text = (latest_plan.objectives if latest_plan else '').strip()
            if objectives_text:
                objective_lines = [line.strip('-• ').strip() for line in objectives_text.splitlines() if line.strip()]
                if not objective_lines:
                    objective_lines = [s.strip() for s in objectives_text.split('.') if s.strip()]
            else:
                objective_lines = []

            if not objective_lines:
                objective_lines = [
                    f"Understand the key idea behind {topic}",
                    "Apply the concept through guided examples",
                    "Check understanding with quick formative questions",
                ]

            Slide.objects.create(
                presentation=deck,
                order=1,
                layout='bullets',
                title='Learning Objectives',
                content='\n'.join(objective_lines[:4]),
                emoji='🎯',
            )

            procedure_bits = []
            if latest_plan:
                for section in [latest_plan.introduction, latest_plan.presentation, latest_plan.evaluation, latest_plan.homework]:
                    clean = (section or '').strip()
                    if clean:
                        procedure_bits.append(clean)

            if procedure_bits:
                procedure_text = procedure_bits[0]
                steps = [line.strip('-• ').strip() for line in procedure_text.splitlines() if line.strip()]
                if not steps:
                    steps = [s.strip() for s in procedure_text.split('.') if s.strip()]
                if not steps:
                    steps = ['Starter activity', 'Main explanation', 'Quick assessment', 'Exit ticket']
            else:
                steps = ['Starter activity', 'Main explanation', 'Guided practice', 'Exit ticket']

            Slide.objects.create(
                presentation=deck,
                order=2,
                layout='bullets',
                title='Lesson Flow',
                content='\n'.join(steps[:4]),
                emoji='🧭',
            )
        else:
            # Seed one title slide
            Slide.objects.create(
                presentation=deck, order=0, layout='title',
                title=title, content='Add your subtitle here', emoji='🚀',
            )
        return redirect('teachers:presentation_editor', pk=deck.pk)

    subjects = Subject.objects.all().order_by('name')
    classes  = Class.objects.all().order_by('name')
    return render(request, 'teachers/presentations/create.html', {
        'subjects': subjects, 'classes': classes,
        'THEME_CHOICES': Presentation.THEME_CHOICES,
        'prefill_subject_id': prefill_subject_id,
        'prefill_class_id': prefill_class_id,
        'prefill_title': prefill_title,
        'prefill_theme': prefill_theme,
        'prefill_lesson_topic': prefill_lesson_topic,
        'prefill_smart_seed': prefill_smart_seed,
    })


@login_required
@require_addon('presentations')
def presentation_editor(request, pk):
    """Main Gamma-style slide editor."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation, Slide
    from academics.models import Subject, Class

    deck   = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    slides = list(deck.slides.all())

    subjects = Subject.objects.all().order_by('name')
    classes  = Class.objects.all().order_by('name')

    import json as _json
    from django.urls import reverse as _reverse
    slides_json = _json.dumps([{
        'id':            s.pk,
        'order':         s.order,
        'layout':        s.layout,
        'title':         s.title,
        'content':       s.content,
        'emoji':         s.emoji,
        'speaker_notes': s.speaker_notes,
        'image_url':     s.image_url,
    } for s in slides])
    _share_path = _reverse('teachers:presentation_share', kwargs={'token': deck.share_token})
    _share_url  = request.build_absolute_uri(_share_path)

    slide_gen_used = get_free_generation_count(request.user, 'slide_gen')
    has_slide_addon = has_addon(request.user, 'aura-slide-generator')

    return render(request, 'teachers/presentations/editor.html', {
        'deck':     deck,
        'slides':   slides,
        'subjects': subjects,
        'classes':  classes,
        'slides_json':    slides_json,
        'share_url':      _share_url,
        'slide_gen_used':  slide_gen_used,
        'slide_gen_limit': FREE_GENERATION_LIMIT,
        'slide_gen_remaining': max(0, FREE_GENERATION_LIMIT - slide_gen_used),
        'has_slide_addon': has_slide_addon,
        'LAYOUT_CHOICES':     Slide.LAYOUT_CHOICES,
        'THEME_CHOICES':      Presentation.THEME_CHOICES,
        'TRANSITION_CHOICES': Presentation.TRANSITION_CHOICES,
        'EMOJI_LIST': ['🚀','🌱','📚','🧠','💡','🔬','🌍','🎯','🏆','✏️','🧪','🌊','💻','🎨','📊'],
    })


@login_required
@require_addon('presentations')
def presentation_delete(request, pk):
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    if request.method == 'POST':
        deck.delete()
        messages.success(request, 'Presentation deleted.')
        return redirect('teachers:presentation_list')
    return render(request, 'teachers/presentations/confirm_delete.html', {'deck': deck})


@login_required
@require_addon('presentations')
def presentation_present(request, pk):
    """Fullscreen slideshow view — read-only."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation
    from django.urls import reverse
    from django.utils import timezone
    deck   = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    slides = list(deck.slides.all())
    share_path = reverse('teachers:presentation_share', kwargs={'token': deck.share_token})
    share_url  = request.build_absolute_uri(share_path)
    # Track presentation usage
    Presentation.objects.filter(pk=pk).update(
        times_presented=models.F('times_presented') + 1,
        last_presented_at=timezone.now(),
    )
    return render(request, 'teachers/presentations/present.html', {
        'deck': deck, 'slides': slides, 'share_url': share_url,
    })


@login_required
@require_addon('presentations')
def presentation_duplicate(request, pk):
    """Duplicate a deck and all its slides, redirect to the new editor."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation, Slide
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    if request.method == 'POST':
        new_deck = Presentation.objects.create(
            teacher=teacher,
            title=deck.title + ' (copy)',
            theme=deck.theme,
            transition=deck.transition,
            subject=deck.subject,
            school_class=deck.school_class,
        )
        for slide in deck.slides.all():
            Slide.objects.create(
                presentation=new_deck,
                order=slide.order,
                layout=slide.layout,
                title=slide.title,
                content=slide.content,
                speaker_notes=slide.speaker_notes,
                emoji=slide.emoji,
                image_url=slide.image_url,
            )
        messages.success(request, f'"{new_deck.title}" created.')
        return redirect('teachers:presentation_editor', pk=new_deck.pk)
    return redirect('teachers:presentation_list')


@login_required
def presentation_print(request, pk):
    """Print-friendly handout view — all slides rendered on one page."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    slides = list(deck.slides.all())
    return render(request, 'teachers/presentations/print.html', {
        'deck': deck, 'slides': slides,
    })


# ── Live Session Views ────────────────────────────────────────────────────────

@login_required
@require_POST
def start_live_session(request, pk):
    """Start a live classroom session; returns JSON with the join code."""
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)
    from .models import Presentation, LiveSession
    teacher = get_object_or_404(Teacher, user=request.user)
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    LiveSession.objects.filter(presentation=deck, is_active=True).update(is_active=False)
    code = ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))
    while LiveSession.objects.filter(code=code).exists():
        code = ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))
    session = LiveSession.objects.create(
        presentation=deck,
        code=code,
        is_active=True,
        current_slide_order=0,
    )
    return JsonResponse({'code': code, 'session_id': session.pk})


@login_required
@require_POST
def end_live_session(request, pk):
    """End all active live sessions for this deck."""
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)
    from .models import Presentation, LiveSession
    teacher = get_object_or_404(Teacher, user=request.user)
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    LiveSession.objects.filter(presentation=deck, is_active=True).update(is_active=False)
    return JsonResponse({'ok': True})


@login_required
@require_POST
def update_live_slide(request, pk):
    """Update the current slide index of the active live session."""
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)
    import json as _json
    from .models import Presentation, LiveSession
    teacher = get_object_or_404(Teacher, user=request.user)
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    try:
        data = _json.loads(request.body)
        slide_order = int(data.get('slide_order', 0))
    except Exception:
        return JsonResponse({'error': 'Invalid data'}, status=400)
    LiveSession.objects.filter(presentation=deck, is_active=True).update(current_slide_order=slide_order)
    return JsonResponse({'ok': True})


def live_student(request, code):
    """Student join page — no authentication required."""
    from .models import LiveSession
    session = get_object_or_404(LiveSession, code=code, is_active=True)
    return render(request, 'teachers/presentations/live_student.html', {
        'session': session,
        'deck': session.presentation,
        'code': code,
    })


def live_state(request, code):
    """JSON: returns current session state (active slide, content)."""
    from .models import LiveSession

    def _extract_student_notes(raw_notes):
        text = str(raw_notes or '').strip()
        if not text:
            return ''

        # Remove quiz answer prefix if present.
        lines = text.splitlines()
        if lines and lines[0].strip().upper().startswith('ANSWER:'):
            text = '\n'.join(lines[1:]).strip()
            if not text:
                return ''

        lower = text.lower()
        teacher_tag = 'teacher cue:'
        student_tag = 'student notes:'
        student_idx = lower.find(student_tag)

        if student_idx != -1:
            return text[student_idx + len(student_tag):].strip()

        # If explicitly structured but missing student section, hide teacher-only content.
        if lower.find(teacher_tag) != -1:
            return ''

        # Legacy plain notes are treated as student-visible.
        return text

    try:
        session = LiveSession.objects.get(code=code)
    except LiveSession.DoesNotExist:
        return JsonResponse({'is_active': False})
    if not session.is_active:
        return JsonResponse({'is_active': False})
    slides = list(session.presentation.slides.all())
    current_slide = None
    for s in slides:
        if s.order == session.current_slide_order:
            current_slide = s
            break
    if current_slide is None and slides:
        current_slide = slides[0]
    state = {
        'is_active': True,
        'current_slide_order': session.current_slide_order,
        'participant_count': session.responses.values('student_name').distinct().count(),
    }
    if current_slide:
        state.update({
            'slide_pk': current_slide.pk,
            'layout': current_slide.layout,
            'title': current_slide.title,
            'content': current_slide.content,
            'emoji': current_slide.emoji,
            'image_url': current_slide.image_url,
            'student_notes': _extract_student_notes(current_slide.speaker_notes),
        })
    return JsonResponse(state)


@require_POST
def live_vote(request, code):
    """Student submits a poll/quiz vote. No auth required."""
    import json as _json
    from .models import LiveSession, Slide, PollResponse
    try:
        data = _json.loads(request.body)
    except Exception:
        return JsonResponse({'error': 'Invalid JSON'}, status=400)
    try:
        session = LiveSession.objects.get(code=code, is_active=True)
    except LiveSession.DoesNotExist:
        return JsonResponse({'error': 'Session not found or ended'}, status=404)
    slide_pk = data.get('slide_pk')
    choice = str(data.get('choice', '')).upper()[:1]
    student_name = str(data.get('student_name', 'Anonymous'))[:100] or 'Anonymous'
    if not slide_pk or choice not in 'ABCD':
        return JsonResponse({'error': 'Invalid vote data'}, status=400)
    try:
        slide = Slide.objects.get(pk=slide_pk)
    except Slide.DoesNotExist:
        return JsonResponse({'error': 'Slide not found'}, status=404)
    PollResponse.objects.update_or_create(
        session=session, slide=slide, student_name=student_name,
        defaults={'choice': choice},
    )
    return JsonResponse({'ok': True})


def live_results(request, code, slide_pk):
    """Return live poll result counts for a given slide."""
    from .models import LiveSession, PollResponse
    try:
        session = LiveSession.objects.get(code=code)
    except LiveSession.DoesNotExist:
        return JsonResponse({'error': 'Session not found'}, status=404)
    responses = PollResponse.objects.filter(session=session, slide_id=slide_pk)
    counts = {'A': 0, 'B': 0, 'C': 0, 'D': 0}
    for r in responses:
        if r.choice in counts:
            counts[r.choice] += 1
    total = sum(counts.values())
    return JsonResponse({'counts': counts, 'total': total})


@require_POST
def log_slide_time(request, code):
    """Student reports seconds spent on a slide. No authentication required."""
    import json as _json
    from .models import LiveSession
    try:
        data = _json.loads(request.body)
    except Exception:
        return JsonResponse({'error': 'Invalid JSON'}, status=400)
    slide_pk = str(data.get('slide_pk', '') or '').strip()
    try:
        seconds = int(data.get('seconds', 0))
    except (TypeError, ValueError):
        seconds = 0
    # Silently ignore bad data rather than erroring on the student
    if not slide_pk or seconds <= 0 or seconds > 3600:
        return JsonResponse({'ok': True})
    try:
        session = LiveSession.objects.get(code=code)
    except LiveSession.DoesNotExist:
        return JsonResponse({'ok': True})
    time_data = session.slide_time_data or {}
    time_data[slide_pk] = time_data.get(slide_pk, 0) + seconds
    session.slide_time_data = time_data
    session.save(update_fields=['slide_time_data'])
    return JsonResponse({'ok': True})


def _detect_slide_layout(title, index, total):
    """Pick a layout based on slide title keywords and position."""
    if index == 0:
        return 'title'
    if index == total - 1:
        return 'summary'
    t = title.lower()
    if any(w in t for w in ('practice', 'question', 'quiz', 'check your', 'check understanding')):
        return 'quiz'
    if any(w in t for w in ('vocabulary', 'definition', 'key term', 'key vocab')):
        return 'two_col'
    if any(w in t for w in ('comparison', 'compare', 'vs', 'versus', 'did you know')):
        return 'two_col'
    if any(w in t for w in ('quote', 'saying', 'proverb')):
        return 'quote'
    if any(w in t for w in ('stat', 'figure', 'number', 'big stat')):
        return 'big_stat'
    if any(w in t for w in ('poll', 'vote', 'survey')):
        return 'poll'
    return 'bullets'


@login_required
@require_POST
def presentation_api(request):
    """
    AJAX API for the slide editor.  Actions:
      save_slide    — update a single slide's fields
      add_slide     — append a new blank slide
      delete_slide  — remove a slide, renumber
      reorder       — save new slide order
      ai_generate   — generate slides from topic via AuraGenEngine
      update_deck   — rename / retheme deck
    """
    import json
    from .models import Presentation, Slide
    from django.db import transaction

    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)
    teacher = get_object_or_404(Teacher, user=request.user)

    try:
        data   = json.loads(request.body)
        action = data.get('action')
        deck_id = data.get('deck_id')
        deck = get_object_or_404(Presentation, pk=deck_id, teacher=teacher)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=400)

    # ── save_slide ──────────────────────────────────────────────────────────
    if action == 'save_slide':
        slide_id = data.get('slide_id')
        slide = get_object_or_404(Slide, pk=slide_id, presentation=deck)
        slide.title         = data.get('title',         slide.title)
        slide.content       = data.get('content',       slide.content)
        slide.speaker_notes = data.get('speaker_notes', slide.speaker_notes)
        slide.layout        = data.get('layout',        slide.layout)
        slide.emoji         = data.get('emoji',         slide.emoji)
        slide.image_url     = data.get('image_url',     slide.image_url)
        slide.save()
        deck.save()   # bump updated_at
        return JsonResponse({'ok': True})

    # ── add_slide ───────────────────────────────────────────────────────────
    elif action == 'add_slide':
        max_order = deck.slides.aggregate(m=models.Max('order'))['m'] or -1
        slide = Slide.objects.create(
            presentation=deck,
            order=max_order + 1,
            layout=data.get('layout', 'bullets'),
            title=data.get('title', 'New Slide'),
            content=data.get('content', ''),
            emoji=data.get('emoji', ''),
        )
        deck.save()
        return JsonResponse({
            'ok':      True,
            'slide_id': slide.pk,
            'order':   slide.order,
            'layout':  slide.layout,
            'title':   slide.title,
            'content': slide.content,
            'emoji':   slide.emoji,
            'speaker_notes': slide.speaker_notes,
            'image_url': slide.image_url,
        })

    # ── delete_slide ────────────────────────────────────────────────────────
    elif action == 'delete_slide':
        slide_id = data.get('slide_id')
        slide = get_object_or_404(Slide, pk=slide_id, presentation=deck)
        slide.delete()
        # renumber
        with transaction.atomic():
            for i, s in enumerate(deck.slides.order_by('order')):
                s.order = i
                s.save(update_fields=['order'])
        deck.save()
        return JsonResponse({'ok': True})

    # ── reorder ─────────────────────────────────────────────────────────────
    elif action == 'reorder':
        order_list = data.get('order', [])   # list of slide PKs in new order
        with transaction.atomic():
            for i, sid in enumerate(order_list):
                Slide.objects.filter(pk=sid, presentation=deck).update(order=i)
        deck.save()
        return JsonResponse({'ok': True})

    # ── update_deck ─────────────────────────────────────────────────────────
    elif action == 'update_deck':
        if 'title' in data:
            deck.title = data['title'] or deck.title
        if 'theme' in data and data['theme'] in dict(Presentation.THEME_CHOICES):
            deck.theme = data['theme']
        if 'transition' in data and data['transition'] in dict(Presentation.TRANSITION_CHOICES):
            deck.transition = data['transition']
        deck.save()
        return JsonResponse({'ok': True, 'title': deck.title, 'theme': deck.theme, 'transition': deck.transition})

    # ── ai_generate ─────────────────────────────────────────────────────────
    elif action == 'ai_generate':
        ok, err = check_freemium_limit(request.user, 'aura-slide-generator', 'slide_gen')
        if not ok:
            return JsonResponse(err, status=403)
        topic      = data.get('topic', '').strip()
        subject_id = data.get('subject_id')
        class_id   = data.get('class_id')
        if not topic:
            return JsonResponse({'error': 'Topic is required'}, status=400)

        from academics.models import Subject, Class
        subject_name = 'General Studies'
        class_name   = 'General'
        if subject_id:
            s = Subject.objects.filter(pk=subject_id).first()
            if s: subject_name = s.name
        if class_id:
            c = Class.objects.filter(pk=class_id).first()
            if c: class_name = c.name

        from tenants.ai_quota import check_and_consume, QuotaExceeded
        try:
            check_and_consume(request.tenant, request.user.id, 'slide_gen')
        except QuotaExceeded as e:
            return JsonResponse({'error': e.user_message, 'error_code': 'quota_exceeded', 'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost}, status=429)
        from teachers.services.aura_gen_engine import AuraGenEngine
        result = AuraGenEngine.generate_slides_outline(topic, subject_name, class_name)
        raw_slides = result.get('slides', [])

        # Replace all existing slides with AI-generated ones
        with transaction.atomic():
            deck.slides.all().delete()
            created = []
            for i, s in enumerate(raw_slides):
                bullets = s.get('bullets', [])
                content = '\n'.join(bullets)
                layout  = _detect_slide_layout(s.get('title', ''), i, len(raw_slides))
                slide   = Slide.objects.create(
                    presentation=deck,
                    order=i,
                    layout=layout,
                    title=s.get('title', ''),
                    content=content,
                    speaker_notes=s.get('notes', ''),
                    emoji=s.get('emoji', ''),
                )
                created.append({
                    'slide_id':      slide.pk,
                    'order':         slide.order,
                    'layout':        slide.layout,
                    'title':         slide.title,
                    'content':       slide.content,
                    'emoji':         slide.emoji,
                    'speaker_notes': slide.speaker_notes,
                })
            if data.get('update_title') and topic:
                deck.title = topic
            deck.save()

        return JsonResponse({
            'ok':         True,
            'slides':     created,
            'deck_title': deck.title,
            'activities': result.get('activities', []),
        })

    # ── duplicate_slide ─────────────────────────────────────────────────────
    elif action == 'duplicate_slide':
        slide_id = data.get('slide_id')
        source = get_object_or_404(Slide, pk=slide_id, presentation=deck)
        max_order = deck.slides.aggregate(m=models.Max('order'))['m'] or -1
        new_slide = Slide.objects.create(
            presentation=deck,
            order=max_order + 1,
            layout=source.layout,
            title=source.title + ' (copy)',
            content=source.content,
            speaker_notes=source.speaker_notes,
            emoji=source.emoji,
            image_url=source.image_url,
        )
        deck.save()
        return JsonResponse({
            'ok': True,
            'slide_id': new_slide.pk,
            'order':    new_slide.order,
            'layout':   new_slide.layout,
            'title':    new_slide.title,
            'content':  new_slide.content,
            'emoji':    new_slide.emoji,
            'speaker_notes': new_slide.speaker_notes,
            'image_url': new_slide.image_url,
        })

    # ── from_lesson_plan ─────────────────────────────────────────────────────────
    elif action == 'from_lesson_plan':
        ok, err = check_freemium_limit(request.user, 'aura-slide-generator', 'slide_gen')
        if not ok:
            return JsonResponse(err, status=403)
        plan_id = data.get('plan_id')
        if not plan_id:
            return JsonResponse({'error': 'plan_id is required'}, status=400)
        plan = get_object_or_404(LessonPlan, pk=plan_id, teacher=teacher)
        b7_meta = plan.b7_meta if isinstance(plan.b7_meta, dict) else {}
        indicator = (
            (b7_meta.get('indicator') if isinstance(b7_meta, dict) else '')
            or (b7_meta.get('perf_indicator') if isinstance(b7_meta, dict) else '')
            or ''
        )

        school_info = SchoolInfo.objects.first()
        demographic_bits = []
        if school_info:
            if school_info.name:
                demographic_bits.append(f"School: {school_info.name}")
            if school_info.address:
                demographic_bits.append(f"Address: {school_info.address}")
            if school_info.motto:
                demographic_bits.append(f"Motto: {school_info.motto}")
        if teacher.city:
            demographic_bits.append(f"Teacher city: {teacher.city}")
        if teacher.region:
            demographic_bits.append(f"Teacher region: {teacher.get_region_display()}")
        if teacher.hometown:
            demographic_bits.append(f"Teacher hometown: {teacher.hometown}")
        if teacher.preferred_language:
            demographic_bits.append(f"Preferred language: {teacher.get_preferred_language_display()}")

        plan_dict = {
            'topic':        plan.topic,
            'subject':      plan.subject.name if plan.subject else 'General',
            'class_name':   plan.school_class.name if plan.school_class else 'General',
            'week':         plan.week_number,
            'indicator':    indicator,
            'objectives':   plan.objectives,
            'introduction': plan.introduction,
            'presentation': plan.presentation,
            'evaluation':   plan.evaluation,
            'homework':     plan.homework,
            'demographic_context': '; '.join([x for x in demographic_bits if x]),
        }
        from teachers.services.aura_gen_engine import AuraGenEngine
        from tenants.ai_quota import check_and_consume, QuotaExceeded
        try:
            check_and_consume(request.tenant, request.user.id, 'slide_gen')
        except QuotaExceeded as e:
            return JsonResponse({'error': e.user_message, 'error_code': 'quota_exceeded', 'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost}, status=429)
        result = AuraGenEngine.generate_slides_from_lesson_plan(plan_dict)
        raw_slides = result.get('slides', [])
        with transaction.atomic():
            deck.slides.all().delete()
            created = []
            for i, s in enumerate(raw_slides):
                bullets = s.get('bullets', [])
                content = '\n'.join(bullets)
                layout = _detect_slide_layout(s.get('title', ''), i, len(raw_slides))
                slide = Slide.objects.create(
                    presentation=deck, order=i, layout=layout,
                    title=s.get('title', ''), content=content,
                    speaker_notes=s.get('notes', ''),
                    emoji=s.get('emoji', ''),
                )
                created.append({
                    'slide_id': slide.pk, 'order': slide.order, 'layout': slide.layout,
                    'title': slide.title, 'content': slide.content,
                    'emoji': slide.emoji, 'speaker_notes': slide.speaker_notes,
                })
            deck.title = plan.topic
            deck.save()
        return JsonResponse({
            'ok': True, 'slides': created,
            'deck_title': deck.title, 'activities': result.get('activities', []),
        })

    # ── suggest_bullets ──────────────────────────────────────────────────────
    elif action == 'suggest_bullets':
        ok, err = check_freemium_limit(request.user, 'aura-slide-generator', 'slide_gen')
        if not ok:
            return JsonResponse(err, status=403)
        title = data.get('title', '').strip()
        if not title:
            return JsonResponse({'error': 'title is required'}, status=400)
        subject_name = data.get('subject', 'General')
        from tenants.ai_quota import check_and_consume, QuotaExceeded
        try:
            check_and_consume(request.tenant, request.user.id, 'slide_gen')
        except QuotaExceeded as e:
            return JsonResponse({'error': e.user_message, 'error_code': 'quota_exceeded', 'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost}, status=429)
        from teachers.services.aura_gen_engine import AuraGenEngine
        result = AuraGenEngine.suggest_slide_bullets(title, subject_name)
        return JsonResponse({'ok': True, 'bullets': result.get('bullets', [])})

    # ── suggest_layouts ──────────────────────────────────────────────────────
    elif action == 'suggest_layouts':
        ok, err = check_freemium_limit(request.user, 'aura-slide-generator', 'slide_gen')
        if not ok:
            return JsonResponse(err, status=403)
        slides_qs = deck.slides.order_by('order')
        slides_data = [
            {'slide_id': s.pk, 'order': s.order, 'title': s.title, 'content': s.content}
            for s in slides_qs
        ]
        if not slides_data:
            return JsonResponse({'error': 'No slides in deck.'}, status=400)
        from tenants.ai_quota import check_and_consume, QuotaExceeded
        try:
            check_and_consume(request.tenant, request.user.id, 'slide_gen')
        except QuotaExceeded as e:
            return JsonResponse({'error': e.user_message, 'error_code': 'quota_exceeded', 'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost}, status=429)
        from teachers.services.aura_gen_engine import AuraGenEngine
        result = AuraGenEngine.suggest_slide_layouts(slides_data)
        updates = result.get('updates', [])
        with transaction.atomic():
            for upd in updates:
                sid    = upd.get('slide_id')
                layout = upd.get('layout', '')
                if sid and layout:
                    Slide.objects.filter(pk=sid, presentation=deck).update(layout=layout)
            deck.save()
        return JsonResponse({'ok': True, 'updates': updates, 'count': len(updates)})

    # ── harmonize_deck ──────────────────────────────────────────────────────
    elif action == 'harmonize_deck':
        ok, err = check_freemium_limit(request.user, 'aura-slide-generator', 'slide_gen')
        if not ok:
            return JsonResponse(err, status=403)
        slides_qs = deck.slides.order_by('order')
        slides_data = [
            {'slide_id': s.pk, 'title': s.title, 'content': s.content}
            for s in slides_qs
        ]
        if not slides_data:
            return JsonResponse({'error': 'No slides in deck.'}, status=400)
        from teachers.services.aura_gen_engine import AuraGenEngine
        result = AuraGenEngine.harmonize_deck(slides_data)
        # Return updates for frontend preview/confirmation — user applies explicitly
        return JsonResponse({
            'ok':      True,
            'updates': result.get('updates', []),
            'summary': result.get('summary', ''),
        })

    # ── apply_template ───────────────────────────────────────────────────────
    elif action == 'apply_template':
        template_id = data.get('template_id', '')
        TEMPLATES = {
            'lesson_intro': {'slides': [
                {'order':0,'layout':'title',   'emoji':'📚','title':'[Lesson Title]','content':'[Subject] · [Class] · [Teacher]','speaker_notes':''},
                {'order':1,'layout':'bullets', 'emoji':'🎯','title':'Learning Objectives','content':'By the end of this lesson, students will be able to...\nUnderstand the key concepts of [topic]\nApply [skill] in real-world situations\nAnalyse and evaluate [concept]','speaker_notes':''},
                {'order':2,'layout':'bullets', 'emoji':'🔤','title':'Key Vocabulary','content':'[Term 1]: definition here\n[Term 2]: definition here\n[Term 3]: definition here\n[Term 4]: definition here','speaker_notes':''},
                {'order':3,'layout':'two_col', 'emoji':'📋','title':"Today's Agenda",'content':"First Half\nWarm-up activity\nDirect instruction\nGuided practice\nSecond Half\nIndependent practice\nGroup discussion\nExit ticket",'speaker_notes':''},
                {'order':4,'layout':'quote',   'emoji':'💭','title':'','content':'Think about it: [provocative question related to the topic]\n— [Your Name]','speaker_notes':''},
                {'order':5,'layout':'summary', 'emoji':'📝','title':'Lesson Recap','content':'Key concept reviewed\nObjectives met\nActivity completed\nExit ticket collected\nHomework assigned','speaker_notes':''},
            ]},
            'quiz': {'slides': [
                {'order':0,'layout':'title',    'emoji':'🏆','title':'Quiz Time!','content':'[Topic] Review · [Class]','speaker_notes':''},
                {'order':1,'layout':'big_stat', 'emoji':'❓','title':'Question 1','content':'[Q1 — type your question here]\n[Optional hint or sub-text]','speaker_notes':'Read the question aloud twice.'},
                {'order':2,'layout':'bullets',  'emoji':'✅','title':'Answer 1','content':'[Correct answer]\nExplanation: [why this is correct]\nKey takeaway: [one sentence summary]','speaker_notes':''},
                {'order':3,'layout':'big_stat', 'emoji':'❓','title':'Question 2','content':'[Q2 — type your question here]\n[Optional hint or sub-text]','speaker_notes':''},
                {'order':4,'layout':'bullets',  'emoji':'✅','title':'Answer 2','content':'[Correct answer]\nExplanation: [why this is correct]\nKey takeaway: [one sentence summary]','speaker_notes':''},
                {'order':5,'layout':'summary',  'emoji':'🎉','title':'Quiz Complete!','content':'Great effort!\nReview your answers\nAsk questions\nApply what you learned\nSee you next time!','speaker_notes':''},
            ]},
            'timeline': {'slides': [
                {'order':0,'layout':'title',   'emoji':'⏳','title':'[Event] Timeline','content':'From [Start Year] to [End Year]','speaker_notes':''},
                {'order':1,'layout':'bullets', 'emoji':'📌','title':'[Year/Period 1]','content':'[Event name]\nWhat happened: [brief description]\nKey figures: [names]\nImpact: [consequence]','speaker_notes':''},
                {'order':2,'layout':'bullets', 'emoji':'📌','title':'[Year/Period 2]','content':'[Event name]\nWhat happened: [brief description]\nKey figures: [names]\nImpact: [consequence]','speaker_notes':''},
                {'order':3,'layout':'bullets', 'emoji':'📌','title':'[Year/Period 3]','content':'[Event name]\nWhat happened: [brief description]\nKey figures: [names]\nImpact: [consequence]','speaker_notes':''},
                {'order':4,'layout':'quote',   'emoji':'💬','title':'','content':'[Famous quote from the period]\n[Person, Year]','speaker_notes':''},
                {'order':5,'layout':'summary', 'emoji':'📖','title':'What We Learned','content':'Events were connected\nCauses lead to effects\nKey figures shaped history\nLessons are still relevant today','speaker_notes':''},
            ]},
            'debate': {'slides': [
                {'order':0,'layout':'title',   'emoji':'⚖️','title':'Debate: [Motion]','content':'This House Believes That [statement]','speaker_notes':''},
                {'order':1,'layout':'bullets', 'emoji':'✅','title':'Arguments FOR','content':'Point 1: [main argument]\nEvidence: [data or example]\nPoint 2: [main argument]\nEvidence: [data or example]\nPoint 3: [main argument]','speaker_notes':''},
                {'order':2,'layout':'bullets', 'emoji':'❌','title':'Arguments AGAINST','content':'Point 1: [counter-argument]\nEvidence: [data or example]\nPoint 2: [counter-argument]\nEvidence: [data or example]\nPoint 3: [counter-argument]','speaker_notes':''},
                {'order':3,'layout':'two_col', 'emoji':'🔄','title':'Rebuttals','content':'FOR Rebuttal\nAddress: [counter-point]\nResponse: [why it fails]\nAGAINST Rebuttal\nAddress: [for-point]\nResponse: [why it fails]','speaker_notes':''},
                {'order':4,'layout':'quote',   'emoji':'💭','title':'','content':'[Relevant quote about the topic]\n[Source]','speaker_notes':''},
                {'order':5,'layout':'summary', 'emoji':'🏁','title':'Verdict & Takeaways','content':'Key learning points\nStrong arguments considered\nEvidence was critical\nOpen to different perspectives','speaker_notes':''},
            ]},
            'lab_report': {'slides': [
                {'order':0,'layout':'title',    'emoji':'🔬','title':'[Experiment Title]','content':'[Subject] · [Class] · [Date]','speaker_notes':''},
                {'order':1,'layout':'bullets',  'emoji':'💡','title':'Aim & Hypothesis','content':'Aim: [what you are trying to find out]\nHypothesis: I predict that...\nBecause: [scientific reasoning]\nVariables — Independent: [x], Dependent: [y], Controlled: [list]','speaker_notes':''},
                {'order':2,'layout':'bullets',  'emoji':'⚗️','title':'Materials & Method','content':'Step 1: [first step]\nStep 2: [second step]\nStep 3: [third step]\nStep 4: [fourth step]\nSafety: [precautions]','speaker_notes':''},
                {'order':3,'layout':'big_stat', 'emoji':'📊','title':'Key Result','content':'[Main finding, e.g. "42°C"]\n[What this result means]','speaker_notes':''},
                {'order':4,'layout':'bullets',  'emoji':'🔍','title':'Analysis','content':'The results show that...\nThis supports / contradicts the hypothesis because...\nSources of error: [list]\nImprovements: [suggestions]','speaker_notes':''},
                {'order':5,'layout':'summary',  'emoji':'📝','title':'Conclusion','content':'Hypothesis was supported / rejected\nKey evidence collected\nMethod could be improved\nFurther research needed on [topic]','speaker_notes':''},
            ]},
            'book_review': {'slides': [
                {'order':0,'layout':'title',   'emoji':'📖','title':'[Book Title]','content':'By [Author] · Review by [Student/Class]','speaker_notes':''},
                {'order':1,'layout':'bullets', 'emoji':'📋','title':'Plot Summary','content':'Setting: [time and place]\nOpening: [how the story begins]\nMiddle: [main events]\nClimax: [turning point]\nEnding: [resolution]','speaker_notes':''},
                {'order':2,'layout':'two_col', 'emoji':'👥','title':'Key Characters','content':'[Protagonist]\nName: [character name]\nRole: [what they do]\nTrait: [key personality]\n[Antagonist]\nName: [character name]\nRole: [what they do]\nTrait: [key personality]','speaker_notes':''},
                {'order':3,'layout':'bullets', 'emoji':'💡','title':'Major Themes','content':'Theme 1: [e.g. Friendship] — [how it appears]\nTheme 2: [e.g. Courage] — [how it appears]\nTheme 3: [e.g. Identity] — [how it appears]','speaker_notes':''},
                {'order':4,'layout':'quote',   'emoji':'💬','title':'','content':'[Favourite quote from the book]\n— [Character/Narrator, Page/Chapter]','speaker_notes':''},
                {'order':5,'layout':'big_stat','emoji':'⭐','title':'My Rating','content':'[e.g. 4/5 Stars]\n[One sentence verdict]','speaker_notes':''},
            ]},
            'chapter': {'slides': [
                {'order':0,'layout':'title',    'emoji':'📑','title':'Chapter [#]: [Title]','content':'[Subject] · [Class]','speaker_notes':''},
                {'order':1,'layout':'bullets',  'emoji':'🔁','title':'What We Covered','content':'[Key point 1]\n[Key point 2]\n[Key point 3]\n[Key point 4]\n[Key point 5]','speaker_notes':''},
                {'order':2,'layout':'two_col',  'emoji':'🧠','title':'Key Concepts','content':'[Concept 1]\n[Definition]\n[Example]\n[Concept 2]\n[Definition]\n[Example]','speaker_notes':''},
                {'order':3,'layout':'big_stat', 'emoji':'✏️','title':'Worked Example','content':'[Formula, equation, or method]\n[Show the process here]','speaker_notes':''},
                {'order':4,'layout':'summary',  'emoji':'📝','title':'Practice & Review','content':'Attempt past questions\nReview class notes\nCreate a mind map\nPractise with a partner\nAsk your teacher for help','speaker_notes':''},
            ]},
            'project': {'slides': [
                {'order':0,'layout':'title',   'emoji':'🚀','title':'[Project Name]','content':'[Team / Class] · [Date]','speaker_notes':''},
                {'order':1,'layout':'bullets', 'emoji':'🌍','title':'Background & Problem','content':'The problem we are solving: [description]\nWhy it matters: [impact]\nCurrent situation: [context]\nOur approach: [brief overview]','speaker_notes':''},
                {'order':2,'layout':'bullets', 'emoji':'🎯','title':'Goals & Objectives','content':'Goal 1: [measurable objective]\nGoal 2: [measurable objective]\nGoal 3: [measurable objective]\nSuccess metric: [how we will measure success]','speaker_notes':''},
                {'order':3,'layout':'two_col', 'emoji':'🔧','title':'Method & Tools','content':'Methods\n[Method 1]\n[Method 2]\n[Method 3]\nTools & Resources\n[Tool 1]\n[Tool 2]\n[Tool 3]','speaker_notes':''},
                {'order':4,'layout':'bullets', 'emoji':'📅','title':'Timeline','content':'Week 1–2: Planning\nWeek 3–4: Research / Development\nWeek 5–6: Implementation\nWeek 7: Testing / Review\nWeek 8: Presentation / Submission','speaker_notes':''},
                {'order':5,'layout':'summary', 'emoji':'📦','title':'Deliverables','content':'Final report submitted\nPresentation delivered\nPrototype ready\nTeam evaluation complete\nLessons learned documented','speaker_notes':''},
            ]},
        }
        if template_id not in TEMPLATES:
            return JsonResponse({'error': f'Unknown template: {template_id}'}, status=400)
        tpl_data = TEMPLATES[template_id]
        with transaction.atomic():
            deck.slides.all().delete()
            new_slides = []
            for s_data in tpl_data['slides']:
                sl = Slide.objects.create(
                    presentation=deck,
                    order=s_data['order'],
                    layout=s_data['layout'],
                    emoji=s_data.get('emoji', ''),
                    title=s_data['title'],
                    content=s_data['content'],
                    speaker_notes=s_data.get('speaker_notes', ''),
                )
                new_slides.append(sl)
        return JsonResponse({
            'ok': True,
            'slides': [{
                'slide_id': sl.pk,
                'order':    sl.order,
                'layout':   sl.layout,
                'title':    sl.title,
                'content':  sl.content,
                'emoji':    sl.emoji,
                'speaker_notes': sl.speaker_notes,
                'image_url': sl.image_url,
            } for sl in new_slides],
        })

    # ── refine_slide ─────────────────────────────────────────────────────────
    elif action == 'refine_slide':
        ok, err = check_freemium_limit(request.user, 'aura-slide-generator', 'slide_gen')
        if not ok:
            return JsonResponse(err, status=403)
        slide_id    = data.get('slide_id')
        instruction = data.get('instruction', '').strip()
        slide       = get_object_or_404(Slide, pk=slide_id, presentation=deck)
        from teachers.services.aura_gen_engine import AuraGenEngine
        subject_name = deck.subject.name if deck.subject else 'General'
        result = AuraGenEngine.refine_slide(
            {'title': slide.title, 'content': slide.content, 'layout': slide.layout},
            instruction,
            subject_name,
        )
        slide.title   = result['title']
        slide.content = result['content']
        slide.save()
        deck.save()
        return JsonResponse({'ok': True, 'title': result['title'], 'content': result['content']})

    return JsonResponse({'error': f'Unknown action: {action}'}, status=400)


@login_required
@requires_addon_freemium('aura-slide-generator', action_type='slide_gen')
def presentation_generate_from_doc(request):
    """
    Accepts a multipart POST with:
      deck_id   — int
      document  — .pdf or .docx file
    Extracts the text, calls AuraGenEngine, replaces deck slides, returns JSON.
    """
    from .models import Presentation, Slide

    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)

    teacher = get_object_or_404(Teacher, user=request.user)
    deck_id = request.POST.get('deck_id')
    uploaded = request.FILES.get('document')

    if not deck_id or not uploaded:
        return JsonResponse({'error': 'deck_id and document file are required'}, status=400)

    deck = get_object_or_404(Presentation, pk=deck_id, teacher=teacher)

    filename = uploaded.name.lower()
    if not (filename.endswith('.pdf') or filename.endswith('.docx') or filename.endswith('.doc')):
        return JsonResponse({'error': 'Only .pdf and .docx files are supported'}, status=400)

    # ── Extract text ────────────────────────────────────────────────────────
    try:
        if filename.endswith('.pdf'):
            import pypdf
            reader = pypdf.PdfReader(uploaded)
            text_parts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
            document_text = '\n'.join(text_parts)
        else:
            from docx import Document as DocxDocument
            doc = DocxDocument(uploaded)
            document_text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as exc:
        return JsonResponse({'error': f'Could not read file: {exc}'}, status=400)

    if len(document_text.strip()) < 80:
        return JsonResponse({'error': 'The document appears to be empty or unreadable. Please try a different file.'}, status=400)

    # ── Generate slides ──────────────────────────────────────────────────────
    from teachers.services.aura_gen_engine import AuraGenEngine
    from django.db import transaction

    result = AuraGenEngine.generate_slides_from_document(document_text, uploaded.name)
    raw_slides = result.get('slides', [])

    with transaction.atomic():
        deck.slides.all().delete()
        created = []
        for i, s in enumerate(raw_slides):
            bullets = s.get('bullets', [])
            content = '\n'.join(bullets)
            layout = _detect_slide_layout(s.get('title', ''), i, len(raw_slides))
            slide = Slide.objects.create(
                presentation=deck,
                order=i,
                layout=layout,
                title=s.get('title', ''),
                content=content,
                speaker_notes=s.get('notes', ''),
                emoji=s.get('emoji', ''),
            )
            created.append({
                'slide_id':      slide.pk,
                'order':         slide.order,
                'layout':        slide.layout,
                'title':         slide.title,
                'content':       slide.content,
                'emoji':         slide.emoji,
                'speaker_notes': slide.speaker_notes,
            })
        deck.save()

    return JsonResponse({
        'ok':         True,
        'slides':     created,
        'deck_title': deck.title,
        'activities': result.get('activities', []),
    })


@login_required
@requires_addon_freemium('aura-slide-generator', action_type='slide_gen')
def presentation_from_youtube(request):
    """
    POST JSON: {deck_id: int, youtube_url: str}
    Extracts a YouTube transcript, then generates a slide deck via AuraGenEngine.
    """
    import json as _json
    import re
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)
    try:
        body = _json.loads(request.body)
    except (ValueError, _json.JSONDecodeError):
        return JsonResponse({'error': 'Invalid JSON body'}, status=400)

    from .models import Presentation, Slide
    teacher  = get_object_or_404(Teacher, user=request.user)
    deck_id  = body.get('deck_id')
    yt_url   = (body.get('youtube_url') or '').strip()

    if not deck_id or not yt_url:
        return JsonResponse({'error': 'deck_id and youtube_url are required'}, status=400)

    deck = get_object_or_404(Presentation, pk=deck_id, teacher=teacher)

    # Extract video ID from standard URLs, shortened links, and Shorts
    m = re.search(r'(?:v=|youtu\.be/|shorts/|embed/)([a-zA-Z0-9_-]{11})', yt_url)
    if not m:
        return JsonResponse(
            {'error': 'Could not read video ID from URL. Paste a standard YouTube link.'},
            status=400,
        )
    video_id = m.group(1)

    # Fetch transcript
    try:
        from youtube_transcript_api import (
            YouTubeTranscriptApi,
            TranscriptsDisabled,
            NoTranscriptFound,
        )
        entries = YouTubeTranscriptApi.get_transcript(video_id)
        transcript_text = ' '.join(e['text'] for e in entries)
    except TranscriptsDisabled:
        return JsonResponse(
            {'error': 'This video has transcripts disabled. Try a different video.'},
            status=400,
        )
    except NoTranscriptFound:
        return JsonResponse(
            {'error': 'No transcript found for this video (it may not have captions).'},
            status=400,
        )
    except Exception as exc:
        return JsonResponse({'error': f'Could not fetch transcript: {exc}'}, status=400)

    if len(transcript_text.strip()) < 100:
        return JsonResponse({'error': 'Transcript too short to generate slides.'}, status=400)

    # Best-effort fetch of video title
    import html as _html
    import requests as _req
    video_title = f'YouTube Video'
    try:
        resp = _req.get(
            f'https://www.youtube.com/watch?v={video_id}',
            headers={'User-Agent': 'Mozilla/5.0'},
            timeout=6,
        )
        tm = re.search(r'<title>(.*?) - YouTube</title>', resp.text)
        if tm:
            video_title = _html.unescape(tm.group(1))
    except Exception:
        pass

    # Generate slides from transcript
    from teachers.services.aura_gen_engine import AuraGenEngine
    from django.db import transaction
    result     = AuraGenEngine.generate_slides_from_document(transcript_text, video_title)
    raw_slides = result.get('slides', [])

    with transaction.atomic():
        deck.slides.all().delete()
        created = []
        for i, s in enumerate(raw_slides):
            content = '\n'.join(s.get('bullets', []))
            layout  = _detect_slide_layout(s.get('title', ''), i, len(raw_slides))
            slide   = Slide.objects.create(
                presentation=deck, order=i, layout=layout,
                title=s.get('title', ''), content=content,
                speaker_notes=s.get('notes', ''),
                emoji=s.get('emoji', ''),
            )
            created.append({
                'slide_id':      slide.pk,
                'order':         slide.order,
                'layout':        slide.layout,
                'title':         slide.title,
                'content':       slide.content,
                'emoji':         slide.emoji,
                'speaker_notes': slide.speaker_notes,
                'image_url':     slide.image_url,
            })
        deck.title = video_title
        deck.save()

    return JsonResponse({
        'ok':         True,
        'slides':     created,
        'deck_title': deck.title,
        'activities': result.get('activities', []),
    })


@login_required
def presentation_lesson_plans(request):
    """Return teacher's lesson plans as JSON for the slide deck AI modal."""
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)
    teacher = get_object_or_404(Teacher, user=request.user)
    plans = LessonPlan.objects.filter(teacher=teacher).select_related(
        'subject', 'school_class'
    ).order_by('-date_added')[:60]
    data = [{
        'id':         p.pk,
        'topic':      p.topic,
        'week':       p.week_number,
        'subject':    p.subject.name if p.subject else '',
        'class_name': p.school_class.name if p.school_class else '',
    } for p in plans]
    return JsonResponse({'plans': data})


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE DECK v1.9 — Session Analytics, Share Links, Image Upload
# ═══════════════════════════════════════════════════════════════════════════════

def _build_time_data(session, deck):
    """Return a sorted list of {slide, seconds, duration, pct} dicts for time-on-slide display."""
    raw = session.slide_time_data or {}
    entries = []
    for pk_str, secs in raw.items():
        if not secs or secs <= 0:
            continue
        try:
            slide_obj = deck.slides.filter(pk=int(pk_str)).first()
        except (ValueError, TypeError):
            continue
        if not slide_obj:
            continue
        m, s = divmod(int(secs), 60)
        entries.append({
            'slide':    slide_obj,
            'seconds':  secs,
            'duration': f'{m}m {s}s' if m else f'{s}s',
            'pct':      0,
        })
    entries.sort(key=lambda x: x['slide'].order)
    if entries:
        max_s = max(e['seconds'] for e in entries)
        for e in entries:
            e['pct'] = round(e['seconds'] / max_s * 100) if max_s else 0
    return entries


@login_required
def presentation_session_report(request, pk):
    """Post-session analytics: all past live sessions for a deck with per-slide poll results."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation, LiveSession, PollResponse

    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    sessions = (
        LiveSession.objects.filter(presentation=deck)
        .prefetch_related('responses__slide')
        .order_by('-created_at')
    )

    # Build per-session analytics
    session_data = []
    for sess in sessions:
        participants = sess.responses.values('student_name').distinct().count()
        # Poll slides that received at least one vote
        slide_stats = []
        slide_ids_with_votes = (
            sess.responses.values_list('slide_id', flat=True).distinct()
        )
        for slide_id in slide_ids_with_votes:
            slide = sess.presentation.slides.filter(pk=slide_id).first()
            if not slide or slide.layout not in ('poll', 'quiz'):
                continue
            counts = {'A': 0, 'B': 0, 'C': 0, 'D': 0}
            for resp in sess.responses.filter(slide_id=slide_id):
                if resp.choice in counts:
                    counts[resp.choice] += 1
            total = sum(counts.values())
            leading = max(counts, key=counts.get) if total else None
            slide_stats.append({
                'slide': slide,
                'total': total,
                'leading': leading,
                'choices': [
                    {
                        'letter': letter,
                        'count': counts[letter],
                        'pct': round(counts[letter] / total * 100) if total else 0,
                    }
                    for letter in 'ABCD'
                ],
            })
        session_data.append({
            'session': sess,
            'participants': participants,
            'slide_stats': slide_stats,
            'time_data': _build_time_data(sess, deck),
        })

    return render(request, 'teachers/presentations/session_report.html', {
        'deck': deck,
        'session_data': session_data,
    })


def presentation_share(request, token):
    """Public read-only slideshow view — no login required."""
    from .models import Presentation
    deck = get_object_or_404(Presentation, share_token=token)
    slides = list(deck.slides.all())
    return render(request, 'teachers/presentations/share_view.html', {
        'deck': deck,
        'slides': slides,
    })


@login_required
@require_POST
def presentation_slide_image_upload(request):
    """Accept a multipart image upload and return its served URL."""
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)

    image_file = request.FILES.get('image')
    if not image_file:
        return JsonResponse({'error': 'No image provided'}, status=400)

    # Validate content type
    allowed_types = ('image/jpeg', 'image/png', 'image/gif', 'image/webp')
    if image_file.content_type not in allowed_types:
        return JsonResponse({'error': 'Invalid file type. Only JPEG, PNG, GIF and WebP are accepted.'}, status=400)

    # Limit to 5 MB
    if image_file.size > 5 * 1024 * 1024:
        return JsonResponse({'error': 'File too large. Maximum size is 5 MB.'}, status=400)

    import os
    import uuid as _uuid
    from django.core.files.storage import default_storage
    from django.core.files.base import ContentFile

    ext = os.path.splitext(image_file.name)[1].lower() or '.jpg'
    filename = f"slide_images/{_uuid.uuid4().hex}{ext}"
    path = default_storage.save(filename, ContentFile(image_file.read()))
    url = request.build_absolute_uri(default_storage.url(path))
    return JsonResponse({'url': url})


@login_required
@require_POST
@requires_addon_freemium('aura-slide-generator', action_type='slide_gen')
def generate_slide_image(request):
    """Generate an AI image for a slide using OpenAI DALL\u00b7E."""
    if request.user.user_type != 'teacher':
        return JsonResponse({'status': 'error', 'message': 'Access denied'}, status=403)

    slide_id = request.POST.get('slide_id')
    prompt = request.POST.get('prompt', '').strip()
    if not slide_id or not prompt:
        return JsonResponse({'status': 'error', 'message': 'Missing slide_id or prompt.'}, status=400)

    from .models import Slide
    slide = Slide.objects.filter(pk=slide_id).first()
    if not slide:
        return JsonResponse({'status': 'error', 'message': 'Slide not found.'}, status=404)

    import requests as _requests
    api_key = os.environ.get('OPENAI_API_KEY')
    if not api_key:
        return JsonResponse({'status': 'error', 'message': 'OpenAI API key not configured.'}, status=500)

    try:
        dalle_url = 'https://api.openai.com/v1/images/generations'
        payload = {
            'model': 'dall-e-3',
            'prompt': prompt,
            'n': 1,
            'size': '1024x1024',
        }
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json',
        }
        resp = _requests.post(dalle_url, json=payload, headers=headers, timeout=60)
        resp.raise_for_status()
        data = resp.json()
        image_url = data['data'][0]['url'] if 'data' in data and data['data'] else ''
        if not image_url:
            return JsonResponse({'status': 'error', 'message': 'No image returned.'}, status=500)
        slide.image_url = image_url
        slide.save()
        return JsonResponse({'status': 'success', 'image_url': image_url})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': f'Image generation failed: {str(e)}'}, status=500)


@login_required
def presentation_study_guide(request, pk):
    from .models import Presentation
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)

    if request.method == 'POST':
        import json as _json
        action = request.POST.get('action') or ''
        if action == 'generate_ai':
            ok, err = check_freemium_limit(request.user, 'aura-slide-generator', 'study_guide')
            if not ok:
                return JsonResponse(err, status=403)
            slides = list(deck.slides.values('title', 'content', 'layout', 'speaker_notes', 'emoji'))
            from teachers.services.aura_gen_engine import AuraGenEngine
            from tenants.ai_quota import check_and_consume, QuotaExceeded
            try:
                check_and_consume(request.tenant, request.user.id, 'study_guide')
            except QuotaExceeded as e:
                return JsonResponse({'error': e.user_message, 'error_code': 'quota_exceeded', 'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost}, status=429)
            result = AuraGenEngine.generate_study_guide(slides)
            return JsonResponse(result)
        return JsonResponse({'error': 'Unknown action'}, status=400)

    slides = deck.slides.all()
    from academics.models import Class, AcademicYear
    current_year = AcademicYear.objects.filter(is_current=True).first()
    classes = Class.objects.filter(academic_year=current_year) if current_year else Class.objects.none()
    return render(request, 'teachers/presentations/study_guide.html', {
        'deck': deck,
        'slides': slides,
        'classes': classes,
    })


@login_required
@require_POST
def presentation_pulse_launch(request, pk):
    """
    Launch a Pulse Check session from AI-generated study guide questions.
    POST JSON: {q1, q2, q3, chips: [...]}
    Requires the presentation to have a school_class assigned.
    """
    import json as _json
    from .models import Presentation
    from academics.pulse_models import PulseSession

    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)

    teacher = get_object_or_404(Teacher, user=request.user)
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)

    if not deck.school_class_id:
        return JsonResponse(
            {'error': 'This presentation has no class assigned. Edit the presentation to set a class first.'},
            status=400)

    try:
        body = _json.loads(request.body)
    except (ValueError, _json.JSONDecodeError):
        return JsonResponse({'error': 'Invalid JSON body'}, status=400)

    q1    = str(body.get('q1', '')).strip()[:500]
    q2    = str(body.get('q2', '')).strip()[:500]
    q3    = str(body.get('q3', '')).strip()[:500]
    chips = body.get('chips', [])
    if not isinstance(chips, list):
        chips = []
    chips = [str(c)[:200] for c in chips[:6]]

    if not q1 or not q2 or not q3:
        return JsonResponse({'error': 'q1, q2 and q3 are all required'}, status=400)

    # Close any other active sessions targeting the same class
    PulseSession.objects.filter(
        target_class=deck.school_class, status='active'
    ).update(status='closed')
    PulseSession.objects.filter(
        lesson_plan__school_class=deck.school_class, status='active'
    ).update(status='closed')

    session = PulseSession.objects.create(
        lesson_plan=None,
        presentation=deck,
        target_class=deck.school_class,
        teacher=teacher,
        q1_text=q1,
        q2_text=q2,
        q3_text=q3,
        q3_chips=chips or [q3],
    )

    script_name = request.META.get('SCRIPT_NAME', '').rstrip('/')
    from django.urls import reverse
    live_url  = script_name + reverse('teachers:pulse_live',  args=[session.pk])
    close_url = script_name + reverse('teachers:pulse_close', args=[session.pk])

    # Notify students that a live pulse is active for their class
    try:
        from announcements.models import Notification as _Notif
        from students.models import Student as _Stu
        _dashboard_url = script_name + '/dashboard/'
        _students = _Stu.objects.filter(
            current_class=deck.school_class, user__is_active=True
        ).select_related('user')
        topic_label = deck.title[:60] if deck.title else 'Study Guide'
        _Notif.objects.bulk_create([
            _Notif(
                recipient=s.user,
                message=f"🟢 Live Pulse! Check your dashboard to respond — {topic_label}",
                alert_type='general',
                link=_dashboard_url,
            )
            for s in _students
        ], ignore_conflicts=True)
    except Exception:
        pass

    return JsonResponse({
        'ok':             True,
        'session_id':     session.pk,
        'total_students': session.total_students,
        'live_url':       live_url,
        'close_url':      close_url,
    })


@login_required
@require_POST
def presentation_send_as_assignment(request, pk):
    """
    Create a Homework assignment from an AI-generated study guide.
    POST JSON: {class_id, due_date, note, questions: [{question, type, options, answer}], summary}
    """
    import json as _json
    from .models import Presentation
    from homework.models import Homework, Question, Choice
    from academics.models import Class

    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)

    teacher = get_object_or_404(Teacher, user=request.user)
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)

    try:
        body = _json.loads(request.body)
    except (ValueError, _json.JSONDecodeError):
        return JsonResponse({'error': 'Invalid JSON body'}, status=400)

    class_id  = body.get('class_id')
    due_date  = body.get('due_date', '')
    note      = str(body.get('note', '')).strip()[:1000]
    questions = body.get('questions', [])
    summary   = str(body.get('summary', '')).strip()[:2000]

    if not class_id:
        return JsonResponse({'error': 'class_id is required'}, status=400)
    if not due_date:
        return JsonResponse({'error': 'due_date is required'}, status=400)
    if not questions:
        return JsonResponse({'error': 'No questions to assign'}, status=400)

    target_class = get_object_or_404(Class, pk=class_id)

    description_parts = []
    if summary:
        description_parts.append(summary)
    if note:
        description_parts.append(f"\nTeacher note: {note}")
    description_parts.append(f"\n(Generated from AI Study Guide: {deck.title})")
    description = "\n".join(description_parts).strip()

    hw = Homework.objects.create(
        title=f"{deck.title} — Study Guide Assignment",
        description=description,
        teacher=teacher,
        subject=deck.subject,
        target_class=target_class,
        due_date=due_date,
    )

    for q_data in questions[:20]:  # cap at 20 questions
        q_text = str(q_data.get('question', '')).strip()[:500]
        if not q_text:
            continue
        q_type_raw = str(q_data.get('type', 'short')).lower()
        if q_type_raw == 'mcq':
            q_type = 'mcq'
        elif q_type_raw in ('essay', 'long'):
            q_type = 'essay'
        else:
            q_type = 'short'

        correct_answer = str(q_data.get('answer', '')).strip()[:500]
        question = Question.objects.create(
            homework=hw,
            text=q_text,
            question_type=q_type,
            correct_answer=correct_answer,
            points=1,
        )

        if q_type == 'mcq':
            options = q_data.get('options', [])
            if isinstance(options, list):
                for opt in options[:6]:
                    opt_text = str(opt).strip()[:200]
                    if not opt_text:
                        continue
                    is_correct = opt_text.lower() == correct_answer.lower()
                    Choice.objects.create(question=question, text=opt_text, is_correct=is_correct)

    script_name = request.META.get('SCRIPT_NAME', '').rstrip('/')
    from django.urls import reverse as _reverse
    hw_url = script_name + _reverse('homework:homework_detail', args=[hw.pk])

    # Bulk-notify all students in the target class
    try:
        from announcements.models import Notification as _Notif
        from students.models import Student as _Stu
        _students = _Stu.objects.filter(
            current_class=target_class, user__is_active=True
        ).select_related('user')
        _Notif.objects.bulk_create([
            _Notif(
                recipient=s.user,
                message=f"📚 New Assignment: {hw.title}",
                alert_type='general',
                link=hw_url,
            )
            for s in _students
        ], ignore_conflicts=True)
    except Exception:
        pass  # Never block homework creation over a notification failure

    return JsonResponse({'ok': True, 'homework_id': hw.pk, 'homework_url': hw_url})


@login_required
@require_POST
def presentation_send_as_notes(request, pk):
    """
    Share AI-generated key concepts as a ClassNote for the target class.
    POST JSON: {class_id, concepts: [{term, description}, ...], title}
    """
    import json as _json
    from .models import Presentation
    from homework.models import ClassNote
    from academics.models import Class

    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)

    teacher = get_object_or_404(Teacher, user=request.user)
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)

    try:
        body = _json.loads(request.body)
    except (ValueError, _json.JSONDecodeError):
        return JsonResponse({'error': 'Invalid JSON body'}, status=400)

    class_id = body.get('class_id')
    concepts = body.get('concepts', [])
    title = str(body.get('title', f"{deck.title} — Key Concepts")).strip()[:200]

    if not class_id:
        return JsonResponse({'error': 'class_id is required'}, status=400)
    if not concepts:
        return JsonResponse({'error': 'No concepts to send'}, status=400)

    target_class = get_object_or_404(Class, pk=class_id)

    # Sanitise concepts list
    clean_concepts = []
    for c in concepts[:30]:
        term = str(c.get('term', '')).strip()[:150]
        desc = str(c.get('description', '')).strip()[:400]
        if term:
            clean_concepts.append({'term': term, 'description': desc})

    note = ClassNote.objects.create(
        title=title,
        content=clean_concepts,
        teacher=teacher,
        target_class=target_class,
        source_deck=deck,
    )

    # Notify students in the class
    try:
        from announcements.models import Notification as _Notif
        from students.models import Student as _Stu
        from django.urls import reverse as _reverse
        script_name = request.META.get('SCRIPT_NAME', '').rstrip('/')
        notes_url = _reverse('homework:student_notes_list') + f'?note={note.pk}'
        _students = _Stu.objects.filter(
            current_class=target_class, user__is_active=True
        ).select_related('user')
        _Notif.objects.bulk_create([
            _Notif(
                recipient=s.user,
                message=f"📖 New Class Note: {note.title}",
                alert_type='general',
                link=notes_url,
            )
            for s in _students
        ], ignore_conflicts=True)
    except Exception:
        pass

    notes_url = _reverse('homework:student_notes_list') + f'?note={note.pk}'
    return JsonResponse({'ok': True, 'note_id': note.pk, 'notes_url': notes_url})


@login_required
def presentation_export_pptx(request, pk):
    """Export a presentation deck as a .pptx file."""
    if request.user.user_type not in ('teacher', 'admin'):
        messages.error(request, 'Access denied.')
        return redirect('dashboard')
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation
    deck = get_object_or_404(Presentation, pk=pk, teacher=teacher)
    slides_qs = list(deck.slides.all())

    try:
        from pptx import Presentation as PPTXPres
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io as _io
    except ImportError:
        messages.error(request, 'python-pptx is not installed. Run: pip install python-pptx')
        return redirect('teachers:presentation_editor', pk=pk)

    THEME_COLORS = {
        'aurora':   {'bg': (30, 10, 64),   'accent': (124, 58, 237)},
        'midnight': {'bg': (15, 23, 42),   'accent': (99, 102, 241)},
        'forest':   {'bg': (2, 29, 17),    'accent': (52, 211, 153)},
        'coral':    {'bg': (59, 0, 18),    'accent': (251, 113, 133)},
        'slate':    {'bg': (15, 24, 36),   'accent': (148, 163, 184)},
        'ocean':    {'bg': (12, 45, 78),   'accent': (56, 189, 248)},
        'amber':    {'bg': (67, 20, 7),    'accent': (251, 191, 36)},
        'rose':     {'bg': (76, 5, 25),    'accent': (244, 63, 94)},
    }
    theme  = deck.theme if deck.theme in THEME_COLORS else 'aurora'
    colors = THEME_COLORS[theme]
    bg_rgb      = RGBColor(*colors['bg'])
    accent_rgb  = RGBColor(*colors['accent'])
    text_rgb    = RGBColor(241, 245, 249)    # #f1f5f9
    subtext_rgb = RGBColor(148, 163, 184)    # #94a3b8

    prs = PPTXPres()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    W = prs.slide_width
    H = prs.slide_height
    MARGIN = Inches(0.65)

    def _set_bg(slide, rgb):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb

    def _textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, italic=False,
                 color=None, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text            = text
        run.font.size       = Pt(font_size)
        run.font.bold       = bold
        run.font.italic     = italic
        run.font.color.rgb  = color if color is not None else text_rgb
        return txb

    def _bullets_frame(slide, items, left, top, width, height, font_size=18):
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text           = '\u25b8  ' + item
            p.font.size      = Pt(font_size)
            p.font.color.rgb = subtext_rgb

    for s in slides_qs:
        sl = prs.slides.add_slide(blank_layout)
        _set_bg(sl, bg_rgb)
        layout  = s.layout
        title   = s.title or ''
        content = s.content or ''
        bullets = [b.strip() for b in content.split('\n') if b.strip()]
        notes   = s.speaker_notes or ''

        if layout == 'title':
            _textbox(sl, title, MARGIN, Inches(2.0), W - 2*MARGIN, Inches(2.2),
                     font_size=44, bold=True, align=PP_ALIGN.CENTER)
            sub = bullets[0] if bullets else content.strip()
            if sub:
                _textbox(sl, sub, MARGIN, Inches(4.3), W - 2*MARGIN, Inches(1.3),
                         font_size=22, color=subtext_rgb, align=PP_ALIGN.CENTER)

        elif layout == 'big_stat':
            _textbox(sl, title, MARGIN, Inches(0.5), W - 2*MARGIN, Inches(0.9),
                     font_size=18, color=subtext_rgb, align=PP_ALIGN.CENTER)
            stat_val = bullets[0] if bullets else content.strip()
            _textbox(sl, stat_val, MARGIN, Inches(1.4), W - 2*MARGIN, Inches(3.6),
                     font_size=96, bold=True, color=accent_rgb, align=PP_ALIGN.CENTER)
            if len(bullets) > 1:
                _textbox(sl, bullets[1], MARGIN, Inches(5.1), W - 2*MARGIN, Inches(0.9),
                         font_size=16, color=subtext_rgb, align=PP_ALIGN.CENTER)

        elif layout == 'quote':
            _textbox(sl, '\u201c' + title + '\u201d', MARGIN, Inches(1.4), W - 2*MARGIN, Inches(3.8),
                     font_size=32, italic=True, align=PP_ALIGN.CENTER)
            author = bullets[0] if bullets else ''
            if author:
                _textbox(sl, '\u2014 ' + author, MARGIN, Inches(5.5), W - 2*MARGIN, Inches(0.9),
                         font_size=16, color=subtext_rgb, align=PP_ALIGN.CENTER)

        elif layout == 'two_col':
            _textbox(sl, title, MARGIN, MARGIN, W - 2*MARGIN, Inches(0.95), font_size=28, bold=True)
            mid   = W / 2
            col_w = mid - MARGIN - Inches(0.15)
            half  = max(len(bullets) // 2, 1)
            _bullets_frame(sl, bullets[:half],  MARGIN,               Inches(1.55), col_w, H - Inches(2.4))
            _bullets_frame(sl, bullets[half:],  mid + Inches(0.15),   Inches(1.55), col_w, H - Inches(2.4))

        elif layout in ('poll', 'quiz'):
            _textbox(sl, title, MARGIN, MARGIN, W - 2*MARGIN, Inches(1.3), font_size=30, bold=True)
            y = Inches(1.65)
            for b in bullets[:6]:
                _textbox(sl, b, MARGIN, y, W - 2*MARGIN, Inches(0.8), font_size=18, color=subtext_rgb)
                y += Inches(0.92)

        elif layout == 'video':
            _textbox(sl, title, MARGIN, MARGIN, W - 2*MARGIN, Inches(1.1), font_size=32, bold=True)
            note = content.strip() or '(no URL)'
            _textbox(sl, '\U0001f3ac  Video: ' + note, MARGIN, Inches(2.0), W - 2*MARGIN, Inches(1.5),
                     font_size=14, italic=True, color=subtext_rgb)

        elif layout == 'image':
            _textbox(sl, title, MARGIN, MARGIN, W - 2*MARGIN, Inches(1.1), font_size=32, bold=True)
            img_note = s.image_url[:80] if s.image_url else '(no image)'
            _textbox(sl, '\U0001f5bc\ufe0f  Image: ' + img_note, MARGIN, Inches(2.0), W - 2*MARGIN, Inches(1.5),
                     font_size=14, italic=True, color=subtext_rgb)
            cap = bullets[0] if bullets else ''
            if cap:
                _textbox(sl, cap, MARGIN, Inches(6.2), W - 2*MARGIN, Inches(0.8),
                         font_size=14, color=subtext_rgb, align=PP_ALIGN.CENTER)

        else:  # bullets, summary
            _textbox(sl, title, MARGIN, MARGIN, W - 2*MARGIN, Inches(1.15), font_size=32, bold=True)
            if bullets:
                _bullets_frame(sl, bullets, MARGIN, Inches(1.6), W - 2*MARGIN, H - Inches(2.5))
            elif content.strip():
                _textbox(sl, content.strip(), MARGIN, Inches(1.6), W - 2*MARGIN, H - Inches(2.5),
                         font_size=16, color=subtext_rgb)

        if notes:
            sl.notes_slide.notes_text_frame.text = notes

    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    import re as _re
    safe = _re.sub(r'[^\w\s-]', '', deck.title).strip().replace(' ', '_')[:60] or 'presentation'
    resp = HttpResponse(
        buf.read(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )
    resp['Content-Disposition'] = f'attachment; filename="{safe}.pptx"'
    return resp


@login_required
def presentation_bulk_action(request):
    """Bulk delete or duplicate presentations via JSON POST."""
    import json as _json
    if request.method != 'POST':
        return JsonResponse({'error': 'POST required'}, status=405)
    if request.user.user_type not in ('teacher', 'admin'):
        return JsonResponse({'error': 'Access denied'}, status=403)
    teacher = get_object_or_404(Teacher, user=request.user)
    from .models import Presentation, Slide
    try:
        data   = _json.loads(request.body)
        action = data.get('action')
        pks    = [int(p) for p in data.get('pks', [])]
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=400)
    if not pks:
        return JsonResponse({'error': 'No presentations selected'}, status=400)

    decks = Presentation.objects.filter(pk__in=pks, teacher=teacher)

    if action == 'delete':
        count = decks.count()
        decks.delete()
        return JsonResponse({'ok': True, 'count': count})

    elif action == 'duplicate':
        count = 0
        for deck in decks:
            new_deck = Presentation.objects.create(
                teacher=teacher,
                title=deck.title + ' (copy)',
                theme=deck.theme,
                transition=deck.transition,
                subject=deck.subject,
                school_class=deck.school_class,
            )
            for slide in deck.slides.all():
                Slide.objects.create(
                    presentation=new_deck,
                    order=slide.order, layout=slide.layout,
                    title=slide.title, content=slide.content,
                    speaker_notes=slide.speaker_notes,
                    emoji=slide.emoji, image_url=slide.image_url,
                )
            count += 1
        return JsonResponse({'ok': True, 'count': count})

    return JsonResponse({'error': 'Invalid action'}, status=400)


@login_required
def my_fee_tasks(request):
    """Teacher view: fee structures I am assigned to collect."""
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied.')
        return redirect('dashboard')

    from finance.models import FeeStructure, StudentFee
    from decimal import Decimal as _Decimal
    from django.db.models import Count, Q, Sum
    from django.db.models.functions import Coalesce
    from django.db.utils import ProgrammingError, OperationalError

    try:
        teacher = Teacher.objects.get(user=request.user)
    except Teacher.DoesNotExist:
        messages.error(request, 'Your teacher profile was not found. Please contact an administrator.')
        return redirect('dashboard')
    except (ProgrammingError, OperationalError):
        messages.error(
            request,
            'Fee tasks are temporarily unavailable for this school because teacher tables are not ready yet. '
            'Please ask your administrator to run tenant migrations.'
        )
        return redirect('dashboard')

    try:
        structures = (
            FeeStructure.objects
            .filter(assigned_collector=teacher)
            .select_related('head', 'class_level', 'academic_year')
            .annotate(
                total_students=Count('studentfee', distinct=True),
                paid_students=Count(
                    'studentfee',
                    filter=Q(studentfee__status='paid'),
                    distinct=True,
                ),
                unpaid_students=Count(
                    'studentfee',
                    filter=Q(studentfee__status__in=['unpaid', 'partial']),
                    distinct=True,
                ),
                total_collected=Coalesce(
                    Sum('studentfee__payments__amount'),
                    _Decimal('0'),
                ),
            )
            .order_by('class_level__name', '-id')
        )
    except (ProgrammingError, OperationalError):
        messages.error(
            request,
            'Fee tasks are temporarily unavailable for this school because finance/teacher tables are not ready yet. '
            'Please ask your administrator to run tenant migrations.'
        )
        return redirect('dashboard')

    return render(request, 'teachers/my_fee_tasks.html', {'structures': structures})


@login_required
def fee_task_detail(request, structure_id):
    """Teacher view: students who have paid and those who owe for a specific fee structure."""
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied.')
        return redirect('dashboard')

    from finance.models import FeeStructure, StudentFee
    from decimal import Decimal as _Decimal
    from django.db.models import Sum, Q
    from django.db.models.functions import Coalesce
    from django.db.utils import ProgrammingError, OperationalError
    from django.shortcuts import get_object_or_404

    try:
        teacher = Teacher.objects.get(user=request.user)
    except Teacher.DoesNotExist:
        messages.error(request, 'Your teacher profile was not found.')
        return redirect('dashboard')
    except (ProgrammingError, OperationalError):
        messages.error(request, 'Fee data is temporarily unavailable. Ask your administrator to run tenant migrations.')
        return redirect('dashboard')

    structure = get_object_or_404(
        FeeStructure.objects.select_related('head', 'class_level', 'academic_year', 'assigned_collector'),
        id=structure_id,
        assigned_collector=teacher,
    )

    try:
        student_fees = (
            StudentFee.objects
            .filter(fee_structure=structure)
            .select_related('student__user', 'student__current_class')
            .prefetch_related('payments')
            .annotate(
                total_paid_amount=Coalesce(Sum('payments__amount'), _Decimal('0')),
            )
            .order_by('student__user__last_name', 'student__user__first_name')
        )
    except (ProgrammingError, OperationalError):
        messages.error(request, 'Fee data is temporarily unavailable.')
        return redirect('teachers:my_fee_tasks')

    paid_fees = [sf for sf in student_fees if sf.status == 'paid']
    owing_fees = [sf for sf in student_fees if sf.status in ('unpaid', 'partial')]

    total_collected = sum(sf.total_paid_amount for sf in paid_fees)
    total_outstanding = sum(sf.amount_payable - sf.total_paid_amount for sf in owing_fees)

    return render(request, 'teachers/fee_task_detail.html', {
        'structure': structure,
        'paid_fees': paid_fees,
        'owing_fees': owing_fees,
        'total_collected': total_collected,
        'total_outstanding': total_outstanding,
        'paid_count': len(paid_fees),
        'owing_count': len(owing_fees),
        'total_count': len(list(student_fees)),
    })


# ---------------------------------------------------------------------------
# Teacher Add-On Store
# ---------------------------------------------------------------------------


# Shared mapping of add-on slugs → Django URL names for launching
ADDON_LAUNCH_URLS = {
    'task-board': 'teachers:addon_task_board',
    'cpd-tracker': 'teachers:addon_cpd_tracker',
    'observation-notes': 'teachers:addon_observation_notes',
    'rubric-designer': 'teachers:addon_rubric_designer',
    'study-guide-builder': 'teachers:addon_study_guide',
    'random-picker': 'teachers:addon_random_picker',
    'countdown-timer': 'teachers:addon_countdown_timer',
    'noise-meter': 'teachers:addon_noise_meter',
    'stem-activity-pack': 'teachers:addon_stem_pack',
    'creative-arts-kit': 'teachers:addon_creative_arts',
    'aura-slide-generator': 'teachers:presentation_list',
    'smart-planner-pro': 'teachers:aura_command_center',
    'exercise-maker': 'teachers:my_classes',
    'grade-insight-dashboard': 'teachers:analytics_dashboard',
    'quick-report-writer': 'teachers:lesson_plan_list',
    'report-card-writer': 'teachers:addon_report_card',
    'exam-question-bank': 'teachers:addon_question_bank',
    'behavior-sel-tracker': 'teachers:addon_behavior_tracker',
    'differentiated-lesson-ai': 'teachers:addon_differentiated',
    'live-quiz-engine': 'teachers:addon_live_quiz',
}


@login_required
def teacher_store(request):
    """Browse & manage personal teacher add-ons."""
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    from teachers.models import TeacherAddOn, TeacherAddOnPurchase

    addons = TeacherAddOn.objects.filter(is_active=True)
    purchased_ids = set(
        TeacherAddOnPurchase.objects.filter(
            teacher=request.user, is_active=True
        ).values_list('addon_id', flat=True)
    )

    # Group by category + set launch URL for owned add-ons
    categories = {}
    for addon in addons:
        addon.is_purchased = addon.id in purchased_ids
        if addon.is_purchased:
            url_name = ADDON_LAUNCH_URLS.get(addon.slug)
            if url_name:
                try:
                    addon.launch_url = reverse(url_name)
                except Exception:
                    addon.launch_url = None
            else:
                addon.launch_url = None
        label = addon.get_category_display()
        categories.setdefault(label, []).append(addon)

    return render(request, 'teachers/teacher_store.html', {
        'categories': categories,
        'purchased_count': len(purchased_ids),
        'total_count': addons.count(),
    })


@login_required
def teacher_store_purchase(request, addon_id):
    """Purchase (activate) a teacher add-on — free add-ons activate instantly,
    paid ones require Paystack verification via teacher_store_verify."""
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    if request.method != 'POST':
        return redirect('teachers:teacher_store')

    from teachers.models import TeacherAddOn, TeacherAddOnPurchase

    addon = get_object_or_404(TeacherAddOn, id=addon_id, is_active=True)

    if addon.is_free or addon.price <= 0:
        # Free add-ons: activate immediately
        obj, created = TeacherAddOnPurchase.objects.get_or_create(
            teacher=request.user,
            addon=addon,
            defaults={'is_active': True},
        )
        if not created and not obj.is_active:
            obj.is_active = True
            obj.save(update_fields=['is_active'])
        messages.success(request, f'"{addon.name}" activated!')
        return redirect('teachers:teacher_store')

    # Paid add-ons: return JSON with Paystack params for inline popup
    from django.conf import settings
    import uuid
    ref = f"TA-{request.user.id}-{addon.id}-{uuid.uuid4().hex[:8]}"
    return JsonResponse({
        'paystack': True,
        'public_key': settings.PAYSTACK_PUBLIC_KEY,
        'email': request.user.email or f'{request.user.username}@school.local',
        'amount': int(addon.price * 100),  # kobo/pesewas
        'currency': settings.PAYSTACK_CURRENCY,
        'reference': ref,
        'addon_name': addon.name,
    })


@login_required
def teacher_store_verify(request):
    """POST — verify a Paystack payment and activate the add-on."""
    if request.user.user_type != 'teacher':
        return JsonResponse({'ok': False, 'error': 'Access denied'}, status=403)
    if request.method != 'POST':
        return JsonResponse({'ok': False, 'error': 'POST only'}, status=405)

    import requests as http_requests
    from django.conf import settings
    from teachers.models import TeacherAddOn, TeacherAddOnPurchase

    body = json.loads(request.body) if request.content_type == 'application/json' else request.POST
    reference = body.get('reference', '')
    addon_id = body.get('addon_id')

    if not reference or not addon_id:
        return JsonResponse({'ok': False, 'error': 'Missing data'}, status=400)

    addon = get_object_or_404(TeacherAddOn, id=addon_id, is_active=True)

    # Verify with Paystack API
    secret = settings.PAYSTACK_SECRET_KEY
    if secret:
        resp = http_requests.get(
            f'https://api.paystack.co/transaction/verify/{reference}',
            headers={'Authorization': f'Bearer {secret}'},
            timeout=15,
        )
        data = resp.json()
        if not data.get('status') or data.get('data', {}).get('status') != 'success':
            return JsonResponse({'ok': False, 'error': 'Payment verification failed'}, status=402)
        amount_paid = Decimal(str(data['data']['amount'])) / 100
    else:
        # Dev mode: no Paystack key configured — accept reference at face value
        amount_paid = addon.price

    obj, created = TeacherAddOnPurchase.objects.get_or_create(
        teacher=request.user,
        addon=addon,
        defaults={
            'is_active': True,
            'payment_reference': reference,
            'amount_paid': amount_paid,
        },
    )
    if not created:
        obj.is_active = True
        obj.payment_reference = reference
        obj.amount_paid = amount_paid
        obj.save(update_fields=['is_active', 'payment_reference', 'amount_paid'])

    return JsonResponse({'ok': True, 'addon_name': addon.name})


@login_required
def teacher_store_cancel(request, addon_id):
    """Deactivate a purchased teacher add-on."""
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    if request.method != 'POST':
        return redirect('teachers:teacher_store')

    from teachers.models import TeacherAddOn, TeacherAddOnPurchase

    addon = get_object_or_404(TeacherAddOn, id=addon_id)
    TeacherAddOnPurchase.objects.filter(
        teacher=request.user, addon=addon
    ).update(is_active=False)

    messages.info(request, f'"{addon.name}" deactivated.')
    return redirect('teachers:teacher_store')


# ── Paystack Webhook (server-to-server) ────────────────────────────
from django.views.decorators.csrf import csrf_exempt
import hashlib
import hmac


@csrf_exempt
def paystack_teacher_webhook(request):
    """Paystack sends charge.success events here. HMAC-verified, idempotent."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST only'}, status=405)

    secret = getattr(settings, 'PAYSTACK_SECRET_KEY', '')
    if not secret:
        return JsonResponse({'error': 'not configured'}, status=503)

    # Verify HMAC-SHA512 signature
    sig = request.headers.get('X-Paystack-Signature', '')
    expected = hmac.new(secret.encode(), request.body, hashlib.sha512).hexdigest()
    if not hmac.compare_digest(sig, expected):
        return JsonResponse({'error': 'bad signature'}, status=403)

    try:
        payload = json.loads(request.body)
    except (json.JSONDecodeError, ValueError):
        return JsonResponse({'error': 'bad json'}, status=400)

    event = payload.get('event')
    data = payload.get('data', {})

    if event == 'charge.success':
        reference = data.get('reference', '')
        # Teacher add-on references start with "TA-"
        if reference.startswith('TA-'):
            _handle_teacher_addon_payment(reference, data)

    return JsonResponse({'ok': True})


def _handle_teacher_addon_payment(reference, data):
    """Activate a teacher add-on from a confirmed Paystack charge."""
    from teachers.models import TeacherAddOn, TeacherAddOnPurchase
    from decimal import Decimal

    # Parse reference: TA-{user_id}-{addon_id}-{uuid}
    parts = reference.split('-')
    if len(parts) < 3:
        logger.warning("Webhook: bad teacher addon reference format: %s", reference)
        return
    try:
        user_id = int(parts[1])
        addon_id = int(parts[2])
    except (ValueError, IndexError):
        logger.warning("Webhook: unparseable teacher addon reference: %s", reference)
        return

    from django.contrib.auth import get_user_model
    User = get_user_model()
    try:
        user = User.objects.get(id=user_id, user_type='teacher')
        addon = TeacherAddOn.objects.get(id=addon_id, is_active=True)
    except (User.DoesNotExist, TeacherAddOn.DoesNotExist):
        logger.warning("Webhook: teacher/addon not found for ref %s", reference)
        return

    # Deduplicate by payment_reference
    if TeacherAddOnPurchase.objects.filter(payment_reference=reference).exists():
        logger.info("Webhook: duplicate payment ref %s, skipping", reference)
        return

    amount_paid = Decimal(str(data.get('amount', 0))) / 100

    obj, created = TeacherAddOnPurchase.objects.get_or_create(
        teacher=user,
        addon=addon,
        defaults={
            'is_active': True,
            'payment_reference': reference,
            'amount_paid': amount_paid,
        },
    )
    if not created and not obj.is_active:
        obj.is_active = True
        obj.payment_reference = reference
        obj.amount_paid = amount_paid
        obj.expires_at = None
        obj.save(update_fields=['is_active', 'payment_reference', 'amount_paid', 'expires_at'])
        logger.info("Webhook: re-activated addon %s for user %s (ref %s)", addon.slug, user_id, reference)
    elif created:
        logger.info("Webhook: activated addon %s for user %s (ref %s)", addon.slug, user_id, reference)


# ── My Add-Ons ─────────────────────────────────────────────────────

@login_required
def my_addons(request):
    """Teacher's personal add-on dashboard — active, expired, purchase history."""
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')

    from teachers.models import TeacherAddOn, TeacherAddOnPurchase
    from django.utils import timezone

    now = timezone.now()
    purchases = (
        TeacherAddOnPurchase.objects
        .filter(teacher=request.user)
        .select_related('addon')
        .order_by('-purchased_at')
    )

    active = []
    expired = []
    inactive = []
    for p in purchases:
        if p.is_active and (p.expires_at is None or p.expires_at > now):
            p.status_label = 'Active'
            if p.expires_at:
                p.days_left = (p.expires_at - now).days
            active.append(p)
        elif p.is_active and p.expires_at and p.expires_at <= now:
            p.status_label = 'Expired'
            expired.append(p)
        else:
            p.status_label = 'Deactivated'
            inactive.append(p)

    # Map addon slug → URL name so "Launch" links work
    from django.urls import reverse
    for p in active:
        url_name = ADDON_LAUNCH_URLS.get(p.addon.slug)
        if url_name:
            try:
                p.launch_url = reverse(url_name)
            except Exception:
                p.launch_url = None
        else:
            p.launch_url = None

    total_spent = sum(p.amount_paid for p in purchases)
    total_boost = sum(p.addon.quota_boost for p in active if p.addon.quota_boost)

    return render(request, 'teachers/my_addons.html', {
        'active': active,
        'expired': expired,
        'inactive': inactive,
        'total_spent': total_spent,
        'total_active': len(active),
        'total_boost': total_boost,
    })


# ── Start Trial ────────────────────────────────────────────────────

@login_required
def teacher_store_trial(request, addon_id):
    """Start a free trial for a paid add-on that supports trial_days > 0."""
    if request.user.user_type != 'teacher':
        messages.error(request, 'Access denied')
        return redirect('dashboard')
    if request.method != 'POST':
        return redirect('teachers:teacher_store')

    from teachers.models import TeacherAddOn, TeacherAddOnPurchase
    from django.utils import timezone
    from datetime import timedelta

    addon = get_object_or_404(TeacherAddOn, id=addon_id, is_active=True)

    if addon.trial_days <= 0:
        messages.warning(request, f'"{addon.name}" does not offer a free trial.')
        return redirect('teachers:teacher_store')

    # Check if already purchased or trialed
    existing = TeacherAddOnPurchase.objects.filter(
        teacher=request.user, addon=addon
    ).first()
    if existing:
        if existing.is_active:
            messages.info(request, f'You already have "{addon.name}".')
        else:
            messages.warning(request, f'You have already used a trial for "{addon.name}". Purchase to re-activate.')
        return redirect('teachers:teacher_store')

    TeacherAddOnPurchase.objects.create(
        teacher=request.user,
        addon=addon,
        is_active=True,
        expires_at=timezone.now() + timedelta(days=addon.trial_days),
        payment_reference=f'TRIAL-{request.user.id}-{addon.id}',
        amount_paid=0,
    )
    messages.success(request, f'🎉 Free {addon.trial_days}-day trial for "{addon.name}" activated!')
    return redirect('teachers:teacher_store')


# ---------------------------------------------------------------------------
# Add-On Feature Views
# ---------------------------------------------------------------------------

@login_required
@requires_addon('task-board')
def addon_task_board(request):
    """Kanban-style task board for teachers."""
    from teachers.models import TaskCard
    teacher = request.user
    COLUMNS = [('todo', 'To Do'), ('doing', 'In Progress'), ('done', 'Done')]

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'add':
            TaskCard.objects.create(
                teacher=teacher,
                title=request.POST.get('title', '').strip()[:200] or 'Untitled',
                column='todo',
                priority=request.POST.get('priority', 'medium'),
                due_date=request.POST.get('due_date') or None,
            )
        elif action == 'move':
            card = get_object_or_404(TaskCard, id=request.POST.get('card_id'), teacher=teacher)
            new_col = request.POST.get('column', 'todo')
            if new_col in dict(COLUMNS):
                card.column = new_col
                card.save(update_fields=['column'])
        elif action == 'delete':
            TaskCard.objects.filter(id=request.POST.get('card_id'), teacher=teacher).delete()
        return redirect('teachers:addon_task_board')

    cards = TaskCard.objects.filter(teacher=teacher)
    columns = [(key, label, cards.filter(column=key)) for key, label in COLUMNS]
    return render(request, 'teachers/addon_task_board.html', {'columns': columns, 'cards': cards})


@login_required
@requires_addon('cpd-tracker')
def addon_cpd_tracker(request):
    """Continuing Professional Development tracker."""
    from teachers.models import CPDEntry
    teacher = request.user

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'add':
            CPDEntry.objects.create(
                teacher=teacher,
                title=request.POST.get('title', '').strip()[:200] or 'Untitled',
                category=request.POST.get('category', 'other'),
                hours=Decimal(request.POST.get('hours', '1') or '1'),
                date=request.POST.get('date') or date.today(),
                notes=request.POST.get('notes', '').strip(),
            )
        elif action == 'delete':
            CPDEntry.objects.filter(id=request.POST.get('entry_id'), teacher=teacher).delete()
        return redirect('teachers:addon_cpd_tracker')

    entries = CPDEntry.objects.filter(teacher=teacher)
    total_hours = entries.aggregate(t=Sum('hours'))['t'] or 0
    return render(request, 'teachers/addon_cpd_tracker.html', {
        'entries': entries, 'total_hours': total_hours,
    })


@login_required
@requires_addon('observation-notes')
def addon_observation_notes(request):
    """Classroom observation notes for peer review / mentoring."""
    from teachers.models import ObservationNote
    teacher = request.user

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'add':
            cls_id = request.POST.get('observed_class')
            ObservationNote.objects.create(
                teacher=teacher,
                observed_class_id=cls_id if cls_id else None,
                date=request.POST.get('date') or date.today(),
                strengths=request.POST.get('strengths', '').strip(),
                growth_areas=request.POST.get('growth_areas', '').strip(),
                action_plan=request.POST.get('action_plan', '').strip(),
                is_private=request.POST.get('is_private') == 'on',
            )
        elif action == 'delete':
            ObservationNote.objects.filter(id=request.POST.get('note_id'), teacher=teacher).delete()
        return redirect('teachers:addon_observation_notes')

    notes = ObservationNote.objects.filter(teacher=teacher).select_related('observed_class')
    classes = Class.objects.all()
    return render(request, 'teachers/addon_observation_notes.html', {
        'notes': notes, 'classes': classes,
    })


@login_required
@requires_addon('rubric-designer')
def addon_rubric_designer(request):
    """Design assessment rubrics with JSON criteria."""
    from teachers.models import Rubric
    teacher = request.user

    edit_id = request.GET.get('edit')
    rubric = None
    if edit_id:
        rubric = get_object_or_404(Rubric, id=edit_id, teacher=teacher)

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'save':
            title = request.POST.get('title', '').strip()[:200] or 'Untitled Rubric'
            subject_id = request.POST.get('subject') or None
            criteria_json = request.POST.get('criteria', '[]')
            try:
                criteria = json.loads(criteria_json)
            except json.JSONDecodeError:
                criteria = []
            rid = request.POST.get('rubric_id')
            if rid:
                r = get_object_or_404(Rubric, id=rid, teacher=teacher)
                r.title = title
                r.subject_id = subject_id
                r.criteria = criteria
                r.save(update_fields=['title', 'subject', 'criteria', 'updated_at'])
            else:
                Rubric.objects.create(
                    teacher=teacher, title=title, subject_id=subject_id, criteria=criteria,
                )
            return redirect('teachers:addon_rubric_designer')
        elif action == 'delete':
            Rubric.objects.filter(id=request.POST.get('rubric_id'), teacher=teacher).delete()
            return redirect('teachers:addon_rubric_designer')

    subjects = Subject.objects.all()
    rubrics = Rubric.objects.filter(teacher=teacher).select_related('subject')
    return render(request, 'teachers/addon_rubric_designer.html', {
        'subjects': subjects, 'rubrics': rubrics,
        'rubric': rubric, 'editing': rubric is not None,
        'show_list': not rubric,
    })


@login_required
def addon_study_guide(request):
    """Create study guides from lesson content."""
    from teachers.models import StudyGuide
    teacher = request.user

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'save':
            StudyGuide.objects.create(
                teacher=teacher,
                title=request.POST.get('title', '').strip()[:200] or 'Untitled Guide',
                subject_id=request.POST.get('subject') or None,
                target_class_id=request.POST.get('target_class') or None,
                source_notes=request.POST.get('source_notes', '').strip(),
                content_html=request.POST.get('content_html', '').strip(),
            )
            return redirect('teachers:addon_study_guide')
        elif action == 'delete':
            StudyGuide.objects.filter(id=request.POST.get('guide_id'), teacher=teacher).delete()
            return redirect('teachers:addon_study_guide')

    subjects = Subject.objects.all()
    classes = Class.objects.all()
    guides = StudyGuide.objects.filter(teacher=teacher).select_related('subject', 'target_class')
    sg_used = get_free_generation_count(teacher, 'study_guide')
    return render(request, 'teachers/addon_study_guide.html', {
        'subjects': subjects, 'classes': classes, 'guides': guides,
        'sg_gen_used': sg_used,
        'sg_gen_limit': FREE_GENERATION_LIMIT,
        'sg_gen_remaining': max(0, FREE_GENERATION_LIMIT - sg_used),
        'has_sg_addon': has_addon(teacher, 'study-guide-builder'),
    })


@login_required
@requires_addon_freemium('study-guide-builder', action_type='study_guide')
def study_guide_ai(request):
    """AJAX endpoint: generate study guide HTML from teacher notes using AI."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST only'}, status=405)

    body = json.loads(request.body) if request.content_type == 'application/json' else request.POST
    notes = body.get('notes', '').strip()
    if not notes:
        return JsonResponse({'error': 'No notes provided'}, status=400)

    # Use OpenAI / Gemini if available, otherwise return a formatted version
    try:
        from tenants.ai_quota import check_and_consume, QuotaExceeded
        try:
            check_and_consume(request.tenant, request.user.id, 'study_guide')
        except QuotaExceeded as e:
            return JsonResponse({'error': e.user_message, 'error_code': 'quota_exceeded', 'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost}, status=429)

        import google.generativeai as genai
        api_key = os.environ.get('GEMINI_API_KEY', '')
        if api_key:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-2.0-flash')
            prompt = (
                "You are a Ghanaian JHS teacher assistant. Turn the following lesson notes into a "
                "well-structured HTML study guide for students. Use <h3>, <ul>, <li>, <strong>, <p> tags. "
                "Include a summary, key points, and revision questions.\n\n"
                f"NOTES:\n{notes[:3000]}"
            )
            resp = model.generate_content(prompt)
            return JsonResponse({'html': resp.text})
    except Exception:
        pass

    # Fallback: simple formatting
    paragraphs = notes.split('\n\n')
    html = '<h3>Study Guide</h3>'
    html += ''.join(f'<p>{p.strip()}</p>' for p in paragraphs if p.strip())
    html += '<h4>Revision Questions</h4><p><em>Review the notes above and test yourself.</em></p>'
    return JsonResponse({'html': html})


@login_required
@requires_addon('random-picker')
def addon_random_picker(request):
    """Random student picker — template is client-side, just needs class list."""
    classes = Class.objects.all()
    return render(request, 'teachers/addon_random_picker.html', {'classes': classes})


@login_required
@requires_addon('countdown-timer')
def addon_countdown_timer(request):
    """Classroom countdown timer — fully client-side."""
    return render(request, 'teachers/addon_countdown_timer.html')


@login_required
@requires_addon('noise-meter')
def addon_noise_meter(request):
    """Visual noise level meter — uses Web Audio API, fully client-side."""
    return render(request, 'teachers/addon_noise_meter.html')


@login_required
@requires_addon('stem-activity-pack')
def addon_stem_pack(request):
    """Curated STEM activity library — static content pack."""
    return render(request, 'teachers/addon_stem_pack.html')


@login_required
@requires_addon('creative-arts-kit')
def addon_creative_arts(request):
    """Curated creative arts activity library — static content pack."""
    return render(request, 'teachers/addon_creative_arts.html')


# ---------------------------------------------------------------------------
# Wave 2 Add-Ons
# ---------------------------------------------------------------------------

# ── 1. Report Card AI Writer ──────────────────────────────────────────────

@login_required
def addon_report_card(request):
    """Report Card AI Writer — generate end-of-term comments for students."""
    from teachers.models import ReportCardComment

    teacher = request.user
    current_year = AcademicYear.objects.filter(is_current=True).first()
    classes = Class.objects.filter(academic_year=current_year) if current_year else Class.objects.none()

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'save_edit':
            rc = get_object_or_404(ReportCardComment, id=request.POST.get('comment_id'), teacher=teacher)
            rc.comment = request.POST.get('comment', '').strip()
            rc.conduct = request.POST.get('conduct', '').strip()
            rc.attitude = request.POST.get('attitude', '').strip()
            rc.is_edited = True
            rc.save(update_fields=['comment', 'conduct', 'attitude', 'is_edited', 'updated_at'])
            messages.success(request, 'Comment saved.')
            return redirect('teachers:addon_report_card')
        elif action == 'delete':
            ReportCardComment.objects.filter(id=request.POST.get('comment_id'), teacher=teacher).delete()
            return redirect('teachers:addon_report_card')

    sel_class = request.GET.get('class_id')
    sel_term = request.GET.get('term', 'first')
    comments = ReportCardComment.objects.none()
    students_in_class = []

    if sel_class and current_year:
        students_in_class = Student.objects.filter(
            current_class_id=sel_class,
        ).select_related('user').order_by('user__first_name')
        comments = ReportCardComment.objects.filter(
            teacher=teacher, academic_year=current_year, term=sel_term,
            student__in=students_in_class,
        ).select_related('student__user')

    rc_used = get_free_generation_count(teacher, 'report_card')
    return render(request, 'teachers/addon_report_card.html', {
        'classes': classes,
        'sel_class': int(sel_class) if sel_class else None,
        'sel_term': sel_term,
        'comments': comments,
        'students_in_class': students_in_class,
        'current_year': current_year,
        'rc_gen_used': rc_used,
        'rc_gen_limit': FREE_GENERATION_LIMIT,
        'rc_gen_remaining': max(0, FREE_GENERATION_LIMIT - rc_used),
        'has_rc_addon': has_addon(teacher, 'report-card-writer'),
    })


@login_required
@requires_addon_freemium('report-card-writer', action_type='report_card')
def report_card_ai(request):
    """AJAX: AI-generate report card comments for selected students."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST only'}, status=405)

    body = json.loads(request.body) if request.content_type == 'application/json' else request.POST
    student_ids = body.get('student_ids', [])
    term = body.get('term', 'first')
    class_id = body.get('class_id')

    if not student_ids or not class_id:
        return JsonResponse({'error': 'Select students and class.'}, status=400)

    from tenants.ai_quota import check_and_consume, QuotaExceeded
    try:
        check_and_consume(request.tenant, request.user.id, 'report_card')
    except QuotaExceeded as e:
        return JsonResponse({
            'error': e.user_message, 'error_code': 'quota_exceeded',
            'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost,
        }, status=429)

    current_year = AcademicYear.objects.filter(is_current=True).first()
    if not current_year:
        return JsonResponse({'error': 'No active academic year.'}, status=400)

    students = Student.objects.filter(
        id__in=student_ids[:40], current_class_id=class_id,
    ).select_related('user')

    # Build per-student grade summaries
    student_data = []
    for s in students:
        grades = Grade.objects.filter(
            student=s, academic_year=current_year, term=term,
        ).select_related('subject')
        grade_summary = ', '.join(
            f"{g.subject.name}: {g.total_score}/100" for g in grades
        ) or 'No grades recorded'
        avg = grades.aggregate(a=Avg('total_score'))['a'] or 0
        student_data.append({
            'id': s.id, 'name': s.user.get_full_name(),
            'grades': grade_summary, 'avg': round(float(avg), 1),
        })

    # AI generation
    try:
        import google.generativeai as genai
        api_key = os.environ.get('GEMINI_API_KEY', '')
        if api_key:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-2.0-flash')
            students_text = '\n'.join(
                f"- {d['name']}: Average {d['avg']}%. Subjects: {d['grades']}"
                for d in student_data
            )
            prompt = (
                "You are an experienced Ghanaian JHS teacher writing end-of-term report card comments. "
                "For EACH student below, write a personalised 2-3 sentence report comment that:\n"
                "1. Mentions their strengths based on their grades\n"
                "2. Offers constructive areas for improvement\n"
                "3. Ends with encouragement\n"
                "Also suggest a one-word Conduct (e.g. Excellent, Good, Satisfactory) and "
                "Attitude (e.g. Excellent, Very Good, Good).\n\n"
                "Return ONLY valid JSON: [{\"name\":\"…\",\"comment\":\"…\",\"conduct\":\"…\",\"attitude\":\"…\"}, …]\n\n"
                f"STUDENTS:\n{students_text}"
            )
            resp = model.generate_content(prompt)
            text = resp.text.strip()
            # Extract JSON from response
            if '```' in text:
                text = text.split('```')[1]
                if text.startswith('json'):
                    text = text[4:]
            import re
            json_match = re.search(r'\[.*\]', text, re.DOTALL)
            if json_match:
                ai_results = json.loads(json_match.group())
            else:
                ai_results = json.loads(text)

            # Save to database
            from teachers.models import ReportCardComment
            saved = []
            for d in student_data:
                ai_match = next(
                    (r for r in ai_results if r.get('name', '').lower().strip() == d['name'].lower().strip()),
                    None,
                )
                if not ai_match:
                    ai_match = ai_results[student_data.index(d)] if len(ai_results) > student_data.index(d) else {}
                comment_text = ai_match.get('comment', f"{d['name']} has shown effort this term.")
                conduct = ai_match.get('conduct', 'Good')
                attitude = ai_match.get('attitude', 'Good')
                obj, _ = ReportCardComment.objects.update_or_create(
                    student_id=d['id'], academic_year=current_year, term=term,
                    defaults={
                        'teacher': request.user,
                        'comment': comment_text,
                        'conduct': conduct,
                        'attitude': attitude,
                    },
                )
                saved.append({
                    'id': obj.id, 'student_id': d['id'], 'name': d['name'],
                    'comment': comment_text, 'conduct': conduct, 'attitude': attitude,
                })
            return JsonResponse({'results': saved})
    except QuotaExceeded:
        raise
    except Exception as exc:
        logger.exception('Report card AI error')
        return _ai_json_error_response(exc)

    # Fallback
    from teachers.models import ReportCardComment
    saved = []
    for d in student_data:
        cmt = f"{d['name']} achieved an average of {d['avg']}% this term. "
        cmt += "Keep up the hard work." if d['avg'] >= 50 else "More effort is needed next term."
        obj, _ = ReportCardComment.objects.update_or_create(
            student_id=d['id'], academic_year=current_year, term=term,
            defaults={'teacher': request.user, 'comment': cmt, 'conduct': 'Good', 'attitude': 'Good'},
        )
        saved.append({'id': obj.id, 'student_id': d['id'], 'name': d['name'], 'comment': cmt, 'conduct': 'Good', 'attitude': 'Good'})
    return JsonResponse({'results': saved})


# ── 2. Exam & Question Bank Pro ──────────────────────────────────────────

@login_required
def addon_question_bank(request):
    """Question bank management + exam paper assembly."""
    from teachers.models import QuestionBank, ExamPaper

    teacher = request.user
    subjects = Subject.objects.all()
    classes = Class.objects.filter(academic_year__is_current=True)

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'add':
            opts_raw = request.POST.get('options', '')
            opts = [o.strip() for o in opts_raw.split('\n') if o.strip()] if opts_raw.strip() else []
            QuestionBank.objects.create(
                teacher=teacher,
                subject_id=request.POST.get('subject') or None,
                target_class_id=request.POST.get('target_class') or None,
                topic=request.POST.get('topic', '').strip()[:200],
                question_text=request.POST.get('question_text', '').strip(),
                question_format=request.POST.get('question_format', 'mcq'),
                difficulty=request.POST.get('difficulty', 'medium'),
                options=opts,
                correct_answer=request.POST.get('correct_answer', '').strip(),
                explanation=request.POST.get('explanation', '').strip(),
            )
            messages.success(request, 'Question added to bank.')
            return redirect('teachers:addon_question_bank')

        elif action == 'delete':
            QuestionBank.objects.filter(id=request.POST.get('q_id'), teacher=teacher).delete()
            return redirect('teachers:addon_question_bank')

        elif action == 'create_paper':
            q_ids = request.POST.getlist('selected_questions')
            if q_ids:
                paper = ExamPaper.objects.create(
                    teacher=teacher,
                    title=request.POST.get('paper_title', 'Exam Paper').strip()[:200],
                    subject_id=request.POST.get('paper_subject') or None,
                    target_class_id=request.POST.get('paper_class') or None,
                    duration_minutes=int(request.POST.get('duration', 60) or 60),
                    instructions=request.POST.get('instructions', 'Answer ALL questions.').strip(),
                )
                paper.questions.set(QuestionBank.objects.filter(id__in=q_ids, teacher=teacher))
                messages.success(request, f'Exam paper "{paper.title}" created with {len(q_ids)} questions.')
            return redirect('teachers:addon_question_bank')

    # Filter
    f_subject = request.GET.get('subject')
    f_difficulty = request.GET.get('difficulty')
    f_format = request.GET.get('format')
    qs = QuestionBank.objects.filter(teacher=teacher).select_related('subject', 'target_class')
    if f_subject:
        qs = qs.filter(subject_id=f_subject)
    if f_difficulty:
        qs = qs.filter(difficulty=f_difficulty)
    if f_format:
        qs = qs.filter(question_format=f_format)

    papers = ExamPaper.objects.filter(teacher=teacher).select_related('subject', 'target_class')
    qb_used = get_free_generation_count(teacher, 'question_gen')

    return render(request, 'teachers/addon_question_bank.html', {
        'questions': qs, 'papers': papers, 'subjects': subjects, 'classes': classes,
        'f_subject': f_subject, 'f_difficulty': f_difficulty, 'f_format': f_format,
        'qb_gen_used': qb_used, 'qb_gen_limit': FREE_GENERATION_LIMIT,
        'qb_gen_remaining': max(0, FREE_GENERATION_LIMIT - qb_used),
        'has_qb_addon': has_addon(teacher, 'exam-question-bank'),
    })


@login_required
@requires_addon_freemium('exam-question-bank', action_type='question_gen')
def question_bank_ai(request):
    """AJAX: AI-generate questions for the question bank."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST only'}, status=405)

    body = json.loads(request.body) if request.content_type == 'application/json' else request.POST
    topic = body.get('topic', '').strip()
    subject_name = body.get('subject', '').strip()
    num = min(int(body.get('count', 5) or 5), 10)
    fmt = body.get('format', 'mcq')
    difficulty = body.get('difficulty', 'medium')

    if not topic:
        return JsonResponse({'error': 'Provide a topic.'}, status=400)

    from tenants.ai_quota import check_and_consume, QuotaExceeded
    try:
        check_and_consume(request.tenant, request.user.id, 'question_gen')
    except QuotaExceeded as e:
        return JsonResponse({
            'error': e.user_message, 'error_code': 'quota_exceeded',
            'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost,
        }, status=429)

    try:
        import google.generativeai as genai
        api_key = os.environ.get('GEMINI_API_KEY', '')
        if not api_key:
            return JsonResponse({'error': 'AI not configured.'}, status=500)

        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-2.0-flash')

        fmt_desc = {'mcq': 'multiple choice (4 options A-D)', 'fill': 'fill in the blank',
                    'short': 'short answer', 'truefalse': 'true or false', 'essay': 'essay'}.get(fmt, fmt)
        prompt = (
            f"Generate {num} {difficulty} difficulty {fmt_desc} questions on the topic \"{topic}\" "
            f"for {subject_name or 'a'} Ghanaian JHS class.\n\n"
            "Return ONLY valid JSON array:\n"
            '[{"question":"…","options":["A) …","B) …","C) …","D) …"],"correct":"A","explanation":"…"}, …]\n'
            "For non-MCQ, options can be empty array. For true/false, options should be [\"True\",\"False\"]."
        )
        resp = model.generate_content(prompt)
        text = resp.text.strip()
        if '```' in text:
            text = text.split('```')[1]
            if text.startswith('json'):
                text = text[4:]
        import re
        json_match = re.search(r'\[.*\]', text, re.DOTALL)
        questions = json.loads(json_match.group()) if json_match else json.loads(text)
        return JsonResponse({'questions': questions})

    except QuotaExceeded:
        raise
    except Exception as exc:
        logger.exception('Question bank AI error')
        return _ai_json_error_response(exc)


@login_required
@requires_addon('exam-question-bank')
def addon_exam_paper(request, **kwargs):
    """View/print a specific exam paper."""
    from teachers.models import ExamPaper
    paper_id = request.GET.get('id')
    if not paper_id:
        return redirect('teachers:addon_question_bank')
    paper = get_object_or_404(ExamPaper, id=paper_id, teacher=request.user)
    questions = paper.questions.all().order_by('?')  # shuffled
    return render(request, 'teachers/addon_exam_paper.html', {'paper': paper, 'questions': questions})


# ── 3. Behavior & SEL Tracker ────────────────────────────────────────────

@login_required
@requires_addon('behavior-sel-tracker')
def addon_behavior_tracker(request):
    """Behavior logging and SEL tracking."""
    from teachers.models import BehaviorLog

    teacher = request.user
    classes = Class.objects.filter(academic_year__is_current=True)

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'add':
            student_id = request.POST.get('student_id')
            if student_id:
                BehaviorLog.objects.create(
                    teacher=teacher,
                    student_id=student_id,
                    log_type=request.POST.get('log_type', 'positive'),
                    category=request.POST.get('category', 'other'),
                    note=request.POST.get('note', '').strip(),
                    date=request.POST.get('date') or date.today(),
                )
                messages.success(request, 'Behavior logged.')
            return redirect('teachers:addon_behavior_tracker')
        elif action == 'delete':
            BehaviorLog.objects.filter(id=request.POST.get('log_id'), teacher=teacher).delete()
            return redirect('teachers:addon_behavior_tracker')

    # Filter by class
    sel_class = request.GET.get('class_id')
    logs = BehaviorLog.objects.filter(teacher=teacher).select_related('student__user')
    students = Student.objects.none()
    if sel_class:
        students = Student.objects.filter(current_class_id=sel_class).select_related('user').order_by('user__first_name')
        logs = logs.filter(student__current_class_id=sel_class)

    logs = logs[:100]

    # Stats
    total = BehaviorLog.objects.filter(teacher=teacher).count()
    positive = BehaviorLog.objects.filter(teacher=teacher, log_type='positive').count()
    negative = BehaviorLog.objects.filter(teacher=teacher, log_type='negative').count()

    # Flagged students (3+ negative in last 14 days)
    from django.utils.timezone import now
    from datetime import timedelta
    cutoff = now().date() - timedelta(days=14)
    flagged = (
        BehaviorLog.objects.filter(teacher=teacher, log_type='negative', date__gte=cutoff)
        .values('student__id', 'student__user__first_name', 'student__user__last_name')
        .annotate(count=Count('id'))
        .filter(count__gte=3)
        .order_by('-count')
    )

    return render(request, 'teachers/addon_behavior_tracker.html', {
        'classes': classes, 'students': students, 'logs': logs,
        'sel_class': int(sel_class) if sel_class else None,
        'total': total, 'positive': positive, 'negative': negative,
        'flagged': flagged,
    })


# ── 4. Differentiated Lesson AI ──────────────────────────────────────────

@login_required
def addon_differentiated(request):
    """Build differentiated lesson versions with AI."""
    from teachers.models import DifferentiatedLesson

    teacher = request.user
    subjects = Subject.objects.all()
    classes = Class.objects.filter(academic_year__is_current=True)
    lesson_plans = LessonPlan.objects.filter(teacher__user=teacher).select_related('subject', 'school_class')[:50]

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'save':
            DifferentiatedLesson.objects.create(
                teacher=teacher,
                title=request.POST.get('title', '').strip()[:200] or 'Untitled',
                subject_id=request.POST.get('subject') or None,
                target_class_id=request.POST.get('target_class') or None,
                source_content=request.POST.get('source_content', '').strip(),
                foundational_html=request.POST.get('foundational_html', '').strip(),
                grade_level_html=request.POST.get('grade_level_html', '').strip(),
                extension_html=request.POST.get('extension_html', '').strip(),
            )
            messages.success(request, 'Differentiated lesson saved.')
            return redirect('teachers:addon_differentiated')
        elif action == 'delete':
            DifferentiatedLesson.objects.filter(id=request.POST.get('lesson_id'), teacher=teacher).delete()
            return redirect('teachers:addon_differentiated')

    lessons = DifferentiatedLesson.objects.filter(teacher=teacher).select_related('subject', 'target_class')
    dl_used = get_free_generation_count(teacher, 'differentiate')
    return render(request, 'teachers/addon_differentiated.html', {
        'subjects': subjects, 'classes': classes, 'lesson_plans': lesson_plans,
        'lessons': lessons,
        'dl_gen_used': dl_used, 'dl_gen_limit': FREE_GENERATION_LIMIT,
        'dl_gen_remaining': max(0, FREE_GENERATION_LIMIT - dl_used),
        'has_dl_addon': has_addon(teacher, 'differentiated-lesson-ai'),
    })


@login_required
@requires_addon_freemium('differentiated-lesson-ai', action_type='differentiate')
def differentiated_ai(request):
    """AJAX: Generate three differentiated lesson tiers from source content."""
    if request.method != 'POST':
        return JsonResponse({'error': 'POST only'}, status=405)

    body = json.loads(request.body) if request.content_type == 'application/json' else request.POST
    content = body.get('content', '').strip()
    if not content:
        return JsonResponse({'error': 'Provide lesson content.'}, status=400)

    from tenants.ai_quota import check_and_consume, QuotaExceeded
    try:
        check_and_consume(request.tenant, request.user.id, 'differentiate')
    except QuotaExceeded as e:
        return JsonResponse({
            'error': e.user_message, 'error_code': 'quota_exceeded',
            'used': e.used, 'limit': e.limit, 'addon_boost': e.addon_boost,
        }, status=429)

    try:
        import google.generativeai as genai
        api_key = os.environ.get('GEMINI_API_KEY', '')
        if not api_key:
            return JsonResponse({'error': 'AI not configured.'}, status=500)
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-2.0-flash')
        prompt = (
            "You are a Ghanaian JHS teacher assistant specialising in inclusive education.\n"
            "Take the lesson content below and create THREE differentiated versions as HTML:\n\n"
            "1. **Foundational** — for struggling learners. Simpler language, more scaffolding, "
            "concrete examples, step-by-step guidance, visual aids.\n"
            "2. **Grade Level** — for on-track learners. Standard lesson with clear objectives, "
            "practice activities, and assessment.\n"
            "3. **Extension** — for advanced learners. Deeper analysis, open-ended challenges, "
            "research tasks, connections to real-world applications.\n\n"
            "Use HTML tags: <h3>, <h4>, <ul>, <li>, <p>, <strong>, <em>.\n"
            "Return ONLY valid JSON: {\"foundational\":\"<html>…\",\"grade_level\":\"<html>…\",\"extension\":\"<html>…\"}\n\n"
            f"LESSON CONTENT:\n{content[:3000]}"
        )
        resp = model.generate_content(prompt)
        text = resp.text.strip()
        if '```' in text:
            text = text.split('```')[1]
            if text.startswith('json'):
                text = text[4:]
        import re
        json_match = re.search(r'\{.*\}', text, re.DOTALL)
        result = json.loads(json_match.group()) if json_match else json.loads(text)
        return JsonResponse(result)
    except QuotaExceeded:
        raise
    except Exception as exc:
        logger.exception('Differentiated AI error')
        return _ai_json_error_response(exc)


# ── 5. Live Quiz Engine ──────────────────────────────────────────────────

@login_required
@requires_addon('live-quiz-engine')
def addon_live_quiz(request):
    """Live quiz management — create, list, delete quizzes."""
    from teachers.models import LiveQuiz, LiveQuizQuestion

    teacher = request.user
    subjects = Subject.objects.all()
    classes = Class.objects.filter(academic_year__is_current=True)

    if request.method == 'POST':
        action = request.POST.get('action')
        if action == 'create':
            quiz = LiveQuiz.objects.create(
                teacher=teacher,
                title=request.POST.get('title', 'Quick Quiz').strip()[:200],
                subject_id=request.POST.get('subject') or None,
                target_class_id=request.POST.get('target_class') or None,
            )
            # Parse questions from JSON
            q_json = request.POST.get('questions_json', '[]')
            try:
                q_list = json.loads(q_json)
            except (json.JSONDecodeError, TypeError):
                q_list = []
            for i, q in enumerate(q_list[:20]):
                LiveQuizQuestion.objects.create(
                    quiz=quiz, order=i,
                    text=q.get('text', '')[:500],
                    option_a=q.get('a', '')[:300],
                    option_b=q.get('b', '')[:300],
                    option_c=q.get('c', '')[:300],
                    option_d=q.get('d', '')[:300],
                    correct=q.get('correct', 'A')[:1].upper(),
                    time_limit=min(int(q.get('time', 30) or 30), 120),
                )
            messages.success(request, f'Quiz "{quiz.title}" created with {len(q_list)} questions.')
            return redirect('teachers:addon_live_quiz')

        elif action == 'delete':
            LiveQuiz.objects.filter(id=request.POST.get('quiz_id'), teacher=teacher).delete()
            return redirect('teachers:addon_live_quiz')

    quizzes = LiveQuiz.objects.filter(teacher=teacher).select_related('subject', 'target_class')
    return render(request, 'teachers/addon_live_quiz.html', {
        'quizzes': quizzes, 'subjects': subjects, 'classes': classes,
    })


@login_required
@requires_addon('live-quiz-engine')
def live_quiz_run(request, quiz_id):
    """Teacher view: Run/control a live quiz session."""
    from teachers.models import LiveQuiz
    quiz = get_object_or_404(LiveQuiz, id=quiz_id, teacher=request.user)
    questions = quiz.questions.all()
    if request.method == 'POST' and request.POST.get('action') == 'toggle':
        quiz.is_active = not quiz.is_active
        quiz.save(update_fields=['is_active'])
        return redirect('teachers:live_quiz_run', quiz_id=quiz.id)
    return render(request, 'teachers/addon_live_quiz_run.html', {
        'quiz': quiz, 'questions': questions,
    })


@login_required
@requires_addon('live-quiz-engine')
def live_quiz_api(request, quiz_id):
    """AJAX API for teacher: get live results, advance questions."""
    from teachers.models import LiveQuiz, LiveQuizResponse
    quiz = get_object_or_404(LiveQuiz, id=quiz_id, teacher=request.user)
    questions = list(quiz.questions.values('id', 'order', 'text', 'option_a', 'option_b', 'option_c', 'option_d', 'correct', 'time_limit'))
    # Get response counts
    for q in questions:
        resps = LiveQuizResponse.objects.filter(question_id=q['id'])
        q['total'] = resps.count()
        q['correct_count'] = resps.filter(is_correct=True).count()
        q['breakdown'] = {
            'A': resps.filter(choice='A').count(),
            'B': resps.filter(choice='B').count(),
            'C': resps.filter(choice='C').count(),
            'D': resps.filter(choice='D').count(),
        }
    # Leaderboard
    leaderboard = (
        LiveQuizResponse.objects.filter(question__quiz=quiz, is_correct=True)
        .values('player_name')
        .annotate(score=Count('id'))
        .order_by('-score')[:20]
    )
    return JsonResponse({
        'quiz': {'id': quiz.id, 'title': quiz.title, 'is_active': quiz.is_active, 'code': quiz.join_code},
        'questions': questions,
        'leaderboard': list(leaderboard),
    })


def live_quiz_play(request, code):
    """Student view: join and play a live quiz (no login required)."""
    from teachers.models import LiveQuiz
    quiz = get_object_or_404(LiveQuiz, join_code=code.upper())
    if not quiz.is_active:
        return render(request, 'teachers/addon_live_quiz_closed.html', {'quiz': quiz})
    questions = quiz.questions.all()
    return render(request, 'teachers/addon_live_quiz_play.html', {'quiz': quiz, 'questions': questions})


def live_quiz_student_api(request, code):
    """AJAX: student submits an answer."""
    from teachers.models import LiveQuiz, LiveQuizQuestion, LiveQuizResponse
    if request.method != 'POST':
        return JsonResponse({'error': 'POST only'}, status=405)
    quiz = get_object_or_404(LiveQuiz, join_code=code.upper())
    if not quiz.is_active:
        return JsonResponse({'error': 'Quiz is closed.'}, status=403)

    body = json.loads(request.body) if request.content_type == 'application/json' else request.POST
    q_id = body.get('question_id')
    choice = (body.get('choice', '') or '')[:1].upper()
    player = (body.get('player_name', 'Anonymous') or 'Anonymous')[:100]

    if not q_id or not choice:
        return JsonResponse({'error': 'Missing fields.'}, status=400)

    question = get_object_or_404(LiveQuizQuestion, id=q_id, quiz=quiz)
    is_correct = choice == question.correct.upper()

    obj, created = LiveQuizResponse.objects.get_or_create(
        question=question, player_name=player,
        defaults={'choice': choice, 'is_correct': is_correct},
    )
    if not created:
        return JsonResponse({'error': 'Already answered.', 'is_correct': obj.is_correct}, status=200)

    return JsonResponse({'is_correct': is_correct, 'correct_answer': question.correct})
